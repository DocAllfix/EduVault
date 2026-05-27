"""Upgrade nexus_master.pptx visuals — forza font Montserrat ovunque +
aggiunge sfondo brand decorato a TITLE/CLOSING.

Run from project root:
    python scripts/upgrade_template_visuals.py

Modifica IN-PLACE assets/templates/nexus_master.pptx.
Backup automatico → nexus_master.pptx.pre_upgrade.bak
"""
from __future__ import annotations
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt

TEMPLATE = Path("assets/templates/nexus_master.pptx")
BACKUP = Path("assets/templates/nexus_master.pptx.pre_upgrade.bak")

# Brand colors C.F.P. Montessori (presi da theme.css frontend)
BRAND_PINK = RGBColor(0xC8, 0x2E, 0x6E)
BRAND_GREEN = RGBColor(0x76, 0x9E, 0x2E)
BRAND_DARK = RGBColor(0x1F, 0x1F, 0x2C)
BRAND_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)
FONT_NAME = "Montserrat"


def force_font_on_text_frame(tf):
    """Applica Montserrat a tutti i run di un text frame."""
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.name = FONT_NAME


def upgrade_layout(layout, layout_idx: int) -> int:
    """Applica font Montserrat a tutti i shape testuali del layout.
    Ritorna numero shape modificati.
    """
    modified = 0
    for shape in layout.shapes:
        if not shape.has_text_frame:
            continue
        force_font_on_text_frame(shape.text_frame)
        modified += 1
    return modified


def add_brand_background_to_title_closing(layout, is_title: bool) -> None:
    """Aggiunge un grande rettangolo di sfondo brand colorato dietro tutti gli
    altri elementi sui layout TITLE e CLOSING per dare impatto visivo."""
    slide_w = layout.slide_layout.part.package.presentation_part.presentation.slide_width if False else Emu(int(20.0 * 914400))
    slide_h = Emu(int(11.25 * 914400))

    # Sfondo gradient pink → dark (full slide) inserito come PRIMO elemento
    bg = layout.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND_DARK if is_title else BRAND_PINK
    bg.line.fill.background()  # no border

    # Sposta il bg al fondo dello z-order (deve essere PRIMO nello spTree)
    sp_tree = bg._element.getparent()
    sp_tree.insert(2, bg._element)  # 0,1 sono nvGrpSpPr e grpSpPr standard


def main() -> None:
    if not TEMPLATE.is_file():
        raise SystemExit(f"Template not found: {TEMPLATE}")

    # Backup atomico
    if not BACKUP.is_file():
        shutil.copy2(TEMPLATE, BACKUP)
        print(f"Backup creato: {BACKUP}")

    prs = Presentation(str(TEMPLATE))
    print(f"Layouts in template: {len(prs.slide_layouts)}")

    total_modified = 0
    for i, layout in enumerate(prs.slide_layouts):
        modified = upgrade_layout(layout, i)
        total_modified += modified
        print(f"  Layout {i} '{layout.name}': font Montserrat applicato a {modified} shape")

    # ALSO: master slide (background che eredita tutti i layout)
    master = prs.slide_master
    master_modified = 0
    for shape in master.shapes:
        if shape.has_text_frame:
            force_font_on_text_frame(shape.text_frame)
            master_modified += 1
    print(f"  Master: font Montserrat applicato a {master_modified} shape")

    prs.save(str(TEMPLATE))
    print(f"\nTemplate aggiornato: {TEMPLATE}")
    print(f"Totale shape modificati: {total_modified + master_modified}")


if __name__ == "__main__":
    main()
