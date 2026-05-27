"""Validator strutturale per nexus_master_v4_patched.pptx — FIX #30.1.

Conferma analista: livello strutturale (XML assertions) deterministico, in CI
bloccante. Cattura la regressione che ci interessa davvero: la lstStyle che
torna vuota, il placeholder PICTURE rinominato, l'anchor che torna a top,
la dimensione slide alterata.

Niente pixel-diff (esplicitamente sconsigliato dall'analista: rendering
LibreOffice flaky tra versioni).

Esecuzione:
    python scripts/validate_master_v4.py
Exit code:
    0 = template OK, pipeline può procedere
    1 = template KO, dettagli stampati su stderr

Uso CI: aggiungere `python scripts/validate_master_v4.py` come step in
pre-build / pre-commit hook.
"""

from __future__ import annotations

import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

TEMPLATE = Path("assets/templates/nexus_master_v4_patched.pptx")

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

# Layout attesi: (idx, name, placeholder_specs: dict[idx, type_enum_value])
# PP_PLACEHOLDER values: TITLE=1, BODY=2, PICTURE=18, SUBTITLE=4
EXPECTED_LAYOUTS = [
    (0, "NX TITLE",          {0: 1, 1: 4}),       # title + subtitle (SUBTITLE type=4)
    (1, "NX MODULE_OPEN",    {0: 1, 1: 2}),       # module_num (title) + module_title (body)
    (2, "NX CONTENT_TEXT",   {0: 1, 1: 2}),       # title + body
    (3, "NX CONTENT_IMAGE",  {0: 1, 1: 2, 2: 18}),# title + body + picture
    (4, "NX DIAGRAM",        {0: 1, 2: 18}),      # title + picture (body-less)
    (5, "NX QUIZ",           {0: 1}),             # question (title) + cards custom
    (6, "NX CASE_STUDY",     {0: 1}),             # title + 3 sezioni custom
    (7, "NX MODULE_CLOSE",   {0: 1, 1: 2, 2: 2}), # module_num + module_title + recap_body
    (8, "NX RECAP",          {0: 1, 1: 2}),       # title + recap_body
    (9, "NX CLOSING",        {0: 1}),             # title (no subtitle, just tagline)
]

# Body placeholders con bullet ereditato — devono avere lstStyle/lvl1pPr/buChar
# popolato col carattere corretto.
BULLET_PLACEHOLDERS = {
    "NX CONTENT_TEXT":  [("nx_body", "•")],
    "NX CONTENT_IMAGE": [("nx_body", "•")],
    "NX RECAP":         [("nx_recap_body", "✓")],
    "NX MODULE_CLOSE":  [("nx_recap_body", "✓")],
}

# Body placeholders che devono avere anchor=ctr (FIX #30 P2)
ANCHOR_CTR_PLACEHOLDERS = {
    "NX CONTENT_TEXT":  ["nx_body"],
    "NX CONTENT_IMAGE": ["nx_body"],
    "NX RECAP":         ["nx_recap_body"],
    "NX MODULE_CLOSE":  ["nx_recap_body"],
}

# Shape canonici che ogni layout content deve avere (per render coerente)
REQUIRED_SHAPES = {
    "NX CONTENT_TEXT":  {"nx_title", "nx_body", "nx_accent_v", "nx_mini_bar", "nx_ref", "nx_page"},
    "NX CONTENT_IMAGE": {"nx_title", "nx_body", "nx_image_box", "nx_caption", "nx_accent_v", "nx_mini_bar", "nx_ref", "nx_page"},
    "NX DIAGRAM":       {"nx_title", "nx_diagram_box", "nx_caption", "nx_accent_v", "nx_mini_bar", "nx_ref", "nx_page"},
    "NX RECAP":         {"nx_title", "nx_recap_body", "nx_page"},  # nx_page added in patch_master_v4
}


class ValidationError(Exception):
    pass


def check_slide_size(prs: Presentation) -> list[str]:
    """13.333 x 7.5 inches (PPTX widescreen Microsoft standard)."""
    errors = []
    w_in = prs.slide_width / 914400
    h_in = prs.slide_height / 914400
    if abs(w_in - 13.333) > 0.05:
        errors.append(f"Slide width {w_in:.3f}\" != 13.333\"")
    if abs(h_in - 7.5) > 0.05:
        errors.append(f"Slide height {h_in:.3f}\" != 7.5\"")
    return errors


def check_layout_count(prs: Presentation) -> list[str]:
    n = len(prs.slide_layouts)
    if n != 10:
        return [f"Layout count {n} != 10"]
    return []


def check_layout_names_and_placeholders(prs: Presentation) -> list[str]:
    errors = []
    for idx, expected_name, expected_phs in EXPECTED_LAYOUTS:
        if idx >= len(prs.slide_layouts):
            errors.append(f"Layout idx={idx} missing")
            continue
        layout = prs.slide_layouts[idx]
        if layout.name != expected_name:
            errors.append(f"Layout {idx} name {layout.name!r} != {expected_name!r}")
        for ph in layout.placeholders:
            ph_idx = ph.placeholder_format.idx
            ph_type_val = int(ph.placeholder_format.type)
            if ph_idx in expected_phs:
                if ph_type_val != expected_phs[ph_idx]:
                    errors.append(
                        f"Layout {idx} ({layout.name}) placeholder idx={ph_idx} "
                        f"type={ph_type_val} != {expected_phs[ph_idx]}"
                    )
        # Check tutti i placeholder attesi presenti
        present_idxs = {ph.placeholder_format.idx for ph in layout.placeholders}
        missing = set(expected_phs.keys()) - present_idxs
        if missing:
            errors.append(
                f"Layout {idx} ({layout.name}) missing placeholder idx={missing}"
            )
    return errors


def check_bullet_liststyle(prs: Presentation) -> list[str]:
    """Per ogni body placeholder con bullet ereditato: lstStyle/lvl1pPr/buChar
    deve essere popolato col carattere corretto + marL/indent + buFont.
    """
    errors = []
    for layout in prs.slide_layouts:
        if layout.name not in BULLET_PLACEHOLDERS:
            continue
        for shape_name, expected_char in BULLET_PLACEHOLDERS[layout.name]:
            shape = None
            for sh in layout.shapes:
                if sh.name == shape_name:
                    shape = sh
                    break
            if shape is None:
                errors.append(f"Layout {layout.name!r}: shape {shape_name!r} not found")
                continue
            txBody = shape._element.find(f"{P}txBody")
            if txBody is None:
                errors.append(f"{layout.name}/{shape_name}: no txBody")
                continue
            lstStyle = txBody.find(f"{A}lstStyle")
            if lstStyle is None:
                errors.append(f"{layout.name}/{shape_name}: lstStyle missing")
                continue
            lvl1 = lstStyle.find(f"{A}lvl1pPr")
            if lvl1 is None:
                errors.append(f"{layout.name}/{shape_name}: lvl1pPr missing (lstStyle empty)")
                continue
            buChar = lvl1.find(f"{A}buChar")
            if buChar is None:
                errors.append(f"{layout.name}/{shape_name}: buChar missing")
                continue
            char = buChar.get("char")
            if char != expected_char:
                errors.append(
                    f"{layout.name}/{shape_name}: buChar={char!r} != {expected_char!r}"
                )
            buFont = lvl1.find(f"{A}buFont")
            if buFont is None or not buFont.get("typeface"):
                errors.append(f"{layout.name}/{shape_name}: buFont missing/empty")
            marL = lvl1.get("marL")
            indent = lvl1.get("indent")
            if marL is None or indent is None:
                errors.append(
                    f"{layout.name}/{shape_name}: marL/indent missing on lvl1pPr"
                )
    return errors


def check_anchor_ctr(prs: Presentation) -> list[str]:
    """Body placeholders devono avere bodyPr@anchor='ctr'."""
    errors = []
    for layout in prs.slide_layouts:
        if layout.name not in ANCHOR_CTR_PLACEHOLDERS:
            continue
        for shape_name in ANCHOR_CTR_PLACEHOLDERS[layout.name]:
            shape = None
            for sh in layout.shapes:
                if sh.name == shape_name:
                    shape = sh
                    break
            if shape is None:
                continue  # già loggato da check_bullet_liststyle
            txBody = shape._element.find(f"{P}txBody")
            if txBody is None:
                continue
            bodyPr = txBody.find(f"{A}bodyPr")
            anchor = bodyPr.get("anchor") if bodyPr is not None else None
            if anchor != "ctr":
                errors.append(
                    f"{layout.name}/{shape_name}: anchor={anchor!r} != 'ctr'"
                )
    return errors


def check_required_shapes(prs: Presentation) -> list[str]:
    """Ogni layout content deve avere gli shape canonici (logo, ref, page, etc)."""
    errors = []
    for layout in prs.slide_layouts:
        if layout.name not in REQUIRED_SHAPES:
            continue
        present = {sh.name for sh in layout.shapes}
        missing = REQUIRED_SHAPES[layout.name] - present
        if missing:
            errors.append(f"{layout.name}: missing required shapes {missing}")
    return errors


def main() -> int:
    if not TEMPLATE.exists():
        print(f"❌ Template missing: {TEMPLATE}", file=sys.stderr)
        return 1

    prs = Presentation(str(TEMPLATE))
    print(f"=== validate_master_v4.py — {TEMPLATE.name} ===")
    print()

    all_errors = []
    all_errors.extend(("[size]", e) for e in check_slide_size(prs))
    all_errors.extend(("[count]", e) for e in check_layout_count(prs))
    all_errors.extend(("[layout]", e) for e in check_layout_names_and_placeholders(prs))
    all_errors.extend(("[bullet]", e) for e in check_bullet_liststyle(prs))
    all_errors.extend(("[anchor]", e) for e in check_anchor_ctr(prs))
    all_errors.extend(("[shapes]", e) for e in check_required_shapes(prs))

    if not all_errors:
        print("✅ ALL CHECKS PASSED")
        print()
        print(f"Slide size: {prs.slide_width/914400:.3f}\" x {prs.slide_height/914400:.3f}\"")
        print(f"Layouts: {len(prs.slide_layouts)} (expected 10)")
        print(f"Bullet placeholders verified: {sum(len(v) for v in BULLET_PLACEHOLDERS.values())}")
        print(f"Anchor ctr placeholders verified: {sum(len(v) for v in ANCHOR_CTR_PLACEHOLDERS.values())}")
        return 0
    else:
        print(f"❌ {len(all_errors)} VALIDATION ERRORS:")
        print()
        for prefix, err in all_errors:
            print(f"  {prefix} {err}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    sys.exit(main())
