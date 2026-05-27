"""Upgrade nexus_master.pptx applicando le best-practice Anthropic skill pptx.

Modifica IN-PLACE il template per renderlo visivamente professionale:
1. Font hierarchy: Inter per titoli + Calibri per body (specifiche Anthropic)
2. Font sizes esatte: title 36pt, section 24pt, body 16pt, caption 11pt
3. Palette C.F.P. Montessori applicata in modo dominante (60-20-10-10)
4. Rimuove accent lines sotto titoli (anti "AI-generated look")
5. Logo a dimensione visibile ma non invadente (1.5x0.4 inch)
6. Footer minimal su content slides

Backup: nexus_master.pptx.pre_professional.bak
"""
from __future__ import annotations
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches, Emu

TEMPLATE = Path("assets/templates/nexus_master.pptx")
BACKUP = Path("assets/templates/nexus_master.pptx.pre_professional.bak")

# Best-practice Anthropic skill pptx
TITLE_FONT = "Inter"
BODY_FONT = "Calibri"
TITLE_SIZE = Pt(36)
SECTION_SIZE = Pt(24)
BODY_SIZE = Pt(16)
CAPTION_SIZE = Pt(11)
SMALL_SIZE = Pt(10)

# Palette brand C.F.P. Montessori
BRAND_PINK = RGBColor(0xC8, 0x2E, 0x6E)
BRAND_GREEN = RGBColor(0x76, 0x9E, 0x2E)
BRAND_DARK = RGBColor(0x1F, 0x1F, 0x2C)
BRAND_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)
TEXT_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)

# Logo dimensions (FIX #17 retake): 1.5x0.4 inch sufficiente ma non invadente
LOGO_W = Inches(1.5)
LOGO_H = Inches(0.4)
SLIDE_W = Inches(20)
SLIDE_H = Inches(11.25)
LOGO_MARGIN = Inches(0.2)


def style_text_frame(tf, font_name, font_size, color=None, bold=False):
    """Applica font/size/color/bold a tutti i run."""
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = font_size
            if color is not None:
                run.font.color.rgb = color
            run.font.bold = bold


def is_accent_line(shape):
    """Riconosce le accent line decorative (forme molto piatte sotto al titolo)."""
    try:
        h = shape.height / 914400 if shape.height else 0
        w = shape.width / 914400 if shape.width else 0
    except Exception:
        return False
    # Linee orizzontali sottili (altezza < 0.1 inch, larghezza > 5 inch)
    if h < 0.15 and w > 5 and not (shape.has_text_frame and shape.text_frame.text.strip()):
        return True
    return False


def main():
    if not TEMPLATE.is_file():
        raise SystemExit(f"Template not found: {TEMPLATE}")
    if not BACKUP.is_file():
        shutil.copy2(TEMPLATE, BACKUP)
        print(f"Backup creato: {BACKUP}")
    prs = Presentation(str(TEMPLATE))
    print(f"Template: {len(prs.slide_layouts)} layouts, slide {prs.slide_width/914400:.1f}x{prs.slide_height/914400:.1f}")

    accent_lines_removed = 0
    logos_resized = 0
    text_styled = 0

    for li, layout in enumerate(prs.slide_layouts):
        layout_name = layout.name
        print(f"\n=== Layout {li} '{layout_name}' ===")

        # 1. Rimuovi accent lines sotto titoli (anti-AI look per Anthropic)
        shapes_to_remove = []
        for sh in layout.shapes:
            if is_accent_line(sh):
                shapes_to_remove.append(sh)
        for sh in shapes_to_remove:
            sh._element.getparent().remove(sh._element)
            accent_lines_removed += 1
            print(f"  - Removed accent line ({sh.width/914400:.1f}x{sh.height/914400:.1f})")

        # 2. Ingrandisci e riposiziona logo
        for sh in layout.shapes:
            if "PICTURE" in str(sh.shape_type) and sh.name == "Image 0":
                if layout_name in ("TITLE", "CLOSING"):
                    # Logo grande centrato per intro/outro
                    sh.width = Inches(2.5)
                    sh.height = Inches(0.7)
                else:
                    # Content slides: logo medio basso-destra
                    sh.width = LOGO_W
                    sh.height = LOGO_H
                    sh.left = SLIDE_W - LOGO_W - LOGO_MARGIN
                    sh.top = SLIDE_H - LOGO_H - LOGO_MARGIN
                logos_resized += 1
                print(f"  + Logo resized to {sh.width/914400:.1f}x{sh.height/914400:.1f}")

        # 3. Applica font hierarchy ai text shape
        text_shapes = [s for s in layout.shapes if s.has_text_frame]
        for idx, sh in enumerate(text_shapes):
            tf = sh.text_frame
            text = tf.text.strip()
            if not text:
                continue
            # Decisione font/size basata su contesto:
            # - Titoli (large + pochi caratteri) → TITLE_SIZE Inter Bold
            # - Section badge ("CASO STUDIO", "RIEPILOGO") → SECTION_SIZE Inter Bold pink
            # - Body principale → BODY_SIZE Calibri
            # - Caption ("Art. 18...", page num) → CAPTION_SIZE Calibri muted
            w = sh.width / 914400 if sh.width else 0
            h = sh.height / 914400 if sh.height else 0
            n_chars = len(text)

            if h < 0.4 and n_chars < 30 and text.isupper():
                # Section badge
                style_text_frame(tf, TITLE_FONT, SECTION_SIZE, BRAND_PINK, bold=True)
            elif h > 0.4 and n_chars < 80 and idx == 0:
                # Title principale
                style_text_frame(tf, TITLE_FONT, TITLE_SIZE, BRAND_DARK, bold=True)
            elif h < 0.4 and n_chars < 50:
                # Caption/footer (es. "Art. 18 D.Lgs 81/08", "2 / 8")
                style_text_frame(tf, BODY_FONT, CAPTION_SIZE, TEXT_MUTED)
            else:
                # Body principale
                style_text_frame(tf, BODY_FONT, BODY_SIZE, TEXT_BLACK)
            text_styled += 1

    prs.save(str(TEMPLATE))
    print(f"\n=== Summary ===")
    print(f"Accent lines removed: {accent_lines_removed}")
    print(f"Logos resized: {logos_resized}")
    print(f"Text shapes styled: {text_styled}")
    print(f"Template salvato: {TEMPLATE}")


if __name__ == "__main__":
    main()
