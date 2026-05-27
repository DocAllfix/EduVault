"""Build nexus_master_v4.pptx from scratch — FIX #30.0-ter (2026-05-26).

Genera un template PowerPoint NUOVO a dimensione standard 13.333" × 7.5"
(=33.87cm × 19.05cm = 12192000 × 6858000 EMU), con 10 layout veri e
placeholder PowerPoint nativi (idx + type), pronti per il SlideBuilderV2.

Differenze vs v3 (build_master_v3.py):
- v3 partiva da v2 (20×11.25") e CONVERTIVA shape esistenti in placeholder.
  v4 costruisce DA ZERO un master a dimensione standard.
- v3 aveva 8 layout. v4 ne ha 10: aggiunti MODULE_OPEN (idx 1) e MODULE_CLOSE
  (idx 7) per i bookend modulo decisi nel piano FIX #30.
- v3 ereditava la palette Slidesgo. v4 usa solo i 4 colori CFP.

Layout finali (idx → name):
  0 NX TITLE          (corso)
  1 NX MODULE_OPEN    (NUOVO, apertura modulo)
  2 NX CONTENT_TEXT
  3 NX CONTENT_IMAGE
  4 NX DIAGRAM
  5 NX QUIZ
  6 NX CASE_STUDY
  7 NX MODULE_CLOSE   (NUOVO, chiusura modulo)
  8 NX RECAP          (corso)
  9 NX CLOSING

Esecuzione (host, no container):
    python scripts/build_master_v4.py
Output:
    assets/templates/nexus_master_v4.pptx
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Emu, Inches, Pt

# ═══════════════════════════════════════════════════════════════════════════
# CONFIG
# ═══════════════════════════════════════════════════════════════════════════

OUT_PATH = Path("assets/templates/nexus_master_v4.pptx")
LOGO_PATH = Path("assets/brand/cfp_montessori_logo.jpeg")

# Slide size: PPTX widescreen Microsoft standard
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# CFP Montessori palette (FIX #30 decision)
COL_PRIMARY = RGBColor(0xC8, 0x2E, 0x6E)   # rosa/magenta
COL_SECONDARY = RGBColor(0x76, 0x9E, 0x2E)  # verde oliva
COL_TEXT = RGBColor(0x1F, 0x1F, 0x1F)       # nero soft
COL_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COL_PINK_SOFT = RGBColor(0xFD, 0xE7, 0xEE)  # rosa chiaro background
COL_GREEN_SOFT = RGBColor(0xF4, 0xF8, 0xE8)  # verde chiaro background
COL_AMBER_SOFT = RGBColor(0xFF, 0xF4, 0xE0)  # ambra chiaro background
COL_AMBER = RGBColor(0xF3, 0x9C, 0x12)      # ambra (case study risultato)

# Font primario (Inter se installato, altrimenti Segoe UI fallback gestito da PPT)
FONT_FAMILY = "Inter"

# XML namespaces
A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


# ═══════════════════════════════════════════════════════════════════════════
# HELPER: shape naming + placeholder injection (riusa pattern da v3)
# ═══════════════════════════════════════════════════════════════════════════


def set_shape_name(shape, name: str) -> None:
    """Set shape name visible in PowerPoint Selection Pane."""
    sp = shape._element
    nvSpPr = sp.find(f"{P}nvSpPr") or sp.find(f"{P}nvPicPr") or sp.find(f"{P}nvCxnSpPr")
    if nvSpPr is None:
        return
    cNvPr = nvSpPr.find(f"{P}cNvPr")
    if cNvPr is not None:
        cNvPr.set("name", name)


def inject_placeholder(shape, ph_type: str, ph_idx: int) -> None:
    """Convert a normal shape into a PowerPoint placeholder by injecting
    <p:nvSpPr><p:nvPr><p:ph type="..." idx="..."/></p:nvPr></p:nvSpPr>.
    """
    sp = shape._element
    nvSpPr = sp.find(f"{P}nvSpPr")
    if nvSpPr is None:
        return
    nvPr = nvSpPr.find(f"{P}nvPr")
    if nvPr is None:
        nvPr = etree.SubElement(nvSpPr, f"{P}nvPr")
    for old in nvPr.findall(f"{P}ph"):
        nvPr.remove(old)
    ph = etree.SubElement(nvPr, f"{P}ph")
    ph.set("type", ph_type)
    ph.set("idx", str(ph_idx))


def inject_normautofit(shape) -> None:
    """Add <a:normAutofit/> in <a:bodyPr> for native font auto-shrink."""
    sp = shape._element
    txBody = sp.find(f"{P}txBody")
    if txBody is None:
        return
    bodyPr = txBody.find(f"{A}bodyPr")
    if bodyPr is None:
        bodyPr = etree.SubElement(txBody, f"{A}bodyPr")
        txBody.insert(0, bodyPr)
    for tag in ("normAutofit", "spAutoFit", "noAutofit"):
        for old in bodyPr.findall(f"{A}{tag}"):
            bodyPr.remove(old)
    etree.SubElement(bodyPr, f"{A}normAutofit")


def inject_bullet_style(shape, bullet_char: str = "•", color: RGBColor = COL_PRIMARY) -> None:
    """Inject <a:lstStyle><a:lvl1pPr><a:buChar char="•"/></...> so bullets
    inherit automatically on slides using this layout's BODY placeholder.
    """
    sp = shape._element
    txBody = sp.find(f"{P}txBody")
    if txBody is None:
        return
    lstStyle = txBody.find(f"{A}lstStyle")
    if lstStyle is None:
        bodyPr = txBody.find(f"{A}bodyPr")
        lstStyle = etree.Element(f"{A}lstStyle")
        if bodyPr is not None:
            bodyPr.addnext(lstStyle)
        else:
            txBody.insert(0, lstStyle)
    for old in lstStyle.findall(f"{A}lvl1pPr"):
        lstStyle.remove(old)
    lvl1 = etree.SubElement(lstStyle, f"{A}lvl1pPr")
    lvl1.set("marL", "342900")
    lvl1.set("indent", "-342900")
    buClr = etree.SubElement(lvl1, f"{A}buClr")
    srgbClr = etree.SubElement(buClr, f"{A}srgbClr")
    srgbClr.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
    buFont = etree.SubElement(lvl1, f"{A}buFont")
    buFont.set("typeface", "Arial")
    buChar = etree.SubElement(lvl1, f"{A}buChar")
    buChar.set("char", bullet_char)


def set_text(
    shape,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = COL_TEXT,
    align: PP_ALIGN | None = None,
    anchor: MSO_ANCHOR | None = None,
    font_family: str = FONT_FAMILY,
    italic: bool = False,
) -> None:
    """Set text in a shape with consistent styling."""
    tf = shape.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    if anchor is not None:
        tf.vertical_anchor = anchor
    for run in p.runs:
        run.font.name = font_family
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color


def add_rect(
    shapes,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    fill: RGBColor | None,
    line_color: RGBColor | None = None,
    line_width_pt: float = 0.0,
    shape_type: int = MSO_SHAPE.RECTANGLE,
):
    """Add a rectangle (or other shape) with fill + optional border."""
    s = shapes.add_shape(shape_type, left, top, width, height)
    if fill is not None:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color is not None and line_width_pt > 0:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_width_pt)
    else:
        s.line.fill.background()
    return s


def add_textbox(
    shapes,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    text: str = "",
    **text_kwargs,
):
    """Add a textbox with text + style."""
    tb = shapes.add_textbox(left, top, width, height)
    if text:
        set_text(tb, text, **text_kwargs)
    return tb


def add_logo(shapes, left: Emu, top: Emu, width: Emu, height: Emu, name: str = "nx_logo"):
    """Add CFP logo picture."""
    if not LOGO_PATH.exists():
        # Fallback: textbox con "CFP" se logo non c'è (non blocca generazione)
        tb = add_textbox(shapes, left, top, width, height, "CFP", font_size=12, bold=True, color=COL_PRIMARY)
        set_shape_name(tb, name)
        return tb
    pic = shapes.add_picture(str(LOGO_PATH), left, top, width=width, height=height)
    set_shape_name(pic, name)
    return pic


# ═══════════════════════════════════════════════════════════════════════════
# LAYOUT BUILDERS — 10 layout
# ═══════════════════════════════════════════════════════════════════════════


def build_layout_title(layout) -> None:
    """Layout 0: NX TITLE — corso opener."""
    shapes = layout.shapes

    # Top band rosa (full width)
    band = add_rect(shapes, 0, 0, SLIDE_W, Cm(1.0), COL_PRIMARY)
    set_shape_name(band, "nx_band_top")

    # Blob verde decorativo in basso-sinistra
    blob_bl = add_rect(
        shapes, Cm(-3), Cm(13), Cm(10), Cm(10), COL_SECONDARY, shape_type=MSO_SHAPE.OVAL
    )
    set_shape_name(blob_bl, "nx_blob_bl")

    # Blob rosa in alto-destra
    blob_tr = add_rect(
        shapes, Cm(28), Cm(-3), Cm(8), Cm(8), COL_PRIMARY, shape_type=MSO_SHAPE.OVAL
    )
    set_shape_name(blob_tr, "nx_blob_tr")

    # TITLE placeholder (idx=0)
    title = add_textbox(
        shapes, Cm(2.5), Cm(6), Cm(28.87), Cm(3),
        "Titolo del corso", font_size=44, bold=True, color=COL_TEXT,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    set_shape_name(title, "nx_title")
    inject_placeholder(title, "title", 0)

    # SUBTITLE placeholder (idx=1)
    subtitle = add_textbox(
        shapes, Cm(2.5), Cm(9.5), Cm(28.87), Cm(2),
        "Sottotitolo", font_size=22, color=COL_PRIMARY,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    set_shape_name(subtitle, "nx_subtitle")
    inject_placeholder(subtitle, "body", 1)

    # 3 puntini decorativi in basso
    for i, x_cm in enumerate([14.5, 16.5, 18.5]):
        dot = add_rect(
            shapes, Cm(x_cm), Cm(17.5), Cm(0.4), Cm(0.4),
            COL_PRIMARY, shape_type=MSO_SHAPE.OVAL,
        )
        set_shape_name(dot, f"nx_dot{i+1}")

    # Logo (basso-destra)
    add_logo(shapes, Cm(30), Cm(17.3), Cm(2.5), Cm(1.2), "nx_logo")


def build_layout_module_open(layout) -> None:
    """Layout 1: NX MODULE_OPEN — apertura modulo (NUOVO)."""
    shapes = layout.shapes

    # Accent verticale rosa sinistra (banda larga, è un opener)
    accent = add_rect(shapes, 0, 0, Cm(1.0), SLIDE_H, COL_PRIMARY)
    set_shape_name(accent, "nx_accent_v")

    # Numero modulo grande, centrato (TITLE idx=0)
    num = add_textbox(
        shapes, Cm(2), Cm(5), Cm(29.87), Cm(4),
        "MODULO N", font_size=72, bold=True, color=COL_PRIMARY,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    set_shape_name(num, "nx_module_num")
    inject_placeholder(num, "title", 0)

    # Titolo del modulo (BODY idx=1)
    title = add_textbox(
        shapes, Cm(2), Cm(10), Cm(29.87), Cm(3),
        "Titolo del modulo", font_size=32, color=COL_TEXT,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    set_shape_name(title, "nx_module_title")
    inject_placeholder(title, "body", 1)

    # Underline verde sotto al titolo (decorazione)
    underline = add_rect(
        shapes, Cm(13), Cm(13.2), Cm(7.87), Cm(0.3), COL_SECONDARY,
    )
    set_shape_name(underline, "nx_module_underline")

    # Logo basso-destra
    add_logo(shapes, Cm(30), Cm(17.3), Cm(2.5), Cm(1.2), "nx_logo")


def _add_common_content_decorations(shapes) -> None:
    """Add accent_v + mini_bar + ref + page + logo (common to CONTENT_TEXT,
    CONTENT_IMAGE, DIAGRAM layouts).
    """
    accent = add_rect(shapes, 0, 0, Cm(0.2), SLIDE_H, COL_PRIMARY)
    set_shape_name(accent, "nx_accent_v")

    mini = add_rect(
        shapes, Cm(1.5), Cm(3.5), Cm(2.5), Cm(0.15), COL_PRIMARY,
    )
    set_shape_name(mini, "nx_mini_bar")

    # Box riferimento normativo (basso-sinistra, rosa chiaro)
    ref = add_rect(
        shapes, Cm(1), Cm(17.5), Cm(10), Cm(1.0), COL_PINK_SOFT,
        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    set_shape_name(ref, "nx_ref")
    set_text(
        ref, "Riferimento",
        font_size=10, color=COL_PRIMARY,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # Numero pagina (basso-destra)
    page = add_textbox(
        shapes, Cm(24), Cm(17.5), Cm(5), Cm(1.0),
        "000 / 000", font_size=10, color=COL_TEXT,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
    )
    set_shape_name(page, "nx_page")

    # Logo
    add_logo(shapes, Cm(30), Cm(17.3), Cm(2.5), Cm(1.2), "nx_logo")


def build_layout_content_text(layout) -> None:
    """Layout 2: NX CONTENT_TEXT — workhorse, titolo + 4-6 bullets."""
    shapes = layout.shapes

    # TITLE placeholder (idx=0)
    title = add_textbox(
        shapes, Cm(1.5), Cm(1.5), Cm(31), Cm(1.8),
        "Titolo della slide", font_size=28, bold=True, color=COL_TEXT,
    )
    set_shape_name(title, "nx_title")
    inject_placeholder(title, "title", 0)

    # BODY placeholder (idx=1) — supporta bullet ereditati
    body = add_textbox(
        shapes, Cm(1.5), Cm(4), Cm(31), Cm(13),
        "Primo punto", font_size=20, color=COL_TEXT,
    )
    set_shape_name(body, "nx_body")
    inject_placeholder(body, "body", 1)
    inject_normautofit(body)
    inject_bullet_style(body, "•", COL_PRIMARY)

    _add_common_content_decorations(shapes)


def build_layout_content_image(layout) -> None:
    """Layout 3: NX CONTENT_IMAGE — titolo + bullet + immagine destra."""
    shapes = layout.shapes

    # TITLE (idx=0)
    title = add_textbox(
        shapes, Cm(1.5), Cm(1.5), Cm(31), Cm(1.8),
        "Titolo della slide", font_size=28, bold=True, color=COL_TEXT,
    )
    set_shape_name(title, "nx_title")
    inject_placeholder(title, "title", 0)

    # BODY (idx=1) — più stretto per far spazio all'immagine
    body = add_textbox(
        shapes, Cm(1.5), Cm(4), Cm(17), Cm(13),
        "Primo punto", font_size=20, color=COL_TEXT,
    )
    set_shape_name(body, "nx_body")
    inject_placeholder(body, "body", 1)
    inject_normautofit(body)
    inject_bullet_style(body, "•", COL_PRIMARY)

    # PICTURE placeholder (idx=2) — destra
    img_box = add_rect(
        shapes, Cm(19.5), Cm(4), Cm(12.5), Cm(11),
        COL_PINK_SOFT, line_color=COL_PRIMARY, line_width_pt=0.5,
    )
    set_shape_name(img_box, "nx_image_box")
    inject_placeholder(img_box, "pic", 2)
    set_text(
        img_box, "Immagine",
        font_size=14, color=COL_PRIMARY, italic=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # Caption sotto immagine
    caption = add_textbox(
        shapes, Cm(19.5), Cm(15.3), Cm(12.5), Cm(0.8),
        "Didascalia immagine", font_size=10, italic=True, color=COL_TEXT,
        align=PP_ALIGN.CENTER,
    )
    set_shape_name(caption, "nx_caption")

    _add_common_content_decorations(shapes)


def build_layout_diagram(layout) -> None:
    """Layout 4: NX DIAGRAM — titolo + diagramma grande + caption sotto."""
    shapes = layout.shapes

    # TITLE (idx=0)
    title = add_textbox(
        shapes, Cm(1.5), Cm(1.5), Cm(31), Cm(1.8),
        "Titolo della slide", font_size=28, bold=True, color=COL_TEXT,
    )
    set_shape_name(title, "nx_title")
    inject_placeholder(title, "title", 0)

    # PICTURE placeholder (idx=2) — diagram box grande
    dia_box = add_rect(
        shapes, Cm(2), Cm(4), Cm(29.87), Cm(11),
        COL_WHITE, line_color=COL_SECONDARY, line_width_pt=1.0,
    )
    set_shape_name(dia_box, "nx_diagram_box")
    inject_placeholder(dia_box, "pic", 2)
    set_text(
        dia_box, "Diagramma",
        font_size=16, color=COL_SECONDARY, italic=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # Caption obbligatoria SOTTO il diagramma (qui va il testo lungo,
    # non dentro al diagramma — FIX #30.4 Q4-bis)
    caption = add_textbox(
        shapes, Cm(2), Cm(15.3), Cm(29.87), Cm(1.0),
        "Spiegazione del diagramma (max 200 caratteri)",
        font_size=12, italic=True, color=COL_TEXT,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    set_shape_name(caption, "nx_caption")

    _add_common_content_decorations(shapes)


def build_layout_quiz(layout) -> None:
    """Layout 5: NX QUIZ — domanda + 4 cards opzione."""
    shapes = layout.shapes

    # Accent verticale
    accent = add_rect(shapes, 0, 0, Cm(0.2), SLIDE_H, COL_PRIMARY)
    set_shape_name(accent, "nx_accent_v")

    # Question box (idx=0 title)
    question_bg = add_rect(
        shapes, Cm(1.5), Cm(1.5), Cm(31), Cm(2.5), COL_PINK_SOFT,
        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    set_shape_name(question_bg, "nx_question_box")

    title = add_textbox(
        shapes, Cm(1.5), Cm(1.5), Cm(31), Cm(2.5),
        "Quale è la risposta corretta?", font_size=24, color=COL_TEXT, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    set_shape_name(title, "nx_question")
    inject_placeholder(title, "title", 0)

    # 4 cards opzioni in layout 2x2
    card_w = Cm(15)
    card_h = Cm(5.5)
    gap = Cm(1)
    start_x = Cm(1.5)
    start_y = Cm(5)

    positions = [
        ("a", start_x, start_y),
        ("b", start_x + card_w + gap, start_y),
        ("c", start_x, start_y + card_h + gap),
        ("d", start_x + card_w + gap, start_y + card_h + gap),
    ]

    for letter, lx, ly in positions:
        # Card background
        card = add_rect(
            shapes, lx, ly, card_w, card_h, COL_WHITE,
            line_color=COL_PRIMARY, line_width_pt=1.0,
            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
        )
        set_shape_name(card, f"nx_option_card_{letter}")

        # Letter circle rosa
        letter_circle = add_rect(
            shapes, lx + Cm(0.6), ly + Cm(0.6), Cm(1.4), Cm(1.4),
            COL_PRIMARY, shape_type=MSO_SHAPE.OVAL,
        )
        set_shape_name(letter_circle, f"nx_letter_{letter}")

        # Letter text dentro al cerchio
        letter_text = add_textbox(
            shapes, lx + Cm(0.6), ly + Cm(0.6), Cm(1.4), Cm(1.4),
            letter.upper(), font_size=18, bold=True, color=COL_WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        set_shape_name(letter_text, f"nx_letter_text_{letter}")

        # Option text
        option_text = add_textbox(
            shapes, lx + Cm(2.5), ly + Cm(0.5), card_w - Cm(3), card_h - Cm(1),
            f"Opzione {letter.upper()}", font_size=16, color=COL_TEXT,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        set_shape_name(option_text, f"nx_option_{letter}")

    # Correct marker bar (default sotto card A, runtime sposta)
    correct_bar = add_rect(
        shapes, start_x + Cm(0.5), start_y + card_h - Cm(0.4),
        card_w - Cm(1), Cm(0.3), COL_SECONDARY,
    )
    set_shape_name(correct_bar, "nx_correct_bar")

    # Correct marker cerchio ✓ (default su card A, runtime sposta)
    correct_marker = add_rect(
        shapes, start_x + card_w - Cm(1.2), start_y + Cm(0.4),
        Cm(0.8), Cm(0.8), COL_SECONDARY, shape_type=MSO_SHAPE.OVAL,
    )
    set_shape_name(correct_marker, "nx_correct_marker")
    set_text(
        correct_marker, "✓", font_size=14, bold=True, color=COL_WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # Logo
    add_logo(shapes, Cm(30), Cm(17.3), Cm(2.5), Cm(1.2), "nx_logo")


def build_layout_case_study(layout) -> None:
    """Layout 6: NX CASE_STUDY — 3 sezioni colorate."""
    shapes = layout.shapes

    # Banda top "CASO STUDIO"
    band = add_rect(shapes, 0, 0, SLIDE_W, Cm(1.0), COL_TEXT)
    set_shape_name(band, "nx_case_band")
    band_label = add_textbox(
        shapes, Cm(1), Cm(0), Cm(10), Cm(1.0),
        "CASO STUDIO", font_size=14, bold=True, color=COL_WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    set_shape_name(band_label, "nx_case_band_label")

    # Title (idx=0)
    title = add_textbox(
        shapes, Cm(1.5), Cm(1.5), Cm(31), Cm(1.8),
        "Titolo del caso studio", font_size=24, bold=True, color=COL_TEXT,
    )
    set_shape_name(title, "nx_title")
    inject_placeholder(title, "title", 0)

    # 3 sezioni in colonna (verticali, una sotto l'altra)
    sections = [
        ("situazione", "SITUAZIONE", COL_PRIMARY, COL_PINK_SOFT, "Descrizione della situazione iniziale"),
        ("azione", "AZIONE", COL_SECONDARY, COL_GREEN_SOFT, "Cosa si è fatto"),
        ("risultato", "RISULTATO", COL_AMBER, COL_AMBER_SOFT, "Quale risultato si è ottenuto"),
    ]

    sec_x = Cm(1.5)
    sec_y_start = Cm(4)
    sec_w = Cm(31)
    sec_h = Cm(4)
    sec_gap = Cm(0.3)

    for i, (key, label, color, bg, placeholder_txt) in enumerate(sections):
        ly = sec_y_start + Emu(int(sec_h) * i + int(sec_gap) * i)
        # Background section
        bg_shape = add_rect(
            shapes, sec_x, ly, sec_w, sec_h, bg,
            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
        )
        set_shape_name(bg_shape, f"nx_case_bg_{key}")

        # Barra verticale colorata
        bar = add_rect(shapes, sec_x, ly, Cm(0.3), sec_h, color)
        set_shape_name(bar, f"nx_case_bar_{key}")

        # Label
        lbl = add_textbox(
            shapes, sec_x + Cm(0.7), ly + Cm(0.3), Cm(5), Cm(0.7),
            label, font_size=12, bold=True, color=color,
        )
        set_shape_name(lbl, f"nx_case_label_{key}")

        # Testo sezione
        txt = add_textbox(
            shapes, sec_x + Cm(0.7), ly + Cm(1.0), sec_w - Cm(1.4), sec_h - Cm(1.2),
            placeholder_txt, font_size=14, color=COL_TEXT,
        )
        set_shape_name(txt, f"nx_{key}")

    # Ref + page
    ref = add_rect(
        shapes, Cm(1), Cm(17.5), Cm(10), Cm(1.0), COL_PINK_SOFT,
        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    set_shape_name(ref, "nx_ref")
    set_text(ref, "Riferimento", font_size=10, color=COL_PRIMARY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    page = add_textbox(
        shapes, Cm(24), Cm(17.5), Cm(5), Cm(1.0),
        "000 / 000", font_size=10, color=COL_TEXT,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
    )
    set_shape_name(page, "nx_page")

    add_logo(shapes, Cm(30), Cm(17.3), Cm(2.5), Cm(1.2), "nx_logo")


def build_layout_module_close(layout) -> None:
    """Layout 7: NX MODULE_CLOSE — chiusura modulo (NUOVO), 5 spunte ✓."""
    shapes = layout.shapes

    # Banda top verde (chiusura — diversa dal rosa dell'apertura)
    band = add_rect(shapes, 0, 0, SLIDE_W, Cm(1.0), COL_SECONDARY)
    set_shape_name(band, "nx_module_close_band")
    band_label = add_textbox(
        shapes, Cm(1), Cm(0), Cm(15), Cm(1.0),
        "RIEPILOGO MODULO", font_size=14, bold=True, color=COL_WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    set_shape_name(band_label, "nx_module_close_label")

    # Numero modulo (TITLE idx=0)
    num = add_textbox(
        shapes, Cm(1.5), Cm(1.8), Cm(31), Cm(2),
        "RIEPILOGO MODULO N", font_size=36, bold=True, color=COL_SECONDARY,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    set_shape_name(num, "nx_module_num")
    inject_placeholder(num, "title", 0)

    # Titolo modulo (BODY idx=1)
    title = add_textbox(
        shapes, Cm(1.5), Cm(4.2), Cm(31), Cm(1.5),
        "Sintesi del modulo", font_size=18, color=COL_TEXT,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    set_shape_name(title, "nx_module_title")
    inject_placeholder(title, "body", 1)

    # Recap body (BODY idx=2) — 5 spunte ✓
    body = add_textbox(
        shapes, Cm(3), Cm(7), Cm(27), Cm(9),
        "Concetto chiave 1", font_size=18, color=COL_TEXT,
    )
    set_shape_name(body, "nx_recap_body")
    inject_placeholder(body, "body", 2)
    inject_normautofit(body)
    inject_bullet_style(body, "✓", COL_SECONDARY)

    add_logo(shapes, Cm(30), Cm(17.3), Cm(2.5), Cm(1.2), "nx_logo")


def build_layout_recap(layout) -> None:
    """Layout 8: NX RECAP — riepilogo corso (non modulo)."""
    shapes = layout.shapes

    # Accent verticale
    accent = add_rect(shapes, 0, 0, Cm(0.2), SLIDE_H, COL_PRIMARY)
    set_shape_name(accent, "nx_accent_v")

    # Title
    title = add_textbox(
        shapes, Cm(1.5), Cm(1.5), Cm(31), Cm(2),
        "RIEPILOGO DEL CORSO", font_size=32, bold=True, color=COL_TEXT,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    set_shape_name(title, "nx_title")
    inject_placeholder(title, "title", 0)

    # Mini bar
    mini = add_rect(shapes, Cm(14.5), Cm(4), Cm(4), Cm(0.2), COL_PRIMARY)
    set_shape_name(mini, "nx_mini_bar")

    # Body con ✓
    body = add_textbox(
        shapes, Cm(3), Cm(5), Cm(27), Cm(11),
        "Concetto chiave 1", font_size=20, color=COL_TEXT,
    )
    set_shape_name(body, "nx_recap_body")
    inject_placeholder(body, "body", 1)
    inject_normautofit(body)
    inject_bullet_style(body, "✓", COL_SECONDARY)

    add_logo(shapes, Cm(30), Cm(17.3), Cm(2.5), Cm(1.2), "nx_logo")


def build_layout_closing(layout) -> None:
    """Layout 9: NX CLOSING — ultima slide, grazie + tagline."""
    shapes = layout.shapes

    # Background gradient (semplificato: rettangolo pieno rosa primario)
    bg = add_rect(shapes, 0, 0, SLIDE_W, SLIDE_H, COL_PRIMARY)
    set_shape_name(bg, "nx_bg_grad")

    # Banda decorativa verde top
    band = add_rect(shapes, 0, 0, SLIDE_W, Cm(0.8), COL_SECONDARY)
    set_shape_name(band, "nx_band_top")

    # Blob decorativi
    blob1 = add_rect(
        shapes, Cm(-3), Cm(13), Cm(10), Cm(10), COL_SECONDARY,
        shape_type=MSO_SHAPE.OVAL,
    )
    set_shape_name(blob1, "nx_blob")
    blob2 = add_rect(
        shapes, Cm(28), Cm(-3), Cm(8), Cm(8), COL_WHITE,
        shape_type=MSO_SHAPE.OVAL,
    )
    set_shape_name(blob2, "nx_blob_2")

    # Title (idx=0) — testo bianco su sfondo rosa
    title = add_textbox(
        shapes, Cm(2.5), Cm(5.5), Cm(28.87), Cm(3),
        "Grazie per l'attenzione", font_size=40, bold=True, color=COL_WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    set_shape_name(title, "nx_title")
    inject_placeholder(title, "title", 0)

    # Tagline
    tag = add_textbox(
        shapes, Cm(2.5), Cm(9), Cm(28.87), Cm(1.5),
        "C.F.P. Montessori — Formazione Globale",
        font_size=18, italic=True, color=COL_WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    set_shape_name(tag, "nx_tagline")

    # 3 puntini bianchi
    for i, x_cm in enumerate([14.5, 16.5, 18.5]):
        dot = add_rect(
            shapes, Cm(x_cm), Cm(12), Cm(0.4), Cm(0.4),
            COL_WHITE, shape_type=MSO_SHAPE.OVAL,
        )
        set_shape_name(dot, f"nx_dot{i+1}")

    # Logo grande centrato in basso
    add_logo(shapes, Cm(14), Cm(15), Cm(5.87), Cm(2.5), "nx_logo_large")


# ═══════════════════════════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════════════════════════


# Mapping: layout idx → (name, builder function)
LAYOUTS: list[tuple[str, callable]] = [
    ("NX TITLE",         build_layout_title),
    ("NX MODULE_OPEN",   build_layout_module_open),
    ("NX CONTENT_TEXT",  build_layout_content_text),
    ("NX CONTENT_IMAGE", build_layout_content_image),
    ("NX DIAGRAM",       build_layout_diagram),
    ("NX QUIZ",          build_layout_quiz),
    ("NX CASE_STUDY",    build_layout_case_study),
    ("NX MODULE_CLOSE",  build_layout_module_close),
    ("NX RECAP",         build_layout_recap),
    ("NX CLOSING",       build_layout_closing),
]


def _set_layout_name(layout, name: str) -> None:
    """Set the layout display name (visible in PowerPoint Slide Master view)."""
    layout.element.find(f"{P}cSld").set("name", name)


def _clear_default_placeholders(layout) -> None:
    """Remove the default Title/Content placeholders that come with a blank layout."""
    spTree = layout.shapes._spTree
    for sp in list(spTree.findall(f"{P}sp")):
        spTree.remove(sp)


def build() -> None:
    if not LOGO_PATH.exists():
        print(f"[WARN] Logo non trovato a {LOGO_PATH} — userò textbox 'CFP' come fallback")

    # Crea presentazione vuota
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Pulisce tutti i layout default (PowerPoint ne fornisce ~11 con presentazione vuota)
    # Lavoriamo solo sul primo master, e ricicleremo i suoi layout esistenti riempiendoli
    master = prs.slide_masters[0]
    existing_layouts = list(master.slide_layouts)
    n_existing = len(existing_layouts)
    print(f"Layout pre-esistenti nel master: {n_existing}")

    # Strategia: usiamo i layout esistenti (riempiendoli da zero), e se non bastano
    # aggiungiamo via XML. Default PowerPoint dà 11 layout, ne servono 10 → sovrabbondanti.
    if n_existing < len(LAYOUTS):
        raise RuntimeError(f"Servono {len(LAYOUTS)} layout, presenti solo {n_existing}")

    for idx, (name, builder) in enumerate(LAYOUTS):
        layout = existing_layouts[idx]
        # Svuota tutti gli shape pre-esistenti
        _clear_default_placeholders(layout)
        # Rinomina
        _set_layout_name(layout, name)
        # Disegna
        print(f"[{idx}] Building '{name}'...")
        builder(layout)

    # Elimina layout extra non usati (idx 10+)
    if n_existing > len(LAYOUTS):
        for extra_idx in range(len(LAYOUTS), n_existing):
            extra = existing_layouts[extra_idx]
            # Rimuovi dal cSld del master
            rel_id = master.element.find(f"{P}sldLayoutIdLst")
            if rel_id is not None:
                # Trova e rimuovi l'entry corrispondente
                for entry in list(rel_id):
                    # match per rId
                    rid = entry.get(qn("r:id"))
                    if rid and rid in master.part.rels:
                        target = master.part.rels[rid].target_part
                        if target == extra.part:
                            rel_id.remove(entry)
                            break
            print(f"[STRIP] removed extra layout idx={extra_idx}")

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"\nSaved: {OUT_PATH}")


def verify() -> None:
    """Re-open the file and print layout summary for verification."""
    prs = Presentation(str(OUT_PATH))
    print(f"\n=== VERIFY {OUT_PATH.name} ===")
    print(f"Slide size: {prs.slide_width/914400:.3f}\" x {prs.slide_height/914400:.3f}\"  "
          f"(expected 13.333 x 7.5)")
    print(f"Layouts: {len(prs.slide_layouts)}")
    for i, layout in enumerate(prs.slide_layouts):
        ph_count = len(layout.placeholders)
        sh_count = len(layout.shapes)
        print(f"  [{i}] {layout.name:20s} placeholders={ph_count} shapes={sh_count}")
        for ph in layout.placeholders:
            try:
                print(f"        - idx={ph.placeholder_format.idx} "
                      f"type={ph.placeholder_format.type} name={ph.name!r}")
            except Exception:
                pass


if __name__ == "__main__":
    build()
    verify()
