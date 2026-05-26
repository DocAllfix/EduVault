"""Fix sovrapposizioni testo nel template Claude Design nexus_master_v2.

Problema: il template ha 2 set di shape per ogni layout — i placeholder demo
del template originale + le shape con contenuto reale. Quando renderizzato,
si sovrappongono.

Fix: per ogni layout, identifico le shape che contengono testo placeholder
demo (es. "Titolo Slide", "Primo punto", "Sottotitolo / Riferimento", ecc.)
e le CANCELLO. Lascio solo le shape con il contenuto reale italiano sicurezza
(o le shape vuote che il backend riempirà).
"""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation


INPUT = Path("assets/templates/nexus_master_v2_claude_design.pptx")
OUTPUT = Path("assets/templates/nexus_master_v2.pptx")

# Frasi placeholder demo che Claude Design ha lasciato nei layout.
# Quando una shape ha testo che contiene una di queste, è da CANCELLARE.
PLACEHOLDER_PATTERNS = [
    "Titolo Slide",
    "Titolo del Corso",
    "Titolo del corso",
    "Sottotitolo",
    "Sottotitolo / Riferimento normativo",
    "Riferimento normativo",
    "Primo punto",
    "Secondo punto",
    "Terzo punto",
    "Quarto punto",
    "Verifica punto",
    "Quinto punto",
    "Sesto punto",
    # Nuovi pattern trovati in slide 3+ (numerazione "Punto uno", "Punto due"...)
    "Punto uno",
    "Punto due",
    "Punto tre",
    "Punto quattro",
    "Punto cinque",
    "Punto sei",
    "Testo del corpo",
    "Testo body",
    "Caption immagine",
    "Caption diagramma",
    "Didascalia",
    "Domanda quiz",
    "Opzione A",
    "Opzione B",
    "Opzione C",
    "Opzione D",
    "Risposta corretta:",
    "CASO STUDIO PLACEHOLDER",
    "RIEPILOGO PLACEHOLDER",
    "Module ref",
    "Modulo X",
    "X / Y",
    "Tagline",
    "Lorem ipsum",
    "Your text here",
    "Edit this",
    "[Type",
    # Slidesgo-like placeholder generici
    "Situazione iniziale",
    "Azione intrapresa",
    "Risultato ottenuto",
    "Punto chiave",
    "Concetto chiave",
    "Approfondimento",
    "Esempio normativo",
    "Citazione",
    "Grazie!",  # solo il "Grazie!" placeholder demo, il vero "Grazie" del closing è OK
    # Pattern trovati in slide 4, 6, 7 (run 2)
    "Titolo Diagramma",
    "Titolo del Caso Studio",
    "Titolo Caso Studio",
    "Caso Studio:",
    "Primo obiettivo",
    "Secondo concetto",
    "Terzo elemento",
    "Quarto risultato",
    "Quinta nozione",
    "Sesta competenza",
    "obiettivo punto",
    "concetto punto",
    "elemento punto",
    "risultato punto",
    "nozione punto",
    "competenza punto",
    "Primo obiettivo / concetto",
    "Secondo obiettivo / concetto",
    "Terzo obiettivo / concetto",
    # Fix #3 finale — pattern rimasti
    "Quinta nozione punto",
    "Quinta nozione",
    "Quinto risultato",
    "Quinto obiettivo",
    "Sesto risultato",
    "Modulo — Riferimento",
    "Modulo - Riferimento",
    "Figura — didascalia esplicativa",
    "Figura - didascalia esplicativa",
    "Figura - didascalia",
    "Figura — didascalia",
]


def has_placeholder_text(shape) -> bool:
    """True se la shape contiene testo placeholder demo da cancellare.

    Strategia: pattern match strict (exact match O prefix che inizia da inizio
    riga). Evita falsi positivi tipo "Primo Soccorso" che contiene "Primo".
    """
    if not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip()
    if not text:
        return False  # vuoto = OK
    text_lower = text.lower()
    for pattern in PLACEHOLDER_PATTERNS:
        p = pattern.lower()
        # Match esatto
        if text_lower == p:
            return True
        # Match start of line
        if text_lower.startswith(p + " ") or text_lower.startswith(p + "."):
            return True
        # Match isolato a un punto (es. "Quarto punto.")
        if text_lower == p + "." or text_lower == p + " punto.":
            return True
        # Match di una linea singola (multi-line: ogni linea inizia con il pattern)
        lines = [ln.strip().lower() for ln in text.split("\n") if ln.strip()]
        if lines and all(ln == p or ln.startswith(p + " ") or ln.startswith(p + ".") for ln in lines):
            return True
    return False


def fix_layout(layout, layout_name: str) -> int:
    """Svuota il testo placeholder demo SENZA cancellare le shape (così il
    backend python-pptx può scriverci dentro dopo).
    Restituisce numero shape svuotate."""
    cleared = 0
    for shape in layout.shapes:
        if has_placeholder_text(shape):
            try:
                text_sample = shape.text_frame.text.strip()[:60]
                # Svuoto il text frame (mantenendo first run per preservare font)
                shape.text_frame.text = ""
                cleared += 1
                print(f"    - Cleared: {text_sample!r} (shape '{shape.name}')")
            except Exception as e:
                print(f"    ! Error clearing: {e}")
    return cleared


def main():
    if not INPUT.is_file():
        raise SystemExit(f"Input not found: {INPUT}")

    # Backup
    backup = OUTPUT.with_suffix(".pptx.bak_before_overlap_fix")
    if OUTPUT.is_file() and not backup.is_file():
        shutil.copy2(OUTPUT, backup)
        print(f"Backup: {backup}")

    shutil.copy2(INPUT, OUTPUT)
    prs = Presentation(str(OUTPUT))
    print(f"Loaded {len(prs.slide_layouts)} layouts, {len(prs.slides)} slides")
    print()

    # Fix layouts
    print("=== FIX LAYOUTS ===")
    total_layout_fix = 0
    for i, layout in enumerate(prs.slide_layouts):
        print(f"Layout {i} '{layout.name}':")
        removed = fix_layout(layout, layout.name)
        total_layout_fix += removed
        if removed == 0:
            print(f"    (no placeholder demo found)")

    # Fix slides demo (anche le slide hanno gli stessi placeholder)
    print()
    print("=== FIX DEMO SLIDES (already filled with real content) ===")
    total_slide_fix = 0
    for i, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name
        print(f"Slide {i+1} (layout {layout_name}):")
        removed = fix_layout(slide, layout_name)
        total_slide_fix += removed
        if removed == 0:
            print(f"    (no placeholder demo found)")

    prs.save(str(OUTPUT))
    print()
    print(f"=== SUMMARY ===")
    print(f"Placeholder rimossi da layout: {total_layout_fix}")
    print(f"Placeholder rimossi da slide demo: {total_slide_fix}")
    print(f"Salvato: {OUTPUT}")


if __name__ == "__main__":
    main()
