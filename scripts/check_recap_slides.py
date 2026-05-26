"""Check RECAP slides in a PPTX: print text content of every shape on a
RECAP slide, so we can confirm FIX #24 (5-bullet split into nx_recap_text_710..750)
landed correctly.

Run from CONTAINER:
    docker compose exec -T backend python scripts/check_recap_slides.py \\
        /app/output/<course_id>.pptx [<slide_idx_1based>...]
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation


def dump_slide(slide, idx_1based: int) -> None:
    layout_name = slide.slide_layout.name
    print(f"\n{'=' * 78}", flush=True)
    print(f"SLIDE {idx_1based}  layout='{layout_name}'  shapes={len(slide.shapes)}", flush=True)
    print('=' * 78, flush=True)
    for shape in slide.shapes:
        name = shape.name or "<unnamed>"
        kind = str(shape.shape_type).split('.')[-1] if shape.shape_type else "?"
        text = ""
        if shape.has_text_frame:
            text = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
        flag = "  *FILLED*" if text else ""
        print(f"  [{kind:18s}] {name:30s}{flag}", flush=True)
        if text:
            for line in text.split("\n"):
                print(f"      | {line}", flush=True)


def main() -> None:
    if len(sys.argv) < 2:
        print("usage: check_recap_slides.py <pptx_path> [<slide_idx>...]", file=sys.stderr)
        sys.exit(2)

    pptx_path = Path(sys.argv[1])
    prs = Presentation(str(pptx_path))
    total = len(prs.slides)
    print(f"PPTX: {pptx_path.name}  total_slides={total}", flush=True)

    if len(sys.argv) > 2:
        targets = [int(x) for x in sys.argv[2:]]
    else:
        # auto: find first 3 RECAP layout slides
        targets = []
        for i, sl in enumerate(prs.slides, start=1):
            if "RECAP" in sl.slide_layout.name.upper():
                targets.append(i)
                if len(targets) >= 3:
                    break
        if not targets:
            print("no RECAP layout slides found", flush=True)
            sys.exit(1)
        print(f"auto-detected RECAP slide indexes: {targets}", flush=True)

    for idx in targets:
        if 1 <= idx <= total:
            dump_slide(prs.slides[idx - 1], idx)
        else:
            print(f"[skip] slide {idx} out of range (1..{total})", flush=True)


if __name__ == "__main__":
    main()
