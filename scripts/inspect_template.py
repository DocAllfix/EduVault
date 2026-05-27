"""Inspect nexus_master.pptx layout details for upgrade decisions."""
from __future__ import annotations
import sys
from pptx import Presentation
from pptx.util import Emu

prs = Presentation("assets/templates/nexus_master.pptx")
print(f"Layouts: {len(prs.slide_layouts)}")
print(f"Slide size: {prs.slide_width/914400:.1f} x {prs.slide_height/914400:.1f} inch")
print()
for i, layout in enumerate(prs.slide_layouts):
    print(f"=== Layout {i}: '{layout.name}' ===")
    for j, shape in enumerate(layout.shapes):
        try:
            t = shape.text[:80].replace("\n", " | ") if shape.has_text_frame else "(no text)"
            print(f"  [{j}] {shape.shape_type} '{shape.name}' "
                  f"({shape.left/914400:.1f},{shape.top/914400:.1f}) "
                  f"{shape.width/914400:.1f}x{shape.height/914400:.1f}: {t}")
        except Exception as e:
            print(f"  [{j}] ERR: {e}")
    print()
