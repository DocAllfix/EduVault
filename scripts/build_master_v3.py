"""Build nexus_master_v3.pptx from v2 — convert key text shapes into REAL placeholders.

FIX #28.2 (2026-05-26): la causa radice di metà dei nostri bug rendering (bullet non
ereditati, autofit in conflitto, testo placeholder del layout che trapela) è che
nexus_master_v2.pptx ha tutti AUTO_SHAPE/TEXT_BOX custom (zero placeholder PowerPoint
veri). v3 risolve convertendo SOLO i 4 shape testuali principali (nx_title, nx_body,
nx_image_box, nx_diagram_box) nei layout con body multi-bullet in VERI placeholder
PowerPoint:
  - nx_title → TITLE placeholder (idx=0)
  - nx_body → BODY placeholder (idx=1) con <a:normAutofit/> per autofit nativo
  - nx_image_box → PICTURE placeholder (CONTENT_IMAGE solo)
  - nx_diagram_box → PICTURE placeholder (DIAGRAM solo)

Layout interessati: CONTENT_TEXT (1), CONTENT_IMAGE (2), DIAGRAM (3).
Layout invariati (pattern "lista → N shape" che funziona):
  - TITLE (0), CLOSING (7): testo solo, placeholder semplice OK ma non critico
  - QUIZ (4): 4 opzioni in shape separate, custom necessario
  - CASE_STUDY (5): 3 sezioni fisiche separate (situazione/azione/risultato)
  - RECAP (6): 5 shape fisiche (nx_recap_text_710..750) col loro checkmark

Decorativi (nx_accent_v, nx_mini_bar, nx_ref, nx_page, nx_logo, nx_recap_band_label,
nx_case_band_label, le label sezione CASO STUDIO, i checkmark ✓ RECAP) → restano
AUTO_SHAPE/TEXT_BOX/PICTURE invariati.

Esecuzione (HOST, no container — usa solo python-pptx + lxml):
    python scripts/build_master_v3.py
Output: assets/templates/nexus_master_v3.pptx
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

SRC = Path("assets/templates/nexus_master_v2.pptx")
DST = Path("assets/templates/nexus_master_v3.pptx")

# Quali layout convertire e quali shape diventare placeholder.
# (layout_idx, {shape_name: (placeholder_type, placeholder_idx)})
# placeholder_type values: "title", "body", "pic"
LAYOUTS_TO_CONVERT: dict[int, dict[str, tuple[str, int]]] = {
    1: {  # NX CONTENT TEXT
        "nx_title": ("title", 0),
        "nx_body":  ("body",  1),
    },
    2: {  # NX CONTENT IMAGE
        "nx_title":      ("title", 0),
        "nx_body":       ("body",  1),
        "nx_image_box":  ("pic",   2),
    },
    3: {  # NX DIAGRAM
        "nx_title":        ("title", 0),
        "nx_body":         ("body",  1),
        "nx_diagram_box":  ("pic",   2),
    },
}

# Namespace shortcuts
A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"


def _inject_placeholder_props(sp_el, ph_type: str, ph_idx: int) -> None:
    """Iniettare <p:nvSpPr><p:nvPr><p:ph type="..." idx="..."/></p:nvPr></p:nvSpPr>
    nella shape, così PowerPoint la riconosce come VERO placeholder e i suoi figli
    (slide che usano questo layout) ereditano lo stile dal master bodyStyle/titleStyle.
    """
    nvSpPr = sp_el.find(f"{P}nvSpPr")
    if nvSpPr is None:
        return
    nvPr = nvSpPr.find(f"{P}nvPr")
    if nvPr is None:
        nvPr = etree.SubElement(nvSpPr, f"{P}nvPr")

    # Rimuovi <p:ph> precedente se presente (idempotente)
    for old in nvPr.findall(f"{P}ph"):
        nvPr.remove(old)

    # Mappa il nostro tipo logico al type attribute PowerPoint.
    # "title" → type="title" (può essere omesso, default è "body"; lo lasciamo esplicito)
    # "body"  → type="body"
    # "pic"   → type="pic"
    ph = etree.SubElement(nvPr, f"{P}ph")
    ph.set("type", ph_type)
    ph.set("idx", str(ph_idx))


def _inject_normautofit(sp_el) -> None:
    """Per i BODY placeholder: aggiungi <a:normAutofit/> nel <a:bodyPr> così
    PowerPoint riduce il font automaticamente all'apertura quando il testo straborda.
    Cancella eventuale <a:spAutoFit/> precedente. È l'autofit NATIVO — sostituisce
    l'euristica FIX #21 (run.font.size + MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE) che andava
    in conflitto e causava il "PowerPoint ripara all'apertura".
    """
    txBody = sp_el.find(f"{P}txBody")
    if txBody is None:
        return
    bodyPr = txBody.find(f"{A}bodyPr")
    if bodyPr is None:
        bodyPr = etree.Element(f"{A}bodyPr")
        txBody.insert(0, bodyPr)

    # Rimuovi vecchi autofit / spAutoFit
    for tag in ("normAutofit", "spAutoFit", "noAutofit"):
        for old in bodyPr.findall(f"{A}{tag}"):
            bodyPr.remove(old)

    etree.SubElement(bodyPr, f"{A}normAutofit")


def _ensure_body_list_style(sp_el) -> None:
    """Per BODY placeholder: assicura un <a:lstStyle> minimale con bullet level-0
    definito (buFont + buChar •), così quando add_paragraph() viene chiamato dal
    builder il glifo bullet appare anche se il master non eredita lvl1pPr.
    Idempotente: se già presente non duplica.
    """
    txBody = sp_el.find(f"{P}txBody")
    if txBody is None:
        return
    lstStyle = txBody.find(f"{A}lstStyle")
    if lstStyle is None:
        # Inserito subito dopo bodyPr
        bodyPr = txBody.find(f"{A}bodyPr")
        lstStyle = etree.Element(f"{A}lstStyle")
        if bodyPr is not None:
            bodyPr.addnext(lstStyle)
        else:
            txBody.insert(0, lstStyle)
    # Aggiungi lvl1pPr con buChar se non c'è
    if lstStyle.find(f"{A}lvl1pPr") is None:
        lvl1 = etree.SubElement(lstStyle, f"{A}lvl1pPr")
        lvl1.set("marL", "342900")
        lvl1.set("indent", "-342900")
        buFont = etree.SubElement(lvl1, f"{A}buFont")
        buFont.set("typeface", "Arial")
        buChar = etree.SubElement(lvl1, f"{A}buChar")
        buChar.set("char", "•")


def _clear_text_runs(sp_el) -> None:
    """Svuota i <a:t> della shape (rimuove il testo placeholder tipo 'Cosa hai imparato')
    lasciando però la struttura paragrafo/run perché lo stile vada in eredità.
    """
    txBody = sp_el.find(f"{P}txBody")
    if txBody is None:
        return
    for t in txBody.iter(f"{A}t"):
        t.text = ""


def convert(src: Path, dst: Path) -> None:
    prs = Presentation(str(src))
    converted_total = 0
    for layout_idx, mapping in LAYOUTS_TO_CONVERT.items():
        if layout_idx >= len(prs.slide_layouts):
            print(f"[SKIP] layout idx={layout_idx} fuori range")
            continue
        layout = prs.slide_layouts[layout_idx]
        spTree = layout.shapes._spTree  # type: ignore[attr-defined]
        for sp in spTree.findall(f"{P}sp"):
            nvSpPr = sp.find(f"{P}nvSpPr")
            if nvSpPr is None:
                continue
            cNvPr = nvSpPr.find(f"{P}cNvPr")
            name = (cNvPr.get("name") if cNvPr is not None else "") or ""
            if name not in mapping:
                continue
            ph_type, ph_idx = mapping[name]
            _inject_placeholder_props(sp, ph_type, ph_idx)
            _clear_text_runs(sp)
            if ph_type == "body":
                _inject_normautofit(sp)
                _ensure_body_list_style(sp)
            converted_total += 1
            print(f"[OK] layout {layout_idx} ({layout.name}): "
                  f"{name} -> {ph_type} placeholder (idx={ph_idx})")
        # Per i layout PICTURE: lo shape esiste già come AUTO_SHAPE; convertirlo a
        # vero <p:pic> placeholder richiederebbe sostituire l'XML intero. Per ora
        # marcarlo come placeholder TYPE='pic' è sufficiente — il builder lo
        # riconoscerà via slide.placeholders e farà insert_picture sul box.
    dst.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dst))
    print(f"\nSaved: {dst}  ({converted_total} shapes converted to placeholders)")


def verify(path: Path) -> None:
    """Apre il file risultante e stampa il count placeholder per ogni layout
    convertito. DEVE essere > 0 (vs 0 nel v2)."""
    prs = Presentation(str(path))
    print(f"\n=== VERIFY {path.name} ===")
    for idx in [0, 1, 2, 3, 4, 5, 6, 7]:
        if idx >= len(prs.slide_layouts):
            continue
        lay = prs.slide_layouts[idx]
        ph_count = len(lay.placeholders)
        marker = " <-- converted" if idx in LAYOUTS_TO_CONVERT else ""
        print(f"  layout {idx} ({lay.name}): {ph_count} placeholders{marker}")
        for ph in lay.placeholders:
            print(f"     - idx={ph.placeholder_format.idx} "
                  f"type={ph.placeholder_format.type} name={ph.name!r}")


if __name__ == "__main__":
    convert(SRC, DST)
    verify(DST)
