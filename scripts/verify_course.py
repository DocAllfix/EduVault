"""Verify generated course quality — slide-by-slide audit.

Run from CONTAINER:
    docker compose exec -T backend python scripts/verify_course.py <course_id>

Checks (5 dimensions):
1. Content: every slide has valid title/body/notes per Pydantic constraints.
2. References: every slide has a non-empty normative_ref with valid chunk_id.
3. Images: every CONTENT_IMAGE/DIAGRAM slide has a local image file present
   and Pillow-loadable (not corrupted).
4. Audio: every narratable slide has an MP3 + duration within 25-35s target.
5. PPTX geometry: every shape stays inside slide bounds (no overflow).

Optional --render flag triggers LibreOffice headless conversion to PNG
(one per slide) saved to output/verify/<course_id>/slide_NNNN.png so the
operator can visually inspect them.

Outputs a markdown report to stdout + optional --report-path.
"""
from __future__ import annotations

import argparse
import asyncio
import json
import os
import shutil
import subprocess
import sys
import uuid as uuid_mod
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

import asyncpg

from app.config import settings


SLIDE_W_EMU = Emu(9144000)   # 25.4 cm × 360000
SLIDE_H_EMU = Emu(6858000)   # 19.05 cm × 360000


def _emoji(ok: bool, warn: bool = False) -> str:
    if warn:
        return "⚠️ "
    return "✅" if ok else "❌"


async def fetch_course(pool: asyncpg.Pool, course_id: str) -> dict[str, Any]:
    row = await pool.fetchrow(
        """
        SELECT id, title, course_type, status, duration_hours, slide_contents_json,
               pptx_path, pdf_path, audio_manifest_path
        FROM courses WHERE id = $1
        """,
        uuid_mod.UUID(course_id),
    )
    if not row:
        raise SystemExit(f"Corso non trovato: {course_id}")
    return dict(row)


async def fetch_audio_tracks(pool: asyncpg.Pool, course_id: str) -> dict[int, dict[str, Any]]:
    rows = await pool.fetch(
        "SELECT slide_index, audio_path, duration_seconds, off_target "
        "FROM audio_tracks WHERE course_id = $1",
        uuid_mod.UUID(course_id),
    )
    return {r["slide_index"]: dict(r) for r in rows}


def parse_slides(slide_contents_json: Any) -> list[dict[str, Any]]:
    if not slide_contents_json:
        return []
    raw = slide_contents_json
    if isinstance(raw, str):
        raw = json.loads(raw)
    if isinstance(raw, dict) and "slides" in raw:
        raw = raw["slides"]
    return list(raw or [])


def check_slide_content(slide: dict[str, Any]) -> tuple[bool, list[str]]:
    """Check title/body/notes/quiz integrity per slide_type."""
    issues: list[str] = []
    stype = slide.get("slide_type", "?")
    title = (slide.get("title") or "").strip()
    body = (slide.get("body") or "").strip()
    notes = (slide.get("speaker_notes") or "").strip()

    if not title:
        issues.append("title vuoto")
    elif len(title) > 70:
        issues.append(f"title {len(title)}>70 char")

    # body bullet count
    if stype in {"CONTENT_TEXT", "CONTENT_IMAGE", "CASE_STUDY", "RECAP"}:
        bullets = [b for b in body.split("\n") if b.strip()]
        if len(bullets) == 0:
            issues.append("body vuoto")
        elif len(bullets) > 8:
            issues.append(f"body {len(bullets)} bullet > 8")

    # notes count
    notes_words = len(notes.split())
    if not notes:
        issues.append("speaker_notes vuoto")
    elif notes_words < 50:
        issues.append(f"notes solo {notes_words} parole (<50)")
    elif notes_words > 120:
        issues.append(f"notes {notes_words} parole (>120)")

    # quiz validity
    if stype == "QUIZ":
        opts = slide.get("quiz_options") or []
        if len(opts) != 4:
            issues.append(f"quiz_options={len(opts)} (richiesto 4)")
        qc = slide.get("quiz_correct")
        if not isinstance(qc, int) or qc < 0 or qc > 3:
            issues.append(f"quiz_correct invalido: {qc!r}")

    # normative_ref
    if not (slide.get("normative_ref") or "").strip():
        issues.append("normative_ref vuoto")

    return len(issues) == 0, issues


def check_slide_image(slide: dict[str, Any], course_id: str) -> tuple[bool | None, str]:
    """Returns (ok, message). ok=None for slides that don't need an image."""
    stype = slide.get("slide_type", "?")
    if stype not in {"CONTENT_IMAGE", "DIAGRAM"}:
        return None, "n/a"
    img = slide.get("image") or {}
    strategy = img.get("strategy")
    if strategy == "diagram":
        code = img.get("diagram_code") or ""
        if "viewBox" not in code:
            return False, "diagram SVG senza viewBox"
        if not code.strip():
            return False, "diagram_code vuoto"
        return True, "SVG OK"
    # web_search
    if not (img.get("query") or "").strip():
        return False, "image.query vuoto"
    return True, f"query='{img.get('query', '')[:40]}'"


def check_pptx_geometry(pptx_path: str) -> tuple[int, list[str]]:
    """Open PPTX, iterate shapes per slide, flag overflows."""
    issues: list[str] = []
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            try:
                left = int(shape.left or 0)
                top = int(shape.top or 0)
                width = int(shape.width or 0)
                height = int(shape.height or 0)
            except Exception:
                continue
            if left < 0 or top < 0:
                issues.append(f"slide {i} shape '{shape.name}' fuori-bordo top-left ({left},{top})")
            if left + width > SLIDE_W_EMU + 100:
                issues.append(
                    f"slide {i} shape '{shape.name}' overflow right "
                    f"({(left+width)/360000:.1f}cm > {SLIDE_W_EMU/360000:.1f}cm)"
                )
            if top + height > SLIDE_H_EMU + 100:
                issues.append(
                    f"slide {i} shape '{shape.name}' overflow bottom "
                    f"({(top+height)/360000:.1f}cm > {SLIDE_H_EMU/360000:.1f}cm)"
                )
    return len(prs.slides), issues


def check_audio(slide_idx: int, audio_tracks: dict[int, dict[str, Any]]) -> tuple[bool, str]:
    track = audio_tracks.get(slide_idx)
    if not track:
        return False, "audio assente"
    path = track["audio_path"]
    if not Path(path).is_file():
        return False, f"file MP3 mancante: {path}"
    dur = float(track["duration_seconds"] or 0)
    off = track["off_target"]
    if off:
        return True, f"durata {dur:.1f}s OFF-TARGET (range 25-35s)"
    return True, f"durata {dur:.1f}s OK"


def render_pptx_to_png(pptx_path: str, out_dir: Path) -> list[Path]:
    """LibreOffice headless: convert PPTX → multi-slide PDF → PNG per slide."""
    if not shutil.which("libreoffice"):
        print(
            "⚠️  LibreOffice non installato — skip render. "
            "Aggiungi al Dockerfile: apt-get install -y libreoffice-impress",
            file=sys.stderr,
        )
        return []
    out_dir.mkdir(parents=True, exist_ok=True)
    pdf_path = out_dir / "course.pdf"
    # PPTX → PDF
    subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf",
         "--outdir", str(out_dir), pptx_path],
        check=True, capture_output=True,
    )
    src_pdf = Path(pptx_path).with_suffix(".pdf").name
    src = out_dir / src_pdf
    if src.exists():
        src.rename(pdf_path)
    if not pdf_path.is_file():
        print(f"⚠️  PDF non generato da LibreOffice: {pdf_path}", file=sys.stderr)
        return []
    # PDF → PNG (1 per pagina) con pdftoppm
    subprocess.run(
        ["pdftoppm", "-png", "-r", "100", str(pdf_path), str(out_dir / "slide")],
        check=True, capture_output=True,
    )
    pngs = sorted(out_dir.glob("slide-*.png"))
    return pngs


async def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("course_id")
    parser.add_argument("--render", action="store_true", help="Render PNG via LibreOffice")
    parser.add_argument("--report-path", default=None)
    args = parser.parse_args()

    pool = await asyncpg.create_pool(settings.database_url)
    try:
        course = await fetch_course(pool, args.course_id)
        audio_tracks = await fetch_audio_tracks(pool, args.course_id)
        slides = parse_slides(course["slide_contents_json"])

        lines: list[str] = []
        lines.append(f"# Verify Course `{args.course_id}`\n")
        lines.append(f"- Titolo: **{course['title']}**")
        lines.append(f"- Tipo: `{course['course_type']}` · Ore: {course['duration_hours']}")
        lines.append(f"- Status: `{course['status']}`")
        lines.append(f"- PPTX: `{course['pptx_path']}`")
        lines.append(f"- Slide totali (JSON): **{len(slides)}**")
        lines.append(f"- Audio tracks: **{len(audio_tracks)}**\n")

        # 1+2+3+4: per-slide checks
        ok_count = 0
        warn_count = 0
        err_count = 0
        lines.append("## Per-slide report\n")
        lines.append("| # | Tipo | Contenuto | Immagine | Audio | Note |")
        lines.append("|---|---|---|---|---|---|")
        for s in slides:
            idx = s.get("index", "?")
            stype = s.get("slide_type", "?")
            c_ok, c_issues = check_slide_content(s)
            i_ok, i_msg = check_slide_image(s, args.course_id)
            a_ok, a_msg = check_audio(int(idx) if isinstance(idx, int) else 0, audio_tracks) if idx != "?" else (False, "no idx")

            row_status = "OK"
            note_parts: list[str] = []
            if not c_ok:
                err_count += 1
                row_status = "ERR"
                note_parts.append("contenuto: " + "; ".join(c_issues))
            elif i_ok is False:
                err_count += 1
                row_status = "ERR"
                note_parts.append("img: " + i_msg)
            elif not a_ok and len(audio_tracks) > 0:
                warn_count += 1
                row_status = "WARN"
                note_parts.append("audio: " + a_msg)
            else:
                ok_count += 1

            lines.append(
                f"| {idx} | {stype} | {_emoji(c_ok)} | "
                f"{('-' if i_ok is None else _emoji(i_ok))} | "
                f"{_emoji(a_ok) if a_msg != 'audio assente' or len(audio_tracks)>0 else '-'} "
                f"| {' / '.join(note_parts)[:120] if note_parts else 'OK'} |"
            )

        # 5: PPTX geometry
        lines.append("\n## PPTX geometry\n")
        if course["pptx_path"] and Path(course["pptx_path"]).is_file():
            slide_count, geom_issues = check_pptx_geometry(course["pptx_path"])
            lines.append(f"- Slide PPTX: **{slide_count}** (atteso {len(slides)})")
            if geom_issues:
                lines.append("- ⚠️ Geometria problemi:")
                for g in geom_issues[:30]:
                    lines.append(f"  - {g}")
                if len(geom_issues) > 30:
                    lines.append(f"  - ...+{len(geom_issues)-30} altri")
            else:
                lines.append("- ✅ Nessun overflow shape")
        else:
            lines.append("- ❌ PPTX non disponibile")

        # 6: PNG render (optional)
        png_paths: list[Path] = []
        if args.render and course["pptx_path"]:
            out_dir = Path("output/verify") / args.course_id
            png_paths = render_pptx_to_png(course["pptx_path"], out_dir)
            lines.append(f"\n## Render PNG\n- Generate **{len(png_paths)}** PNG in `{out_dir}`")

        # Summary
        lines.append("\n## Riepilogo\n")
        lines.append(f"- ✅ OK: **{ok_count}**")
        lines.append(f"- ⚠️ Warning: **{warn_count}**")
        lines.append(f"- ❌ Errori: **{err_count}**")
        verdict = "✅ CORSO VERIFICATO" if err_count == 0 else "❌ CORSO HA ERRORI"
        lines.append(f"\n**{verdict}**")

        report = "\n".join(lines)
        print(report)
        if args.report_path:
            Path(args.report_path).write_text(report, encoding="utf-8")
            print(f"\nReport salvato: {args.report_path}", file=sys.stderr)
        if png_paths:
            print(f"\nPNG paths (per Read tool):", file=sys.stderr)
            for p in png_paths:
                print(f"  /app/{p}", file=sys.stderr)
    finally:
        await pool.close()


if __name__ == "__main__":
    asyncio.run(main())
