"""Render i layout v2 (rebrandati) con contenuto demo."""
import subprocess
from pathlib import Path
from pptx import Presentation

TEMPLATE = Path("assets/templates/nexus_master_v2.pptx")
OUT_DIR = Path("output/verify/v2_preview")

DEMO_PER_LAYOUT = {
    "TITLE": "Primo Soccorso B/C",
    "CONTENT_TEXT": "Obblighi datore di lavoro:\n• Designare incaricati primo soccorso\n• Fornire formazione DM 388\n• Mettere a disposizione cassetta",
    "CONTENT_IMAGE": "I DPI sono l'ultima barriera\nDevono essere marcati CE\nIl lavoratore li deve usare",
    "DIAGRAM": "Flusso emergenza",
    "QUIZ": "Chi nomina il RSPP?",
    "CASE_STUDY": "Caduta da scala",
    "RECAP": "Cosa hai imparato",
    "CLOSING": "Grazie",
}


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    for f in OUT_DIR.glob("*"):
        f.unlink()

    prs = Presentation(str(TEMPLATE))
    # Cancella eventuali slide preesistenti
    for sld_id in list(prs.slides._sldIdLst):
        prs.slides._sldIdLst.remove(sld_id)

    # Aggiungi 1 slide per ogni layout interessante
    for layout in prs.slide_layouts:
        name = layout.name
        if name not in DEMO_PER_LAYOUT:
            continue
        slide = prs.slides.add_slide(layout)
        demo_text = DEMO_PER_LAYOUT[name]
        # Scrivo nel primo text frame
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == "" and not shape.text_frame.text:
                continue
        # Trovo il primo placeholder
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape.text_frame.text = demo_text
                break

    pptx_path = OUT_DIR / "v2_preview.pptx"
    prs.save(str(pptx_path))
    print(f"PPTX: {pptx_path}")
    print(f"Slides created: {len(prs.slides)}")


if __name__ == "__main__":
    main()
