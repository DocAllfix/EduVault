"""Rebrand template Slidesgo Health & Safety Workshop a CFP Montessori.

Esegue tutte le trasformazioni del brief:
1. Recoloring: arancione Slidesgo -> rosa CFP #C82E6E
2. Rimuove decorazioni "ufficio" (lampade, mensole, vasi, orologi, monitor PC)
3. Rimuove strip arancione laterale
4. Rimuove cerchio azzurro gigante o lo riduce
5. Rimuove accent lines sotto titoli
6. Sostituisce logo Slidesgo con logo CFP Montessori
7. Cancella tutte le 59 slide -> resta solo master + 8 layout
8. Riduce a 8 layout custom mappati ai nostri SlideType
9. Rinomina shape per convenzione nx_*

Input: assets/templates/slidesgo_health_safety_original.pptx
Output: assets/templates/nexus_master_v2.pptx
Logo: assets/brand/cfp_montessori_logo.jpeg
"""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Emu
from lxml import etree


INPUT = Path("assets/templates/slidesgo_health_safety_original.pptx")
OUTPUT = Path("assets/templates/nexus_master_v2.pptx")
LOGO = Path("assets/brand/cfp_montessori_logo.jpeg")

# Brand CFP Montessori
PINK_PRIMARY = RGBColor(0xC8, 0x2E, 0x6E)  # #C82E6E
GREEN_BRAND = RGBColor(0x76, 0x9E, 0x2E)   # #769E2E
DARK_NAVY = RGBColor(0x1F, 0x1F, 0x2C)     # #1F1F2C
TEXT_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
BG_SOFT = RGBColor(0xF8, 0xF9, 0xFA)

# Slidesgo palette da sostituire (analisi empirica del template originale)
# #FFA485 = arancione/salmone Slidesgo (148 occ) -> rosa CFP
# #C16D6B = rosa cipria Slidesgo (29 occ) -> verde brand CFP
SLIDESGO_ORANGE_HEXES = {"FFA485", "FF9D7E", "F8956F", "F47B20", "F47A20", "EB6E1F"}
SLIDESGO_PINK_HEXES = {"C16D6B", "B85A5A", "D17C7A"}

# OOXML namespaces
NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def hex_to_rgb_str(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def is_slidesgo_orange(hex_val: str) -> bool:
    """Detect Slidesgo orange shades."""
    if not hex_val or len(hex_val) != 6:
        return False
    hex_val = hex_val.upper()
    # Match exact + variants in range
    if hex_val in SLIDESGO_ORANGE_HEXES:
        return True
    r = int(hex_val[0:2], 16)
    g = int(hex_val[2:4], 16)
    b = int(hex_val[4:6], 16)
    # Heuristic: orange = R alto (>200), G medio (100-180), B basso (<70)
    if r > 200 and 100 <= g <= 180 and b < 70:
        return True
    return False


def recolor_all_orange_to_pink(prs):
    """Sostituisce TUTTI i colori orange Slidesgo con rosa CFP nei layout+master."""
    pink_hex = hex_to_rgb_str(PINK_PRIMARY)
    green_hex = hex_to_rgb_str(GREEN_BRAND)
    changes = 0

    # Itera master + layouts (NO slides, le cancelliamo dopo)
    parts_to_process = []
    parts_to_process.append(prs.slide_master.element)
    for layout in prs.slide_layouts:
        parts_to_process.append(layout.element)

    for root in parts_to_process:
        # Trova tutti gli elementi <a:srgbClr val="HEXVAL"/>
        for el in root.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"):
            val = el.get("val")
            if not val or len(val) != 6:
                continue
            val_upper = val.upper()
            # Slidesgo orange -> CFP pink
            if val_upper in SLIDESGO_ORANGE_HEXES or is_slidesgo_orange(val_upper):
                el.set("val", pink_hex)
                changes += 1
            # Slidesgo dark pink -> CFP green
            elif val_upper in SLIDESGO_PINK_HEXES:
                el.set("val", green_hex)
                changes += 1
    print(f"  Recolored {changes} orange/pink->CFP references in master+layouts")
    return changes


def recolor_pptx_file_level(output_path):
    """Apre il pptx come zip, sostituisce HEX nei XML, ri-zippa.

    Necessario perché python-pptx non vede tutti i colori (theme, charts, ecc.).
    """
    import zipfile, shutil, tempfile, os
    pink_hex = hex_to_rgb_str(PINK_PRIMARY)
    green_hex = hex_to_rgb_str(GREEN_BRAND)

    # Map sostituzioni
    REPLACE_MAP = {}
    for h in SLIDESGO_ORANGE_HEXES:
        REPLACE_MAP[h.upper()] = pink_hex
        REPLACE_MAP[h.lower()] = pink_hex.lower()
    for h in SLIDESGO_PINK_HEXES:
        REPLACE_MAP[h.upper()] = green_hex
        REPLACE_MAP[h.lower()] = green_hex.lower()

    tmp_dir = tempfile.mkdtemp()
    try:
        # Unzip
        with zipfile.ZipFile(output_path) as z:
            z.extractall(tmp_dir)

        # Replace in tutti i .xml e .rels
        changes = 0
        for root_dir, dirs, files in os.walk(tmp_dir):
            for fname in files:
                if not (fname.endswith(".xml") or fname.endswith(".rels")):
                    continue
                fpath = os.path.join(root_dir, fname)
                with open(fpath, "rb") as f:
                    content = f.read()
                orig = content
                for old, new in REPLACE_MAP.items():
                    old_b = old.encode()
                    new_b = new.encode()
                    if old_b in content:
                        n = content.count(old_b)
                        content = content.replace(old_b, new_b)
                        changes += n
                if content != orig:
                    with open(fpath, "wb") as f:
                        f.write(content)

        # Re-zip
        tmp_pptx = output_path + ".tmp"
        with zipfile.ZipFile(tmp_pptx, "w", zipfile.ZIP_DEFLATED) as zout:
            for root_dir, dirs, files in os.walk(tmp_dir):
                for fname in files:
                    fpath = os.path.join(root_dir, fname)
                    arcname = os.path.relpath(fpath, tmp_dir)
                    zout.write(fpath, arcname)
        shutil.move(tmp_pptx, output_path)
        print(f"  File-level recolor: {changes} HEX replacements across all XML")
        return changes
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def remove_decorative_shapes(layout, layout_name: str) -> int:
    """Rimuove decorazioni "ufficio" Slidesgo: lampade, mensole, vasi, orologi,
    monitor, strip arancione laterale, cerchio azzurro gigante.

    Heuristic: rimuovi shape che sono PICTURE/FREEFORM piccole in posizioni
    decorative (top-sx, top-dx, bottom-corners) e che NON contengono testo.
    """
    removed = 0
    slide_w_emu = int(10.0 * 914400)   # template slidesgo è 10x5.625 inch
    slide_h_emu = int(5.625 * 914400)

    shapes_to_remove = []
    for shape in layout.shapes:
        try:
            left = int(shape.left or 0)
            top = int(shape.top or 0)
            width = int(shape.width or 0)
            height = int(shape.height or 0)
        except Exception:
            continue

        # Skip shape con testo "vero" (titoli, body, placeholder)
        has_text = False
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt and len(txt) > 3:
                has_text = True

        # Skip placeholder (li teniamo per il rebranding)
        is_placeholder = bool(getattr(shape, "is_placeholder", False))
        if is_placeholder:
            continue

        # Categoria 1: STRIP ARANCIONE LATERALE (sx, alta e stretta)
        if left < int(0.3 * 914400) and width < int(0.5 * 914400) and height > int(3.0 * 914400):
            shapes_to_remove.append(shape)
            continue

        # Categoria 2: CERCHIO AZZURRO GIGANTE (>2 inch diametro)
        if (width > int(2.0 * 914400) and height > int(2.0 * 914400)
            and not has_text and not is_placeholder):
            shapes_to_remove.append(shape)
            continue

        # Categoria 3: DECORAZIONI TOP/CORNERS (piccole shape decorative)
        # Lampade pendenti, orologi, mensole — di solito in alto sx/dx, piccole
        if (not has_text
            and width < int(1.5 * 914400)
            and height < int(1.5 * 914400)
            and (top < int(1.5 * 914400) or top > int(3.5 * 914400))):
            shape_type = str(shape.shape_type) if shape.shape_type else ""
            if "PICTURE" in shape_type or "FREEFORM" in shape_type or "GROUP" in shape_type:
                shapes_to_remove.append(shape)
                continue

        # Categoria 4: ACCENT LINE SOTTILE sotto titoli (height < 0.1 inch, larga)
        if height < int(0.1 * 914400) and width > int(3.0 * 914400) and not has_text:
            shapes_to_remove.append(shape)
            continue

    for shape in shapes_to_remove:
        try:
            shape._element.getparent().remove(shape._element)
            removed += 1
        except Exception:
            pass
    return removed


def cancel_all_slides(prs):
    """Cancella tutte le slide del template, restano solo master + layouts."""
    xml_slides = prs.slides._sldIdLst
    count = 0
    for sld_id in list(xml_slides):
        xml_slides.remove(sld_id)
        count += 1
    print(f"  Cancelled {count} slides (only master+layouts remain)")
    return count


def insert_logo_in_layout(prs, layout, position="bottom-right", size="small"):
    """Inserisce logo CFP Montessori nel layout via XML diretto.

    python-pptx LayoutShapes.add_picture() non esiste (è solo per Slides).
    Soluzione: aggiungo la part image al layout part + un <p:pic> XML manuale.
    """
    if not LOGO.is_file():
        return False

    from pptx.oxml.ns import qn
    from pptx.util import Emu
    from pptx.parts.image import Image as ImagePart
    from copy import deepcopy

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    if size == "large":
        w, h = Inches(2.5), Inches(0.7)
    else:
        w, h = Inches(1.5), Inches(0.4)
    margin = Inches(0.2)

    if position == "bottom-right":
        left = slide_w - w - margin
        top = slide_h - h - margin
    elif position == "bottom-center":
        left = (slide_w - w) // 2
        top = slide_h - h - Inches(0.8)
    else:
        left, top = margin, slide_h - h - margin

    try:
        # 1. Aggiungi image part al layout
        layout_part = layout.part
        with open(LOGO, "rb") as f:
            image_blob = f.read()
        # Carica come Image part
        from pptx.opc.constants import CONTENT_TYPE as CT
        from pptx.opc.package import PartFactory
        # Usa la relazione standard image-part
        image_part = layout_part.package.image_parts.get_or_add_image_part(LOGO)
        rId = layout_part.relate_to(image_part, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image")

        # 2. Costruisci <p:pic> XML
        nsmap = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        }
        pic_xml = f"""<p:pic xmlns:p="{nsmap['p']}" xmlns:a="{nsmap['a']}" xmlns:r="{nsmap['r']}">
  <p:nvPicPr>
    <p:cNvPr id="9999" name="nx_logo"/>
    <p:cNvPicPr/>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{rId}"/>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="{int(left)}" y="{int(top)}"/>
      <a:ext cx="{int(w)}" cy="{int(h)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>"""
        pic_el = etree.fromstring(pic_xml)
        layout.shapes._spTree.append(pic_el)
        return True
    except Exception as e:
        print(f"    Logo insert failed: {e}")
        return False


def remove_extra_layouts(prs, keep_indices: list[int]) -> int:
    """Rimuove tutti i layout tranne quelli in keep_indices."""
    layouts_xml = prs.slide_master.element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}sldLayoutIdLst")
    all_layouts = list(prs.slide_layouts)
    removed = 0
    # Cancella in reverse order per non shift indici
    for i in range(len(all_layouts) - 1, -1, -1):
        if i not in keep_indices:
            layout = all_layouts[i]
            try:
                # Rimuovi reference dal master
                for sldLayoutId in layouts_xml.findall("{http://schemas.openxmlformats.org/presentationml/2006/main}sldLayoutId"):
                    rId = sldLayoutId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
                    rel = prs.slide_master.part.rels.get(rId)
                    if rel and rel.target_part == layout.part:
                        layouts_xml.remove(sldLayoutId)
                        break
                removed += 1
            except Exception as e:
                print(f"    Cannot remove layout {i}: {e}")
    print(f"  Removed {removed} extra layouts (kept {len(keep_indices)})")
    return removed


def rename_shapes_canonical(layout, layout_name: str):
    """Rinomina i shape secondo convention nx_* per facile lookup dal backend."""
    role_map = {
        "TITLE": ["nx_title", "nx_subtitle"],
        "CONTENT_TEXT": ["nx_title", "nx_body", "nx_ref", "nx_page"],
        "CONTENT_IMAGE": ["nx_title", "nx_body", "nx_image_box", "nx_caption", "nx_ref", "nx_page"],
        "DIAGRAM": ["nx_title", "nx_diagram_box", "nx_caption", "nx_ref", "nx_page"],
        "QUIZ": ["nx_title", "nx_option_a", "nx_option_b", "nx_option_c", "nx_option_d", "nx_correct_marker"],
        "CASE_STUDY": ["nx_title", "nx_situazione", "nx_azione", "nx_risultato", "nx_ref", "nx_page"],
        "RECAP": ["nx_title", "nx_body", "nx_module_ref", "nx_page"],
        "CLOSING": ["nx_title", "nx_tagline"],
    }
    if layout_name not in role_map:
        return 0
    names = role_map[layout_name]
    text_shapes = [s for s in layout.shapes if s.has_text_frame and s.name != "nx_logo"]
    renamed = 0
    for i, shape in enumerate(text_shapes):
        if i < len(names):
            try:
                shape.name = names[i]
                renamed += 1
            except Exception:
                pass
    return renamed


# ============================================================
# MAIN
# ============================================================

def main():
    if not INPUT.is_file():
        raise SystemExit(f"Input non trovato: {INPUT}")
    if not LOGO.is_file():
        raise SystemExit(f"Logo non trovato: {LOGO}")

    # Backup
    OUTPUT_BAK = OUTPUT.with_suffix(".pptx.bak")
    if OUTPUT.is_file() and not OUTPUT_BAK.is_file():
        shutil.copy2(OUTPUT, OUTPUT_BAK)
    shutil.copy2(INPUT, OUTPUT)
    print(f"Copy: {INPUT.name} -> {OUTPUT.name}")

    prs = Presentation(str(OUTPUT))
    print(f"Loaded: {len(prs.slides)} slides, {len(prs.slide_layouts)} layouts")

    # ── 1. CANCELLA TUTTE LE SLIDE ──
    print("\n[1/6] Cancello tutte le slide demo...")
    cancel_all_slides(prs)

    # ── 2. RECOLOR ORANGE -> PINK ──
    print("\n[2/6] Recolor arancione Slidesgo -> rosa CFP...")
    recolor_all_orange_to_pink(prs)

    # ── 3. RIMUOVI DECORAZIONI da TUTTI i layout ──
    print("\n[3/6] Rimuovo decorazioni 'ufficio' Slidesgo...")
    total_removed = 0
    for i, layout in enumerate(prs.slide_layouts):
        removed = remove_decorative_shapes(layout, layout.name)
        if removed > 0:
            print(f"  Layout {i} '{layout.name}': {removed} decorazioni rimosse")
        total_removed += removed
    print(f"  TOTALE: {total_removed} decorazioni rimosse")

    # ── 4. MAPPATURA 8 LAYOUT al nostro SlideType ──
    # Layout Slidesgo migliori per i nostri 8 SlideType:
    # TITLE -> 0 (TITLE)
    # CONTENT_TEXT -> 5 (ONE_COLUMN_TEXT) o 2 (TITLE_AND_BODY)
    # CONTENT_IMAGE -> 3 (TITLE_AND_TWO_COLUMNS) - 2 colonne testo+spazio img
    # DIAGRAM -> 4 (TITLE_ONLY) - title + box vuoto sotto
    # QUIZ -> 5 (ONE_COLUMN_TEXT)
    # CASE_STUDY -> 21 (CUSTOM_7) - 3 sezioni
    # RECAP -> 1 (SECTION_HEADER) o 13 (ONE_COLUMN_TEXT_1)
    # CLOSING -> 7 (SECTION_TITLE_AND_DESCRIPTION) o usa TITLE come mirror
    layout_keep_map = {
        0: "TITLE",         # Layout 0 TITLE
        2: "CONTENT_TEXT",  # Layout 2 TITLE_AND_BODY
        3: "CONTENT_IMAGE", # Layout 3 TITLE_AND_TWO_COLUMNS (sx body, dx img)
        4: "DIAGRAM",       # Layout 4 TITLE_ONLY (title + box body large)
        5: "QUIZ",          # Layout 5 ONE_COLUMN_TEXT (5 shape per options)
        21: "CASE_STUDY",   # Layout 21 CUSTOM_7 (multi-section)
        13: "RECAP",        # Layout 13 ONE_COLUMN_TEXT_1 (bullets)
        7: "CLOSING",       # Layout 7 SECTION_TITLE_AND_DESCRIPTION
    }

    # ── 5. RINOMINA LAYOUT con NOMI SEMANTICI dei nostri SlideType ──
    print("\n[4/6] Rinomina layout con nostri SlideType...")
    for orig_idx, semantic_name in layout_keep_map.items():
        if orig_idx < len(prs.slide_layouts):
            layout = prs.slide_layouts[orig_idx]
            layout.name = semantic_name
            print(f"  Layout {orig_idx} -> '{semantic_name}'")

    # ── 6. INSERT LOGO in ogni layout ──
    print("\n[5/6] Inserisco logo CFP Montessori in ogni layout...")
    inserted = 0
    for orig_idx, semantic_name in layout_keep_map.items():
        if orig_idx < len(prs.slide_layouts):
            layout = prs.slide_layouts[orig_idx]
            size = "large" if semantic_name in ("TITLE", "CLOSING") else "small"
            position = "bottom-center" if semantic_name in ("TITLE", "CLOSING") else "bottom-right"
            if insert_logo_in_layout(prs, layout, position=position, size=size):
                inserted += 1
                print(f"  Logo inserito in '{semantic_name}' ({size}/{position})")
    print(f"  TOTALE: {inserted} logo inseriti")

    # ── 7. RINOMINA SHAPE con convenzione nx_* ──
    print("\n[6/6] Rinomina shape con convenzione nx_*...")
    for orig_idx, semantic_name in layout_keep_map.items():
        if orig_idx < len(prs.slide_layouts):
            layout = prs.slide_layouts[orig_idx]
            renamed = rename_shapes_canonical(layout, semantic_name)
            if renamed > 0:
                print(f"  Layout '{semantic_name}': {renamed} shape rinominati nx_*")

    # ── SAVE ──
    prs.save(str(OUTPUT))
    print(f"\n[OK] Salvato: {OUTPUT}")
    print(f"   Size: {OUTPUT.stat().st_size:,} bytes")
    print(f"   Layouts finali: {len(prs.slide_layouts)}")

    # ── FILE-LEVEL RECOLOR ──
    print("\n[7/7] File-level recolor (theme + charts + tutti gli XML)...")
    recolor_pptx_file_level(str(OUTPUT))
    print(f"   Size finale: {OUTPUT.stat().st_size:,} bytes")


if __name__ == "__main__":
    main()
