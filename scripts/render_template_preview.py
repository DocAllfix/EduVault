"""Render 8 PNG di tutti i layout del template (con testo demo)
per anteprima visiva senza dover lanciare un corso intero.

Output: output/verify/template_preview/slide-001.png ... slide-008.png
"""
from __future__ import annotations
import shutil
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches

TEMPLATE = Path("assets/templates/nexus_master.pptx")
OUT_DIR = Path("output/verify/template_preview")

DEMO_CONTENT = {
    "TITLE": {
        "Text 1": "Primo Soccorso — Gruppi B e C",
        "Text 2": "Formazione 8 ore — D.Lgs 81/08 art. 45 + DM 388/2003",
    },
    "CONTENT_TEXT": {
        "Text 0": "Obblighi del datore di lavoro per il primo soccorso",
        "Text 2": "Designare un numero adeguato di lavoratori incaricati\nFornire la formazione obbligatoria prevista dal DM 388\nMettere a disposizione la cassetta di pronto soccorso conforme\nGarantire procedure di chiamata del 112 NUE\nValutare i rischi specifici aziendali nel DVR",
        "Text 3": "Art. 45 D.Lgs 81/08",
        "Text 4": "2 / 8",
    },
    "CONTENT_IMAGE": {
        "Text 0": "Dispositivi di Protezione Individuale",
        "Text 2": "I DPI sono l'ultima barriera difensiva\nDevono essere marcati CE\nIl lavoratore è obbligato a utilizzarli\nFormazione e addestramento obbligatori",
        "Text 4": "[ immagine: set DPI completo ]",
        "Text 5": "Art. 74-79 D.Lgs 81/08",
        "Text 6": "3 / 8",
    },
    "DIAGRAM": {
        "Text 0": "Flusso del primo soccorso aziendale",
        "Text 3": "Figura 1 — Sequenza operativa: allertamento → autoprotezione → valutazione → intervento",
        "Text 4": "DM 388/2003 art. 3",
        "Text 5": "4 / 8",
    },
    "QUIZ": {
        "Text 0": "Chi nomina il Responsabile del Servizio di Prevenzione e Protezione?",
        "Text 2": "A. Il Rappresentante dei Lavoratori per la Sicurezza (RLS)",
        "Text 3": "B. Il Datore di Lavoro, eventualmente assumendone direttamente i compiti",
        "Text 4": "C. L'Ispettorato Territoriale del Lavoro",
        "Text 5": "D. Il Medico Competente in accordo con il CdA",
        "Text 7": "Risposta corretta: B",
    },
    "CASE_STUDY": {
        "Text 1": "CASO STUDIO",
        "Text 2": "Infortunio in magazzino: caduta da scala portatile",
        "Text 5": "SITUAZIONE",
        "Text 6": "Un magazziniere utilizza una scala doppia non vincolata per prelevare materiale da uno scaffale alto. La scala oscilla, l'operatore cade da 2.5m riportando frattura del polso.",
        "Text 9": "AZIONE",
        "Text 10": "L'azienda aggiorna il DVR, sostituisce le scale con trabattelli certificati EN 1004 e organizza una formazione specifica sul lavoro in quota.",
        "Text 13": "RISULTATO",
        "Text 14": "Nessun infortunio analogo nei 24 mesi successivi. L'audit INAIL riconosce la corretta applicazione delle misure e riduce il premio assicurativo (OT23).",
        "Text 15": "Titolo IV — Capo II, D.Lgs 81/08",
        "Text 16": "6 / 8",
    },
    "RECAP": {
        "Text 1": "RIEPILOGO",
        "Text 2": "Cosa hai imparato in questo modulo",
        "Text 3": "I principi cardine del D.Lgs 81/08 e la gerarchia delle misure di prevenzione\nLe responsabilità del datore di lavoro, RSPP, medico competente e RLS\nLe modalità di valutazione dei rischi e redazione del DVR\nGli obblighi formativi previsti dall'Accordo Stato-Regioni 2011",
        "Text 4": "Modulo 1 — Formazione Generale",
        "Text 5": "7 / 8",
    },
    "CLOSING": {
        "Text 0": "Grazie per l'attenzione",
        "Text 1": "Formazione C.F.P. Montessori",
    },
}


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    # Pulisci output precedenti
    for f in OUT_DIR.glob("*"):
        f.unlink()

    prs = Presentation(str(TEMPLATE))
    # Cancella eventuali slide preesistenti
    for sld_id in list(prs.slides._sldIdLst):
        prs.slides._sldIdLst.remove(sld_id)

    # Per ogni layout, crea una slide con contenuto demo
    for layout in prs.slide_layouts:
        layout_name = layout.name
        if layout_name not in DEMO_CONTENT:
            continue
        slide = prs.slides.add_slide(layout)
        demo = DEMO_CONTENT[layout_name]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name in demo:
                shape.text_frame.text = demo[shape.name]

    pptx_path = OUT_DIR / "template_preview.pptx"
    prs.save(str(pptx_path))
    print(f"PPTX saved: {pptx_path}")

    # Converti in PDF + PNG via LibreOffice + pdftoppm
    subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf",
         "--outdir", str(OUT_DIR), str(pptx_path)],
        check=True, capture_output=True,
    )
    pdf_path = OUT_DIR / "template_preview.pdf"
    subprocess.run(
        ["pdftoppm", "-png", "-r", "100", str(pdf_path), str(OUT_DIR / "slide")],
        check=True, capture_output=True,
    )
    pngs = sorted(OUT_DIR.glob("slide-*.png"))
    print(f"Rendered {len(pngs)} PNG:")
    for p in pngs:
        print(f"  {p}")


if __name__ == "__main__":
    main()
