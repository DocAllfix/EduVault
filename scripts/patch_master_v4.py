"""Patch idempotente per nexus_master_v4_bulletfix.pptx (FIX #30.0-sexies).

Risponde a P2 dell'analista (2026-05-26): centratura verticale del body +
aggiunta nx_page al layout RECAP che ne è privo (causa warning shape_missing).

Caratteristiche:
- Idempotente: rilanciabile più volte senza danni (skip se già applicato).
- Output `nexus_master_v4_patched.pptx` (NON sovrascrive il bulletfix).
- Funzione `set_bullet_liststyle` mantenuta in app/builders/template_ops.py
  per riuso (analista: "una funzione, un posto, importata da chi serve").

Patches applicate:
1. anchor="ctr" su <a:bodyPr> dei placeholder body (nx_body su CONTENT_TEXT/
   CONTENT_IMAGE, nx_recap_body su RECAP/MODULE_CLOSE).
2. Aggiunta shape AUTO_SHAPE `nx_page` al layout RECAP (idx 8) — non c'era
   nel v4_bulletfix.

Esecuzione (host):
    python scripts/patch_master_v4.py

Verifica round-trip dopo l'apertura in PowerPoint che il file regga
(no dialog "documento riparato").
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

SRC = Path("assets/templates/nexus_master_v4_bulletfix.pptx")
DST = Path("assets/templates/nexus_master_v4_patched.pptx")

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

# Shape che ricevono anchor="ctr" (analista P2: body always centered, min-4
# bullet garantito dal validator pydantic → varianza altezza contenuta).
BODY_PLACEHOLDERS_TO_CENTER = {
    "NX CONTENT_TEXT":  ["nx_body"],
    "NX CONTENT_IMAGE": ["nx_body"],
    "NX RECAP":         ["nx_recap_body"],
    "NX MODULE_CLOSE":  ["nx_recap_body"],
}


def patch_anchor_center(sp_el) -> bool:
    """Set anchor='ctr' on bodyPr. Returns True if modified, False if already ctr."""
    txBody = sp_el.find(f"{P}txBody")
    if txBody is None:
        return False
    bodyPr = txBody.find(f"{A}bodyPr")
    if bodyPr is None:
        bodyPr = etree.SubElement(txBody, f"{A}bodyPr")
        txBody.insert(0, bodyPr)
    current = bodyPr.get("anchor")
    if current == "ctr":
        return False
    bodyPr.set("anchor", "ctr")
    return True


def add_nx_page_to_recap(layout) -> bool:
    """Aggiungi nx_page (casella testo '000 / 000') al layout RECAP se assente.

    LayoutShapes non ha add_textbox (a differenza di Slide.shapes), quindi
    costruiamo lo shape direttamente via XML e lo appendiamo al <p:spTree>
    del layout. Coordinate footer destro: X=27cm, Y=17.5cm su slide 33.87×19.05cm.
    """
    # Skip se nx_page già presente
    for sh in layout.shapes:
        if sh.name == "nx_page":
            return False

    # EMU positioning (1 cm = 360000 EMU)
    left_emu = int(27.0 * 360000)
    top_emu = int(17.5 * 360000)
    width_emu = int(5.0 * 360000)
    height_emu = int(0.8 * 360000)

    # Build <p:sp> XML directly (textbox-style)
    sp_xml = f"""<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <p:nvSpPr>
            <p:cNvPr id="999" name="nx_page"/>
            <p:cNvSpPr txBox="1"/>
            <p:nvPr/>
        </p:nvSpPr>
        <p:spPr>
            <a:xfrm>
                <a:off x="{left_emu}" y="{top_emu}"/>
                <a:ext cx="{width_emu}" cy="{height_emu}"/>
            </a:xfrm>
            <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
            <a:noFill/>
        </p:spPr>
        <p:txBody>
            <a:bodyPr wrap="square" rtlCol="0" anchor="ctr"/>
            <a:lstStyle/>
            <a:p>
                <a:pPr algn="r"/>
                <a:r>
                    <a:rPr lang="it-IT" sz="1000">
                        <a:solidFill><a:srgbClr val="1F1F1F"/></a:solidFill>
                        <a:latin typeface="Inter"/>
                    </a:rPr>
                    <a:t>000 / 000</a:t>
                </a:r>
            </a:p>
        </p:txBody>
    </p:sp>"""
    sp_el = etree.fromstring(sp_xml)
    spTree = layout.shapes._spTree
    spTree.append(sp_el)
    return True


def patch(src: Path, dst: Path) -> None:
    if not src.exists():
        raise FileNotFoundError(f"SRC mancante: {src}")
    prs = Presentation(str(src))
    print(f"Loaded: {src}")
    print(f"  Slide size: {prs.slide_width/914400:.3f}\" x {prs.slide_height/914400:.3f}\"")
    print(f"  Layouts: {len(prs.slide_layouts)}")

    n_anchors_patched = 0
    n_anchors_already = 0
    n_recap_page_added = 0

    for layout in prs.slide_layouts:
        layout_name = layout.name
        # P2.1: anchor=ctr
        if layout_name in BODY_PLACEHOLDERS_TO_CENTER:
            target_names = BODY_PLACEHOLDERS_TO_CENTER[layout_name]
            for sh in layout.shapes:
                if sh.name in target_names:
                    if patch_anchor_center(sh._element):
                        n_anchors_patched += 1
                        print(f"  [{layout_name}] {sh.name}: anchor → ctr")
                    else:
                        n_anchors_already += 1
                        print(f"  [{layout_name}] {sh.name}: anchor già ctr (skip)")
        # P2.2: nx_page al layout RECAP
        if layout_name == "NX RECAP":
            if add_nx_page_to_recap(layout):
                n_recap_page_added += 1
                print(f"  [{layout_name}] added nx_page footer textbox")
            else:
                print(f"  [{layout_name}] nx_page già presente (skip)")

    dst.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dst))
    print()
    print(f"Saved: {dst}")
    print(f"  Anchors patched (→ ctr): {n_anchors_patched}")
    print(f"  Anchors already ctr:     {n_anchors_already}")
    print(f"  RECAP nx_page added:     {n_recap_page_added}")


def verify(path: Path) -> None:
    """Verify anchor + nx_page presence."""
    prs = Presentation(str(path))
    print()
    print(f"=== VERIFY {path.name} ===")
    for layout in prs.slide_layouts:
        if layout.name in BODY_PLACEHOLDERS_TO_CENTER:
            for sh in layout.shapes:
                if sh.name in BODY_PLACEHOLDERS_TO_CENTER[layout.name]:
                    txBody = sh._element.find(f"{P}txBody")
                    bodyPr = txBody.find(f"{A}bodyPr") if txBody is not None else None
                    anchor = bodyPr.get("anchor") if bodyPr is not None else None
                    status = "OK" if anchor == "ctr" else "FAIL"
                    print(f"  [{layout.name}] {sh.name}: anchor={anchor!r} [{status}]")
        if layout.name == "NX RECAP":
            has_page = any(sh.name == "nx_page" for sh in layout.shapes)
            status = "OK" if has_page else "FAIL"
            print(f"  [{layout.name}] nx_page present: {has_page} [{status}]")


if __name__ == "__main__":
    patch(SRC, DST)
    verify(DST)
