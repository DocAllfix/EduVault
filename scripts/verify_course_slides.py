"""Verify OGNI slide del PPTX corso. Genera report markdown completo.

Uso:
    docker compose exec -T backend python /app/scripts/verify_course_slides.py \\
        /app/output/<course_id>.pptx
"""
from __future__ import annotations
import sys
from pathlib import Path
from collections import Counter

from pptx import Presentation
from pptx.util import Emu


def shape_bounds(shape):
    """Restituisce (left, top, right, bottom) in EMU."""
    try:
        l = int(shape.left or 0)
        t = int(shape.top or 0)
        w = int(shape.width or 0)
        h = int(shape.height or 0)
        return (l, t, l + w, t + h)
    except Exception:
        return None


def boxes_overlap(b1, b2, tolerance_emu=50000):
    """Due bounding box si sovrappongono se hanno area in comune."""
    if not b1 or not b2:
        return False
    l1, t1, r1, bo1 = b1
    l2, t2, r2, bo2 = b2
    # Tolleranza: piccole sovrapposizioni decorative sono OK
    if r1 - tolerance_emu <= l2 or r2 - tolerance_emu <= l1:
        return False
    if bo1 - tolerance_emu <= t2 or bo2 - tolerance_emu <= t1:
        return False
    return True


def has_logo_in_layout(layout):
    """Verifica se il layout ha il logo (ereditato dalla slide)."""
    for shape in layout.shapes:
        if shape.name == "nx_logo" or ("PICTURE" in str(shape.shape_type) and "logo" in (shape.name or "").lower()):
            return True
    return False


def analyze_slide(slide, slide_idx, layout_name, layout_has_logo=False):
    """Analizza una slide. Restituisce dict con metriche."""
    metrics = {
        "idx": slide_idx,
        "layout": layout_name,
        "shapes_total": len(slide.shapes),
        "text_shapes": 0,
        "text_chars_total": 0,
        "pictures": 0,
        "pictures_non_logo": 0,
        "has_logo": layout_has_logo,
        "has_page_num": False,
        "has_normative_ref": False,
        "has_title": False,
        "has_body": False,
        "has_image_box": False,
        "has_diagram_box": False,
        "text_overlaps": 0,
        "shapes_outside_slide": 0,
        "issues": [],
    }

    # Slide dimensions (20×11.25 inch standard)
    SLIDE_W = int(20.0 * 914400)
    SLIDE_H = int(11.25 * 914400)

    text_shape_bounds = []  # per detection sovrapposizioni testo

    for shape in slide.shapes:
        name = shape.name or ""
        shape_type = str(shape.shape_type) if shape.shape_type else ""
        bounds = shape_bounds(shape)

        # Verifica out-of-slide
        if bounds:
            l, t, r, bo = bounds
            if l < -50000 or t < -50000 or r > SLIDE_W + 50000 or bo > SLIDE_H + 50000:
                metrics["shapes_outside_slide"] += 1
                metrics["issues"].append(f"shape '{name}' out of slide bounds")

        # Pictures
        if "PICTURE" in shape_type:
            metrics["pictures"] += 1
            if name == "nx_logo":
                metrics["has_logo"] = True
            else:
                metrics["pictures_non_logo"] += 1  # immagine reale (Pexels/diagram)

        # PLACEHOLDER DETECTION: shape ancora con testo placeholder
        if shape.has_text_frame:
            txt_low = shape.text_frame.text.strip().lower()
            placeholder_markers = [
                "[ auto_shape", "[auto_shape",
                "immagine inserita dinamicamente",
                "diagramma svg inserito dinamicamente",
                "[ cassetta", "[ immagine", "[ cartello", "[ diagramma",
                "[ foto", "[ illustrazione",
            ]
            if any(m in txt_low for m in placeholder_markers):
                metrics["issues"].append(
                    f"PLACEHOLDER residuo in '{name}': {txt_low[:60]}"
                )

        # Text shapes
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                metrics["text_shapes"] += 1
                metrics["text_chars_total"] += len(text)
                if bounds:
                    text_shape_bounds.append((name, bounds, text[:30]))

                # Riconoscimento ruoli per name
                if name == "nx_title":
                    metrics["has_title"] = True
                elif name == "nx_body":
                    metrics["has_body"] = True
                elif name == "nx_page":
                    metrics["has_page_num"] = True
                elif name == "nx_ref":
                    metrics["has_normative_ref"] = True
                elif name in ("nx_image_box",):
                    metrics["has_image_box"] = True
                elif name in ("nx_diagram_box",):
                    metrics["has_diagram_box"] = True

    # Detection sovrapposizioni testo significative
    for i, (n1, b1, t1) in enumerate(text_shape_bounds):
        for j, (n2, b2, t2) in enumerate(text_shape_bounds):
            if i >= j:
                continue
            # Skip se uno è footer/page-num (di solito stanno separati)
            if any(x in (n1, n2) for x in ("nx_page", "nx_ref", "nx_band_top",
                                            "nx_case_band_label",
                                            "nx_recap_band_label",
                                            "nx_logo")):
                continue
            if boxes_overlap(b1, b2):
                metrics["text_overlaps"] += 1
                metrics["issues"].append(
                    f"overlap: '{n1}'({t1!r}) <-> '{n2}'({t2!r})"
                )

    return metrics


def main(pptx_path):
    prs = Presentation(pptx_path)
    print(f"# Analisi {pptx_path}")
    print(f"\nSlide totali: **{len(prs.slides)}**\n")

    all_metrics = []
    layout_counts = Counter()
    issue_counts = Counter()
    no_text_count = 0
    no_logo_count = 0
    overlap_count = 0
    outside_count = 0

    # Pre-check: quali layout hanno il logo?
    layout_logo_map = {layout.name: has_logo_in_layout(layout) for layout in prs.slide_layouts}

    for i, slide in enumerate(prs.slides, start=1):
        layout_name = slide.slide_layout.name
        layout_counts[layout_name] += 1
        layout_has_logo = layout_logo_map.get(layout_name, False)
        m = analyze_slide(slide, i, layout_name, layout_has_logo)
        all_metrics.append(m)

        if m["text_chars_total"] == 0:
            no_text_count += 1
        if not m["has_logo"]:
            no_logo_count += 1
        if m["text_overlaps"] > 0:
            overlap_count += 1
        if m["shapes_outside_slide"] > 0:
            outside_count += 1
        for issue in m["issues"]:
            issue_counts[issue[:50]] += 1

    print("## Riepilogo aggregato\n")
    print(f"- Slide senza testo: **{no_text_count}** / {len(prs.slides)}")
    print(f"- Slide senza logo: **{no_logo_count}** / {len(prs.slides)}")
    print(f"- Slide con sovrapposizioni testo: **{overlap_count}** / {len(prs.slides)}")
    print(f"- Slide con shape fuori bordi: **{outside_count}** / {len(prs.slides)}")
    print()

    print("## Distribuzione layout\n")
    for layout, n in sorted(layout_counts.items(), key=lambda x: -x[1]):
        print(f"- `{layout}`: {n}")
    print()

    # Slide problematiche
    problematic = [m for m in all_metrics if m["text_overlaps"] > 0
                   or m["shapes_outside_slide"] > 0
                   or m["text_chars_total"] == 0
                   or not m["has_logo"]]
    if problematic:
        print(f"## Slide problematiche ({len(problematic)})\n")
        for m in problematic[:20]:  # max 20 per non sporcare report
            print(f"### Slide {m['idx']} ({m['layout']})")
            print(f"- Text shapes: {m['text_shapes']}, chars: {m['text_chars_total']}")
            print(f"- Logo: {'OK' if m['has_logo'] else 'MANCA'}")
            print(f"- Overlaps: {m['text_overlaps']}")
            print(f"- Outside: {m['shapes_outside_slide']}")
            if m["issues"]:
                print(f"- Issues:")
                for issue in m["issues"][:5]:
                    print(f"    - {issue}")
            print()
        if len(problematic) > 20:
            print(f"\n... e altre {len(problematic)-20} slide problematiche")
    else:
        print("\n## TUTTE LE SLIDE SONO OK\n")

    # Issue patterns
    if issue_counts:
        print("\n## Top pattern issue\n")
        for issue, n in issue_counts.most_common(10):
            print(f"- ({n}) {issue}")

    # Check immagini/diagrammi
    print("\n## Immagini reali (Pexels/diagrammi PNG)\n")
    content_image_slides = [m for m in all_metrics if "CONTENT IMAGE" in m["layout"]]
    diagram_slides = [m for m in all_metrics if "DIAGRAM" in m["layout"]]
    img_with_pic = [m for m in content_image_slides if m["pictures_non_logo"] > 0]
    img_without = [m for m in content_image_slides if m["pictures_non_logo"] == 0]
    diag_with_pic = [m for m in diagram_slides if m["pictures_non_logo"] > 0]
    diag_without = [m for m in diagram_slides if m["pictures_non_logo"] == 0]

    print(f"- CONTENT_IMAGE con immagine reale: **{len(img_with_pic)}/{len(content_image_slides)}** ({100*len(img_with_pic)/max(len(content_image_slides),1):.0f}%)")
    print(f"- DIAGRAM con PNG reso: **{len(diag_with_pic)}/{len(diagram_slides)}** ({100*len(diag_with_pic)/max(len(diagram_slides),1):.0f}%)")

    placeholder_count = sum(1 for m in all_metrics for issue in m["issues"] if "PLACEHOLDER" in issue)
    print(f"- Slide con testo PLACEHOLDER residuo: **{placeholder_count}/{len(prs.slides)}**")

    if img_without:
        print(f"\n### CONTENT_IMAGE senza immagine ({len(img_without)} — Pexels/Openverse falliti):")
        for m in img_without[:10]:
            print(f"  - Slide {m['idx']}")
        if len(img_without) > 10:
            print(f"  ... e altre {len(img_without)-10}")

    if diag_without:
        print(f"\n### DIAGRAM senza PNG ({len(diag_without)} — cairosvg falliti):")
        for m in diag_without[:10]:
            print(f"  - Slide {m['idx']}")
        if len(diag_without) > 10:
            print(f"  ... e altre {len(diag_without)-10}")


if __name__ == "__main__":
    main(sys.argv[1] if len(sys.argv) > 1 else "/app/output/test.pptx")
