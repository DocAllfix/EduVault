"""Tune nexus_master_v3.pptx → nexus_master_v3_tuned.pptx — FIX #30.0-quater.

OBIETTIVO PRIMARIO: portare il template alla DIMENSIONE PPTX STANDARD
13.33×7.5" (vs 20×11.25" attuale). Tutte le coordinate degli shape sono
scalate proporzionalmente (×0.6667). Inoltre:

1. **Slide size** 18288000×10287000 → 12192000×6858000 EMU (=13.33×7.5").
2. **TUTTE le coordinate** (x/y/width/height) di OGNI shape × 0.6667 per
   mantenere posizioni relative (NO distorsione visiva).
3. **Font Inter** (era DM Sans) su tutti i placeholder + theme font scheme.
4. **Body box** anchor=MIDDLE (era TOP) sui layout CONTENT_TEXT /
   CONTENT_IMAGE / DIAGRAM → risolve "vuoto in basso".
5. **Title** allineato CENTER orizzontalmente per coerenza tipografica.

Output: assets/templates/nexus_master_v3_tuned.pptx (sicuro, non sovrascrive).
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SRC = Path("assets/templates/nexus_master_v3.pptx")
DST = Path("assets/templates/nexus_master_v3_tuned.pptx")

FONT_FAMILY = "Inter"

# Scaling factor: 13.33/20 = 0.6667 (anche 7.5/11.25 = 0.6667 → uniforme)
SCALE = 13.333 / 20.0

# Target slide size (PPTX widescreen standard Microsoft)
TARGET_W_EMU = int(13.333 * 914400)  # 12192000
TARGET_H_EMU = int(7.5 * 914400)     # 6858000

# XML namespaces
A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

# Override geometrici POST-SCALING (per layout dove vogliamo migliorare di più che
# solo scalare). Formato: layout_idx -> {shape_name: (left, top, width, height) IN INCHES POST-SCALE}.
# Lasciato vuoto: lo scaling proporzionale (×0.6667) basta per il primo giro.
LAYOUT_GEOMETRY_OVERRIDE: dict[int, dict[str, tuple[float, float, float, float]]] = {}

# Shape che ricevono anchor=MIDDLE (body centrato verticalmente — il fix principale)
SHAPES_TO_CENTER_VERTICALLY = {"nx_body"}  # title resta in alto (consueto)

# Shape che ricevono allineamento orizzontale = CENTER
SHAPES_TO_CENTER_HORIZONTALLY = {"nx_caption"}  # title NO: vogliamo allineato sx classico


def patch_font_on_text(sp_el, font_name: str = FONT_FAMILY) -> int:
    """Force font on all runs in the shape. Returns count of patched runs."""
    count = 0
    txBody = sp_el.find(f"{P}txBody")
    if txBody is None:
        return 0
    # Cerca tutti i <a:rPr> e <a:defRPr> e <a:endParaRPr> per scrivere <a:latin typeface="Inter">
    for rPr_tag in ("rPr", "defRPr", "endParaRPr"):
        for rPr in txBody.iter(f"{A}{rPr_tag}"):
            # Rimuovi vecchio latin
            for old in rPr.findall(f"{A}latin"):
                rPr.remove(old)
            latin = etree.SubElement(rPr, f"{A}latin")
            latin.set("typeface", font_name)
            count += 1
    # E anche dentro <a:lstStyle>/<a:lvl*pPr>/<a:defRPr>
    lstStyle = txBody.find(f"{A}lstStyle")
    if lstStyle is not None:
        for rPr in lstStyle.iter(f"{A}defRPr"):
            for old in rPr.findall(f"{A}latin"):
                rPr.remove(old)
            latin = etree.SubElement(rPr, f"{A}latin")
            latin.set("typeface", font_name)
            count += 1
    return count


def set_anchor_middle(sp_el) -> None:
    """Set vertical anchor to MIDDLE on this shape's bodyPr."""
    txBody = sp_el.find(f"{P}txBody")
    if txBody is None:
        return
    bodyPr = txBody.find(f"{A}bodyPr")
    if bodyPr is None:
        bodyPr = etree.SubElement(txBody, f"{A}bodyPr")
        txBody.insert(0, bodyPr)
    bodyPr.set("anchor", "ctr")


def set_align_center(sp_el) -> None:
    """Set horizontal align to center on all paragraphs."""
    txBody = sp_el.find(f"{P}txBody")
    if txBody is None:
        return
    for p in txBody.findall(f"{A}p"):
        pPr = p.find(f"{A}pPr")
        if pPr is None:
            pPr = etree.Element(f"{A}pPr")
            p.insert(0, pPr)
        pPr.set("algn", "ctr")
    # E nel lstStyle per ereditarietà (slide nuove)
    lstStyle = txBody.find(f"{A}lstStyle")
    if lstStyle is not None:
        for lvl in lstStyle.iter():
            if lvl.tag.endswith("pPr"):
                lvl.set("algn", "ctr")


def set_geometry(sp_el, left_in: float, top_in: float, width_in: float, height_in: float) -> None:
    """Set shape geometry in inches."""
    spPr = sp_el.find(f"{P}spPr")
    if spPr is None:
        return
    xfrm = spPr.find(f"{A}xfrm")
    if xfrm is None:
        xfrm = etree.SubElement(spPr, f"{A}xfrm")
    # off
    off = xfrm.find(f"{A}off")
    if off is None:
        off = etree.SubElement(xfrm, f"{A}off")
    off.set("x", str(int(left_in * 914400)))
    off.set("y", str(int(top_in * 914400)))
    # ext
    ext = xfrm.find(f"{A}ext")
    if ext is None:
        ext = etree.SubElement(xfrm, f"{A}ext")
    ext.set("cx", str(int(width_in * 914400)))
    ext.set("cy", str(int(height_in * 914400)))


def patch_master_font_scheme(prs: Presentation, font_name: str = FONT_FAMILY) -> None:
    """Patch the slide master's theme font scheme so all unstyled text uses Inter."""
    master = prs.slide_masters[0]
    # Find the theme
    theme_part = None
    for rel in master.part.rels.values():
        if "theme" in rel.target_ref:
            theme_part = rel.target_part
            break
    if theme_part is None:
        print("[WARN] Theme part non trovato, skip font scheme patch")
        return
    theme_el = etree.fromstring(theme_part.blob)
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    # Major (titles) e Minor (body) latin font
    for scheme in ("majorFont", "minorFont"):
        for el in theme_el.findall(f".//a:{scheme}/a:latin", ns):
            el.set("typeface", font_name)
    theme_part._blob = etree.tostring(theme_el, xml_declaration=True, encoding="UTF-8", standalone=True)
    print(f"[OK] Theme font scheme: Major/Minor latin = {font_name}")


def scale_shape_geometry(sp_el, scale: float) -> bool:
    """Scale shape's xfrm (offset + extent) by `scale` factor. Returns True if patched."""
    spPr = sp_el.find(f"{P}spPr")
    if spPr is None:
        return False
    xfrm = spPr.find(f"{A}xfrm")
    if xfrm is None:
        return False
    off = xfrm.find(f"{A}off")
    if off is not None:
        x = int(off.get("x", "0"))
        y = int(off.get("y", "0"))
        off.set("x", str(int(x * scale)))
        off.set("y", str(int(y * scale)))
    ext = xfrm.find(f"{A}ext")
    if ext is not None:
        cx = int(ext.get("cx", "0"))
        cy = int(ext.get("cy", "0"))
        ext.set("cx", str(int(cx * scale)))
        ext.set("cy", str(int(cy * scale)))
    return True


def scale_font_sizes(sp_el, scale: float) -> int:
    """Scale all font sizes (a:rPr@sz / a:defRPr@sz / a:endParaRPr@sz) by `scale`.

    PowerPoint stores font sizes in `sz` attribute as 100*pt (es. 18pt = 1800).
    Returns count of patched attributes.
    """
    count = 0
    txBody = sp_el.find(f"{P}txBody")
    if txBody is None:
        return 0
    for rPr_tag in ("rPr", "defRPr", "endParaRPr"):
        for rPr in txBody.iter(f"{A}{rPr_tag}"):
            sz = rPr.get("sz")
            if sz is not None:
                try:
                    new_sz = max(800, int(int(sz) * scale))  # min 8pt
                    rPr.set("sz", str(new_sz))
                    count += 1
                except ValueError:
                    pass
    return count


def tune(src: Path, dst: Path) -> None:
    prs = Presentation(str(src))
    print(f"Loaded: {src}")
    print(f"  Slide size pre: {prs.slide_width/914400:.2f}\" x {prs.slide_height/914400:.2f}\"")
    print(f"  Layouts: {len(prs.slide_layouts)}")
    print(f"  Scaling factor: {SCALE:.4f}")
    print()

    # 1) Theme font scheme master-wide
    patch_master_font_scheme(prs, FONT_FAMILY)

    # 2) Slide size → 13.33×7.5"
    prs.slide_width = TARGET_W_EMU
    prs.slide_height = TARGET_H_EMU
    print(f"[OK] Slide size: {TARGET_W_EMU/914400:.3f}\" x {TARGET_H_EMU/914400:.3f}\"")

    total_font_runs = 0
    total_font_sizes = 0
    total_anchors = 0
    total_aligns = 0
    total_geom_scaled = 0

    # 3) Scaling proporzionale TUTTI gli shape (master + tutti i layout)
    for layout_idx, layout in enumerate(prs.slide_layouts):
        for sh in layout.shapes:
            sp_el = sh._element
            name = sh.name
            # Geometria: scaling proporzionale per TUTTI gli shape
            if scale_shape_geometry(sp_el, SCALE):
                total_geom_scaled += 1
            # Font Inter su tutti
            total_font_runs += patch_font_on_text(sp_el)
            # Scaling font sizes proporzionale
            total_font_sizes += scale_font_sizes(sp_el, SCALE)
            # Anchor middle solo su nx_body
            if name in SHAPES_TO_CENTER_VERTICALLY:
                set_anchor_middle(sp_el)
                total_anchors += 1
            # Align center solo su nx_caption
            if name in SHAPES_TO_CENTER_HORIZONTALLY:
                set_align_center(sp_el)
                total_aligns += 1

    # Anche shape del master (logo, decorazioni)
    master = prs.slide_masters[0]
    for sh in master.shapes:
        sp_el = sh._element
        if scale_shape_geometry(sp_el, SCALE):
            total_geom_scaled += 1
        total_font_runs += patch_font_on_text(sp_el)
        total_font_sizes += scale_font_sizes(sp_el, SCALE)

    # Override geometrici post-scaling (se presenti)
    for layout_idx, overrides in LAYOUT_GEOMETRY_OVERRIDE.items():
        if layout_idx >= len(prs.slide_layouts):
            continue
        layout = prs.slide_layouts[layout_idx]
        for sh in layout.shapes:
            if sh.name in overrides:
                l, t, w, h = overrides[sh.name]
                set_geometry(sh._element, l, t, w, h)
                print(f"  [{layout_idx}] {sh.name}: override -> ({l},{t}) {w}x{h}\"")

    dst.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dst))
    print(f"\nSaved: {dst}")
    print(f"  Geometries scaled: {total_geom_scaled}")
    print(f"  Font sizes scaled: {total_font_sizes}")
    print(f"  Font typeface patched (Inter): {total_font_runs}")
    print(f"  Anchors set to MIDDLE (body): {total_anchors}")
    print(f"  Aligns set to CENTER (caption): {total_aligns}")


def verify(path: Path) -> None:
    prs = Presentation(str(path))
    print(f"\n=== VERIFY {path.name} ===")
    print(f"Slide size: {prs.slide_width/914400:.2f}\" x {prs.slide_height/914400:.2f}\"")
    for idx in [1, 2, 3]:
        if idx >= len(prs.slide_layouts):
            continue
        layout = prs.slide_layouts[idx]
        print(f"\nLayout {idx}: {layout.name}")
        for sh in layout.shapes:
            if sh.name in ("nx_title", "nx_body", "nx_image_box", "nx_diagram_box", "nx_caption"):
                l = int(sh.left)/914400
                t = int(sh.top)/914400
                w = int(sh.width)/914400
                h = int(sh.height)/914400
                anchor = "?"
                align = "?"
                font = "?"
                if sh.has_text_frame:
                    try:
                        anchor = str(sh.text_frame.vertical_anchor)
                    except Exception:
                        pass
                    # Read lstStyle font
                    txBody = sh._element.find(f"{P}txBody")
                    if txBody is not None:
                        lstStyle = txBody.find(f"{A}lstStyle")
                        if lstStyle is not None:
                            for latin in lstStyle.iter(f"{A}latin"):
                                font = latin.get("typeface", "?")
                                break
                print(f"  {sh.name:18s} pos=({l:5.2f},{t:5.2f}) {w:5.2f}x{h:5.2f}\"  anchor={anchor}  font={font}")


if __name__ == "__main__":
    tune(SRC, DST)
    verify(DST)
