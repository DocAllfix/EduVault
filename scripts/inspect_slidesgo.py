"""Analisi del template Health and Safety Workshop di Slidesgo."""
from __future__ import annotations
from pptx import Presentation
from pptx.dml.color import RGBColor

p = Presentation("assets/templates/slidesgo_health_safety_original.pptx")

print(f"=== TEMPLATE SLIDESGO ===")
print(f"Slide totali: {len(p.slides)}")
print(f"Layout disponibili: {len(p.slide_layouts)}")
print(f"Slide size: {p.slide_width/914400:.1f} x {p.slide_height/914400:.1f} inch")
print()

# Layout list
print("=== LAYOUTS ===")
for i, layout in enumerate(p.slide_layouts):
    pic_count = sum(1 for s in layout.shapes if "PICTURE" in str(s.shape_type))
    auto_count = sum(1 for s in layout.shapes if "AUTO_SHAPE" in str(s.shape_type))
    placeh_count = sum(1 for s in layout.shapes if "PLACEHOLDER" in str(s.shape_type))
    text_with_content = sum(1 for s in layout.shapes if s.has_text_frame and s.text_frame.text.strip())
    print(f"  [{i}] '{layout.name}': {len(layout.shapes)} shapes ({pic_count} PICTURE, {auto_count} AUTO, {placeh_count} PLACEHOLDER, {text_with_content} con testo)")

print()
print("=== PRIME 25 SLIDE DEL TEMPLATE ===")
for i, slide in enumerate(p.slides):
    if i >= 25:
        break
    layout_name = slide.slide_layout.name
    # First text in slide
    first_text = ""
    for s in slide.shapes:
        if s.has_text_frame and s.text_frame.text.strip():
            first_text = s.text_frame.text.strip().split("\n")[0][:80]
            break
    pic_count = sum(1 for s in slide.shapes if "PICTURE" in str(s.shape_type))
    print(f"  [{i+1:2}] layout='{layout_name[:30]}' pics={pic_count} '{first_text}'")

# Master colors
print()
print("=== COLORI MASTER THEME ===")
try:
    master = p.slide_master
    if hasattr(master, "color_map"):
        cm = master.color_map
        for attr in dir(cm):
            if not attr.startswith("_"):
                print(f"  {attr}")
except Exception as e:
    print(f"  Error: {e}")

# Fonts usati
print()
print("=== FONTS USATI ===")
fonts = set()
for slide in p.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.name:
                        fonts.add(run.font.name)
print(f"  {fonts}")

# Media count
print()
print("=== MEDIA ===")
import zipfile
z = zipfile.ZipFile("assets/templates/slidesgo_health_safety_original.pptx")
imgs = [n for n in z.namelist() if "media" in n.lower()]
print(f"  Total media files: {len(imgs)}")
for n in imgs[:10]:
    info = z.getinfo(n)
    print(f"    {n}: {info.file_size:,} bytes")
if len(imgs) > 10:
    print(f"    ... e altri {len(imgs)-10}")
