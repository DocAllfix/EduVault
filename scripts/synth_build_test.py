"""Synthetic FASE 4 end-to-end build (Master Plan 4.7).

Generates 30 mock slides with speaker_notes, runs ``ProductionBuilder.build``
with ``outputs=["pptx", "pdf", "audio"]``, and asserts:
    - .pptx exists and contains 30 slides
    - .pdf exists on disk
    - audio/{course_id}/ contains at least one .mp3
    - sync_manifest.json is valid JSON with ``total_tracks`` consistent

The script self-supplies:
    - a synthetic master.pptx (the human #R8 template is not in the repo)
    - a fake asyncpg pool that no-ops on .execute(...) (no live DB — #R2)
    - the default empty ``image_map`` (no web images / SVG diagrams)

edge-tts hits ``speech.platform.bing.com`` directly — no API key. If the
endpoint is unreachable, AudioService's per-slide fallback (BP §07.1
invariant) lets the build complete with 0 audio tracks, and we surface
that as a non-zero exit so the failure is visible.

Run from the repo root:
    python scripts/synth_build_test.py
    python scripts/synth_build_test.py --output-dir /tmp/build
"""

from __future__ import annotations

import argparse
import asyncio
import json
import sys
import tempfile
import uuid
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches

from app.builders.pdf_builder import PdfBuilder
from app.builders.pptx_validator import PptxValidator
from app.builders.production_builder import ProductionBuilder
from app.builders.slide_builder import SlideBuilder
from app.models.core import SlideType
from app.models.pipeline import ImageStrategy, SlideContent

NUM_SLIDES = 30
NUM_MODULES = 3

# Same map used by tests/integration/test_slide_builder.py — works on the
# default python-pptx 11-layout template.
TEST_LAYOUT_MAP = {
    SlideType.TITLE: 0,
    SlideType.CONTENT_TEXT: 1,
    SlideType.CONTENT_IMAGE: 8,
    SlideType.DIAGRAM: 8,
    SlideType.QUIZ: 1,
    SlideType.CASE_STUDY: 1,
    SlideType.RECAP: 1,
    SlideType.CLOSING: 5,
}


class _FakePool:
    """Bare-minimum asyncpg pool stub: just absorbs INSERT/UPDATE calls."""

    def __init__(self) -> None:
        self.calls: list[tuple[Any, ...]] = []

    async def execute(self, sql: str, *params: Any) -> None:
        self.calls.append((sql, *params))


def _build_synthetic_template(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))


def _make_synthetic_slides() -> list[SlideContent]:
    """30 slides across 3 modules, all carrying speaker_notes (required by
    AudioService to emit MP3)."""
    types = [
        SlideType.CONTENT_TEXT,
        SlideType.CONTENT_TEXT,
        SlideType.CASE_STUDY,
        SlideType.QUIZ,
        SlideType.RECAP,
    ]
    slides: list[SlideContent] = []
    per_module = NUM_SLIDES // NUM_MODULES
    for module_index in range(NUM_MODULES):
        for j in range(per_module):
            idx = module_index * per_module + j
            stype = types[j % len(types)]
            slides.append(
                SlideContent(
                    index=idx,
                    module_index=module_index,
                    slide_type=stype,
                    title=f"Slide {idx}",
                    body=(
                        f"Contenuto della slide {idx} su sicurezza sul lavoro "
                        f"per modulo {module_index + 1}."
                    ),
                    speaker_notes=(
                        f"Soffermarsi 20 secondi sulla slide {idx}. "
                        f"Punto chiave del modulo {module_index + 1}: "
                        f"applicazione della normativa di riferimento."
                    ),
                    normative_ref="Art. 1, DM 388/2003",
                    source_chunk_ids=[],
                    image=ImageStrategy(strategy="none"),
                )
            )
    return slides


def _build_production_builder(
    template_path: Path, output_dir: Path
) -> ProductionBuilder:
    """ProductionBuilder bypassing the real template requirement (#R8)."""
    builder = ProductionBuilder.__new__(ProductionBuilder)
    builder.brand_config = {}
    builder.slide_builder = SlideBuilder(
        template_path=template_path,
        output_dir=output_dir,
        layout_map=TEST_LAYOUT_MAP,
    )
    builder.pdf_builder = PdfBuilder(output_dir=output_dir)
    builder.validator = PptxValidator()
    return builder


async def _noop_ws_callback(job_id: str, percent: int, message: str) -> None:
    print(f"  [{percent:>3}%] {message}")


def _check_weasyprint_available() -> str | None:
    """Try a real import of weasyprint and return a human-readable reason
    if GTK runtime is missing (#R12). Returns None when import succeeds.

    The PdfBuilder import inside ProductionBuilder is lazy, so the missing
    GTK only surfaces during build(); checking up-front lets the script
    exit cleanly with actionable instructions instead of a deep traceback.
    """
    try:
        import weasyprint  # noqa: F401
    except OSError as exc:
        return (
            f"WeasyPrint cannot load GTK runtime: {exc}. "
            "This is expected on Windows dev hosts (#R12 in "
            "docs/VERIFICATION_DEBT.md). Re-run inside the Docker image, "
            "where libgobject/cairo/pango are installed:\n"
            "    docker compose run --rm backend python -m scripts.synth_build_test"
        )
    return None


async def _run(output_dir: Path) -> int:
    print(f"FASE 4 synthetic build -> output: {output_dir}")
    output_dir.mkdir(parents=True, exist_ok=True)

    template_path = output_dir / "synthetic_master.pptx"
    _build_synthetic_template(template_path)

    slides = _make_synthetic_slides()
    assert len(slides) == NUM_SLIDES, f"expected {NUM_SLIDES}, got {len(slides)}"

    builder = _build_production_builder(template_path, output_dir)
    pool = _FakePool()
    course_id = f"synth-{uuid.uuid4().hex[:8]}"
    course: dict[str, Any] = {
        "id": course_id,
        "title": "Corso Sintetico FASE 4",
        "duration_hours": 1,
        "target": "discente",
        "outputs": ["pptx", "pdf", "audio"],
    }

    # AudioService writes its MP3s under output/audio/{course_id} (default
    # _AUDIO_ROOT="output/audio"). Redirect it to our output_dir so the
    # script doesn't pollute the project tree.
    from app.services import audio_service

    audio_service._AUDIO_ROOT = output_dir / "audio"

    pptx_path, pdf_path, report = await builder.build(
        slides=slides,
        course=course,
        job_id="synth-job",
        ws_callback=_noop_ws_callback,
        image_map={},
        db=pool,
    )

    # ─────────── verification ───────────
    errors: list[str] = []

    pptx_p = Path(pptx_path)
    if not pptx_p.is_file():
        errors.append(f"PPTX missing: {pptx_path}")
    else:
        prs = Presentation(pptx_path)
        if len(prs.slides) != NUM_SLIDES:
            errors.append(
                f"PPTX slide count mismatch: {len(prs.slides)} vs {NUM_SLIDES}"
            )

    pdf_p = Path(pdf_path)
    if not pdf_p.is_file():
        errors.append(f"PDF missing: {pdf_path}")

    audio_dir = output_dir / "audio" / course_id
    if not audio_dir.is_dir():
        errors.append(f"Audio dir missing: {audio_dir}")
    else:
        mp3s = sorted(audio_dir.glob("*.mp3"))
        if not mp3s:
            errors.append(
                f"No MP3 generated in {audio_dir} — edge-tts unreachable? "
                "All slides hit the per-slide fallback."
            )

        manifest_path = audio_dir / "sync_manifest.json"
        if not manifest_path.is_file():
            errors.append(f"sync_manifest.json missing: {manifest_path}")
        else:
            try:
                manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
            except json.JSONDecodeError as exc:
                errors.append(f"sync_manifest.json invalid JSON: {exc}")
            else:
                if manifest.get("course_id") != course_id:
                    errors.append(
                        f"manifest course_id mismatch: "
                        f"{manifest.get('course_id')} != {course_id}"
                    )
                if manifest.get("total_tracks") != len(mp3s):
                    errors.append(
                        f"manifest total_tracks mismatch: "
                        f"{manifest.get('total_tracks')} != {len(mp3s)}"
                    )

    print("\n----------- BUILD REPORT -----------")
    print(f"PPTX:           {pptx_path}")
    print(f"PDF:            {pdf_path}")
    print(f"Audio dir:      {audio_dir}")
    print(f"Total slides:   {report.get('total_slides')}")
    print(f"Modules done:   {report.get('modules_completed')}")
    print(f"Warnings:       {report.get('warnings')}")
    print(f"Pool execute calls: {len(pool.calls)}")
    print("------------------------------------")

    if errors:
        print("\nVERIFICATION FAILED:")
        for e in errors:
            print(f"  [x] {e}")
        return 1

    print("\nAll checks passed.")
    return 0


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Synthetic FASE 4 build test (PPTX + PDF + Audio)."
    )
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=Path(tempfile.gettempdir()) / "eduvault_synth_build",
        help="Where to write the synthetic build artifacts.",
    )
    args = parser.parse_args(argv)

    missing = _check_weasyprint_available()
    if missing is not None:
        print(f"ENVIRONMENT NOT READY:\n{missing}", file=sys.stderr)
        return 2

    exit_code = asyncio.run(_run(args.output_dir))
    if exit_code == 0:
        print("\nFASE 4 verificata. Pronto per commit.")
    return exit_code


if __name__ == "__main__":
    sys.exit(main())
