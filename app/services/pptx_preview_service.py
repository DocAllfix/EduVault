"""F-STUDIO-UX Step 0 v2 (2026-06-02 / D-207) — PPTX-fedele preview service,
strategia memory-light single-slide.

Renderizza UNA singola slide del PPTX scaricabile come PNG, estraendo PRIMA
quella slide in un mini-PPTX temporaneo (peso ~100-500KB invece di ~25MB del
PPTX completo), poi convertendola con soffice:

  PPTX completo (~25MB, 342 slide)
        │
        ▼  python-pptx: copia + rimuovi tutte le slide tranne N
        │
  mini-PPTX (~100-500KB, 1 slide)
        │
        ▼  soffice --headless --convert-to pdf (memory-light)
        │
  PDF 1 pagina (~50-200KB)
        │
        ▼  pypdfium2 render
        │
  PNG cached: output/previews/{course_id}/{rebuild_token}/{idx:04d}.png

Razionale del cambio strategia v1→v2:
  v1 convertiva l'intero PPTX 342 slide in un PDF unico, riusato per tutte le
  slide. Vantaggio: 1 sola chiamata soffice per corso. Svantaggio: soffice
  caricava in memoria l'intero PPTX (~1GB RAM su corsi 4h con immagini reali),
  causando OOM-kill del container Railway (verificato in prod su corso
  af08e1d1, status 502 + "Stopping Container" nei log).
  v2 limita ogni invocazione soffice a 1 sola slide → memoria proporzionale a
  1 slide (~50-100MB). Trade-off: 1 chiamata soffice per slide invece di 1
  per corso → caricamento iniziale Course Studio piu` lento (5-10s/slide la
  prima volta), ma navigazione fluida dopo (PNG cached). Il caller dovrebbe
  pre-fetchare le slide visibili.

Cache layout:
    output/previews/{course_id}/{rebuild_token}/
        {idx:04d}.png          ← cache della singola slide
        ...
    (NO PDF intermedio cached: ogni slide ha il suo workflow temp + cleanup)

LibreOffice è già installato nel Dockerfile (libreoffice-impress + libreoffice-core,
linee 17-18). Zero apt-deps nuove.

Fallback: se soffice fallisce / OOM / timeout → PreviewRenderError → endpoint
fa graceful fallback a ``preview_source="pdf_dispensa"`` (testo-only).
"""

from __future__ import annotations

import asyncio
import os
import subprocess
from pathlib import Path

import structlog

logger = structlog.get_logger(__name__)

# Timeout soffice — mini-PPTX single-slide deve essere veloce (5-15s tipico).
# Margine 60s per cold-start LibreOffice (caricamento font, JVM, ecc.).
_SOFFICE_TIMEOUT_S = 60

# Render DPI: stesso scale del path legacy (1.6 → ~153 DPI, leggibile su retina
# senza far esplodere la dimensione della cache).
_PDFIUM_RENDER_SCALE = 1.6

# Lock async per serializzare le invocazioni soffice nello stesso processo —
# LibreOffice headless ha problemi noti con istanze concorrenti che condividono
# lo stesso user profile. Per process pool / multi-worker servirebbe un lock
# filesystem-level, ma per il deploy single-worker Railway è sufficiente.
_soffice_lock = asyncio.Lock()


class PreviewRenderError(RuntimeError):
    """Raised when the LibreOffice → PDF conversion fails irrecoverably.

    Caller (endpoint preview) può catturare questo errore per fall-back al path
    legacy ``pdf_dispensa`` senza propagare 500 all'utente.
    """


async def render_pptx_slide_to_png(
    *,
    pptx_path: Path,
    cache_dir: Path,
    slide_index: int,
) -> Path:
    """Renderizza la slide ``slide_index`` del PPTX come PNG, riusando la cache.

    v2 (2026-06-02): strategia single-slide memory-light. Estrae la slide N in
    un mini-PPTX temporaneo (peso ridotto), converte solo quella con soffice,
    poi cancella il temp. Soffice usa proporzionalmente meno RAM (1 slide vs
    342) → niente OOM-kill su container small-tier.

    Args:
        pptx_path: percorso del file ``.pptx`` (deve esistere su disco — il caller
            estrae da ``courses.pptx_path`` e valida l'esistenza).
        cache_dir: directory di cache della preview (caller esistente in
            ``courses.py:966`` calcola già il path con ``rebuild_token``):
            ``output/previews/{course_id}/{rebuild_token}/``.
        slide_index: indice 0-based della slide nel PPTX completo.

    Returns:
        ``Path`` al file PNG. Esiste su disco al ritorno.

    Raises:
        PreviewRenderError: se soffice fallisce, OOM, o produce PDF malformato.
        FileNotFoundError: se ``pptx_path`` non esiste.
        IndexError: se ``slide_index`` non e` un indice valido del PPTX.
    """
    if not pptx_path.is_file():
        raise FileNotFoundError(f"PPTX non trovato: {pptx_path}")

    cache_dir.mkdir(parents=True, exist_ok=True)
    png_path = cache_dir / f"{slide_index:04d}.png"

    # Fast path: PNG già in cache → nessuna conversione necessaria.
    if png_path.is_file():
        return png_path

    # Workspace temp per estrazione single-slide + conversione + cleanup.
    # Tutto wrapped in lock async + try/finally per garantire cleanup anche
    # su errore parziale (file orfani in cache_dir confondono i lookup).
    work_dir = cache_dir / f".tmp_{slide_index:04d}"

    async with _soffice_lock:
        # Re-check dopo l'acquisizione del lock: un'altra richiesta potrebbe
        # aver gia` generato il PNG mentre eravamo in attesa.
        if png_path.is_file():
            return png_path
        try:
            # 1) extract single slide → mini.pptx (lavoro python-pptx,
            #    sincrono ma veloce su singola slide ~50-200ms).
            mini_pptx = work_dir / "single_slide.pptx"
            await asyncio.to_thread(
                _extract_single_slide_to_pptx,
                pptx_path,
                slide_index,
                mini_pptx,
            )

            # 2) soffice mini-pptx → mini-pdf (memory-light).
            await _convert_single_pptx_to_pdf(mini_pptx, work_dir)

            mini_pdf = work_dir / "single_slide.pdf"
            if not mini_pdf.is_file():
                raise PreviewRenderError(
                    f"soffice non ha prodotto il PDF atteso: {mini_pdf}"
                )

            # 3) pdfium render del singolo PDF (1 pagina) → PNG cached.
            await asyncio.to_thread(
                _render_pdf_page_to_png, mini_pdf, 0, png_path
            )
        finally:
            # Cleanup workspace temp anche su error per non accumulare detrito.
            await asyncio.to_thread(_safe_rmtree, work_dir)

    return png_path


def _extract_single_slide_to_pptx(
    source_pptx: Path, slide_index: int, dest_pptx: Path
) -> None:
    """Crea un mini-PPTX con SOLO la slide ``slide_index``.

    Strategia memory-light: copia il file PPTX originale come ZIP (è un OOXML
    package), poi modifica internamente per tenere solo la slide voluta.

    Fallback: usa python-pptx per duplicare il file + rimuovere tutte le slide
    tranne N. Funziona ma richiede di caricare l'intero file in memoria.

    Args:
        source_pptx: percorso PPTX originale (read-only).
        slide_index: indice 0-based della slide da preservare.
        dest_pptx: percorso destinazione mini-PPTX (singola slide).

    Raises:
        IndexError: se slide_index non e` un indice valido.
    """
    from pptx import Presentation
    from pptx.oxml.ns import qn

    dest_pptx.parent.mkdir(parents=True, exist_ok=True)

    # Carica + remove all slides except target.
    prs = Presentation(str(source_pptx))
    total = len(prs.slides)
    if slide_index < 0 or slide_index >= total:
        raise IndexError(
            f"slide_index {slide_index} fuori range (0..{total - 1})"
        )

    # Trick python-pptx: rimuovere slide via XML perché l'API non espone .remove.
    # Ref: https://github.com/scanny/python-pptx/issues/67
    sldIdLst = prs.slides._sldIdLst
    # Lista degli elements <p:sldId> (uno per slide).
    sldId_elements = list(sldIdLst)
    target_sldId = sldId_elements[slide_index]
    # Rimuovi tutti gli altri sldId (la slide voluta resta).
    for sldId in sldId_elements:
        if sldId is not target_sldId:
            sldIdLst.remove(sldId)
            # Rimuovi anche la part relationship corrispondente per non
            # lasciare relazioni orfane.
            rId = sldId.get(qn("r:id"))
            try:
                prs.part.drop_rel(rId)
            except Exception:  # noqa: BLE001
                # drop_rel non e` garantito esistere; ignorare e` safe.
                pass

    prs.save(str(dest_pptx))


def _safe_rmtree(path: Path) -> None:
    """Rimuove ricorsivamente ``path`` ignorando errori (best-effort cleanup)."""
    import shutil

    if path.exists():
        shutil.rmtree(path, ignore_errors=True)


async def _convert_single_pptx_to_pdf(pptx_path: Path, out_dir: Path) -> None:
    """Esegue ``soffice --headless --convert-to pdf`` su un mini-PPTX 1-slide.

    soffice produce il PDF nella ``out_dir`` con lo stesso basename del PPTX
    (es. ``single_slide.pptx`` → ``single_slide.pdf``).

    Memory profile: 1 slide → ~50-150MB RAM di soffice headless (vs ~1GB per
    342 slide nel path v1). Niente OOM su Railway tier.

    Profile utente isolato per evitare "Office is already running" su istanze
    concorrenti che condividono /root/.config/libreoffice.
    """
    profile_dir = out_dir / ".soffice_profile"
    profile_dir.mkdir(parents=True, exist_ok=True)
    user_profile_url = f"-env:UserInstallation=file://{profile_dir.as_posix()}"

    cmd = [
        "soffice",
        user_profile_url,
        "--headless",
        "--norestore",
        "--nologo",
        "--nofirststartwizard",
        "--convert-to",
        "pdf",
        "--outdir",
        str(out_dir),
        str(pptx_path),
    ]

    logger.info(
        "soffice_convert_started",
        pptx=str(pptx_path),
        outdir=str(out_dir),
        memory_strategy="single_slide_v2",
    )

    proc = await asyncio.create_subprocess_exec(
        *cmd,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        env={**os.environ, "HOME": str(profile_dir)},
    )
    try:
        stdout_b, stderr_b = await asyncio.wait_for(
            proc.communicate(), timeout=_SOFFICE_TIMEOUT_S
        )
    except asyncio.TimeoutError:
        proc.kill()
        await proc.wait()
        raise PreviewRenderError(
            f"soffice timeout dopo {_SOFFICE_TIMEOUT_S}s su {pptx_path.name}"
        )

    if proc.returncode != 0:
        raise PreviewRenderError(
            f"soffice exit_code={proc.returncode} stderr={stderr_b.decode(errors='replace')[:500]}"
        )

    # soffice scrive il PDF con stesso basename del mini-PPTX.
    produced_pdf = out_dir / f"{pptx_path.stem}.pdf"
    if not produced_pdf.is_file():
        raise PreviewRenderError(
            f"soffice non ha prodotto il PDF atteso: {produced_pdf}"
        )

    logger.info(
        "soffice_convert_completed",
        pdf=str(produced_pdf),
        size_bytes=produced_pdf.stat().st_size,
    )


def _render_pdf_page_to_png(
    pdf_path: Path, page_index: int, png_out: Path
) -> None:
    """Rendering CPU-bound di una singola pagina PDF → PNG via pypdfium2.

    Funzione sincrona: il chiamante la incapsula in ``asyncio.to_thread``.
    """
    import pypdfium2 as pdfium  # type: ignore[import-untyped]

    pdf = pdfium.PdfDocument(str(pdf_path))
    try:
        if page_index < 0 or page_index >= len(pdf):
            raise PreviewRenderError(
                f"slide {page_index} fuori range PDF ({len(pdf)} pagine)"
            )
        page = pdf[page_index]
        pil_image = page.render(scale=_PDFIUM_RENDER_SCALE).to_pil()
        pil_image.save(png_out, format="PNG", optimize=True)
    finally:
        pdf.close()
