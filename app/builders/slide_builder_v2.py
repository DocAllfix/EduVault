"""SlideBuilderV2 — XML-level text substitution via python-pptx + lxml (MIT).

FASE 3 vast-hopping-sketch: risolve il GAP shape_map del SlideBuilder v1
(test E.03/E.04) che cercava placeholder PowerPoint canonici inesistenti nel
template Claude Design (solo AUTO_SHAPE). V2 lavora direttamente sulle shape
del layout, sostituendo il testo placeholder originale (es. "Formazione Generale
dei Lavoratori") con il contenuto LLM-generated.

Approccio: per ogni SlideType, mappiamo le shape semantiche (title/body/footer/
options/...) all'INDICE numerico dello shape nel layout — derivato da
``scripts/inspect_pptx_template.py`` su ``nexus_master.pptx``. La mappa è
hardcoded qui (deterministica, ~50 righe) invece che caricata da JSON: scelta
karpathy-guidelines (no astrazione prematura per single-use data).

Distinto da SlideBuilder v1 (deprecato non rimosso):
- v1 usa ``slide.placeholders`` API → fallisce su AUTO_SHAPE
- v2 usa ``slide.shapes[idx]`` + iterazione text_frame → funziona su qualsiasi shape

Sincronizzazione: il chiamante (``ProductionBuilder`` con ``asyncio.to_thread``)
mantiene il vincolo REI-3 ``Semaphore(1)``: python-pptx + lxml restano non
thread-safe, qui non cambiamo nulla.

Licenza pulita: codice nostro, ispirato al pattern OOXML standard Microsoft
(zip+xml manipulation). NESSUNA copia dal toolkit Anthropic skill `pptx`
(licenza proprietaria). Dipendenze: python-pptx (MIT) + lxml (BSD-like), già
in pyproject.toml.
"""

from __future__ import annotations

import copy
import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import structlog
from pptx import Presentation

from app.models.core import SlideType
from app.models.pipeline import SlideContent

logger = structlog.get_logger()


# FIX #30.0-sexies (2026-05-26): nexus_master_v4_patched.pptx — base v4_bulletfix
# dell'analista + patch_master_v4.py idempotente che aggiunge:
#  - anchor="ctr" su 4 body placeholder (CONTENT_TEXT/CONTENT_IMAGE/RECAP/MODULE_CLOSE)
#    per centratura verticale (conferma analista P2: min-4 bullet garantito → varianza
#    contenuta, ctr dà inquadramento bilanciato).
#  - nx_page footer textbox al layout RECAP (mancava nel v4_bulletfix, causa warning
#    shape_missing — analista P1 addendum).
DEFAULT_TEMPLATE = Path("assets/templates/nexus_master_v4_patched.pptx")
DEFAULT_OUTPUT_DIR = Path("output")
IMAGE_MISSING_FALLBACK = "[Immagine non disponibile]"


# Mapping SlideType → indice layout in nexus_master_v4_bulletfix.pptx.
# Indici cambiati vs v3 per accomodare i 2 nuovi layout MODULE_OPEN/CLOSE:
#   0 NX TITLE | 1 NX MODULE_OPEN | 2 NX CONTENT_TEXT | 3 NX CONTENT_IMAGE
#   4 NX DIAGRAM | 5 NX QUIZ | 6 NX CASE_STUDY | 7 NX MODULE_CLOSE
#   8 NX RECAP | 9 NX CLOSING
# MODULE_OPEN/CLOSE saranno gestiti in FIX #30.2 (nuove SlideType + bookend pacing).
DEFAULT_LAYOUT_MAP: dict[SlideType, int] = {
    SlideType.TITLE: 0,
    SlideType.MODULE_OPEN: 1,    # FIX #30.2: bookend apertura modulo
    SlideType.CONTENT_TEXT: 2,
    SlideType.CONTENT_IMAGE: 3,
    SlideType.DIAGRAM: 4,
    SlideType.QUIZ: 5,
    SlideType.CASE_STUDY: 6,
    SlideType.MODULE_CLOSE: 7,   # FIX #30.2: bookend chiusura modulo
    SlideType.RECAP: 8,
    SlideType.CLOSING: 9,
}


# FIX #28.3a (2026-05-26): SHAPE_MAP esteso a DUE tipi di target per ruolo,
# secondo il design pulito dell'analista ("un solo writer role-based + un solo
# resolver che incapsula la differenza idx/nome"):
#   - str "nx_xxx"     → target AUTO_SHAPE risolto per NAME (legacy, robusto per
#                        shape custom RECAP/CASE_STUDY/QUIZ/decorativi)
#   - ("ph", int)      → target PLACEHOLDER risolto per IDX (stabile al
#                        round-trip PowerPoint, che può rinominare i placeholder).
#                        Usato per BODY/TITLE/PICTURE veri (v3 template).
# Il resolver _resolve_shape unifica i due path; il rendering logic non lo sa.
SHAPE_TARGET = "str | tuple[str, int] | None"

SHAPE_MAP: dict[SlideType, dict[str, object]] = {
    SlideType.TITLE: {
        "title": "nx_title",
        "subtitle": "nx_subtitle",
    },
    SlideType.CONTENT_TEXT: {
        # v3: nx_title e nx_body sono veri placeholder PowerPoint (idx 0 e 1).
        # Bullet ereditato + normAutofit nativo + nessun overlap dal layout.
        "title": ("ph", 0),
        "body":  ("ph", 1),
        "normative_ref": "nx_ref",
        "page_num": "nx_page",
    },
    SlideType.CONTENT_IMAGE: {
        "title": ("ph", 0),
        "body":  ("ph", 1),
        "image":  ("ph", 2),  # PICTURE placeholder (insert_picture sopra)
        "image_caption": "nx_caption",
        "normative_ref": "nx_ref",
        "page_num": "nx_page",
    },
    SlideType.DIAGRAM: {
        # DIAGRAM body-less per design (caption ≠ body multi-bullet — analista 2026-05-26).
        # nx_caption resta AUTO_SHAPE: una didascalia descrittiva NON deve ereditare bullet.
        "title": ("ph", 0),
        "image": ("ph", 2),
        "diagram_caption": "nx_caption",
        "normative_ref": "nx_ref",
        "page_num": "nx_page",
    },
    SlideType.QUIZ: {
        # FIX #30.0-quinquies: v4 ha rinominato il box domanda da nx_title a nx_question.
        "title": "nx_question",
        "option_a": "nx_option_a",
        "option_b": "nx_option_b",
        "option_c": "nx_option_c",
        "option_d": "nx_option_d",
        "correct_marker": "nx_correct_marker",
    },
    SlideType.CASE_STUDY: {
        "title": "nx_title",
        "situazione": "nx_situazione",
        "azione": "nx_azione",
        "risultato": "nx_risultato",
        "normative_ref": "nx_ref",
        "page_num": "nx_page",
    },
    SlideType.RECAP: {
        # FIX #30.0-sexies (2026-05-26, conferma analista P1): v4 ha collassato
        # i 5 box separati (nx_recap_text_710..750 + 5 checkmark) in UN solo
        # placeholder BODY idx=1 (nx_recap_body) con lstStyle/lvl1pPr/buChar=✓.
        # Pattern identico a CONTENT_TEXT — un body multi-bullet che eredita
        # il glifo (✓ invece di •) dal template. Stesso writer
        # `_write_shape_by_role_bullets` con `slide_content.bullets` paragrafo
        # per paragrafo (NO newline join — ogni \n sarebbe soft line break NON
        # nuovo paragrafo, e l'ereditarietà del bullet si applica per
        # paragrafo).
        "title": ("ph", 0),
        "body":  ("ph", 1),
        # nx_page / nx_module_ref non esistono nel layout RECAP di v4_bulletfix
        # (sarà aggiunto in patch_master_v4 — vedi todo P2). Per ora il warning
        # `shape_missing` resta atteso e silenziato in BuildReport.
        "page_num": "nx_page",
    },
    SlideType.CLOSING: {
        "title": "nx_title",
        "tagline": "nx_tagline",
    },
    # FIX #30.2 (2026-05-26): bookends. Layout v4 MODULE_OPEN ha
    # nx_module_num (title ph idx=0, font 72pt) + nx_module_title (body ph
    # idx=1, font 32pt). MODULE_CLOSE ha nx_module_num + nx_module_title +
    # nx_recap_body (body ph idx=2 con bullet ✓ ereditati).
    SlideType.MODULE_OPEN: {
        "module_num":   ("ph", 0),  # "MODULO N"
        "module_title": ("ph", 1),  # titolo modulo
    },
    SlideType.MODULE_CLOSE: {
        "module_num":   ("ph", 0),  # "RIEPILOGO MODULO N"
        "module_title": ("ph", 1),  # sintesi modulo
        "body":         ("ph", 2),  # 5 bullet ✓ (riusa _write_shape_by_role_bullets)
    },
}


@dataclass
class BuildReport:
    """Per-build counters for diagnostics."""

    slides_built: int = 0
    shapes_written: int = 0
    images_inserted: int = 0
    image_fallbacks: int = 0
    warnings: list[str] = field(default_factory=list)


def _is_local_path(value: str | None) -> bool:
    """True iff ``value`` is a non-empty local file path candidate.

    FIX #30.0-septies (2026-05-26): NON serve che `value` ESISTA su disco —
    serve solo che ABBIA forma di path locale (no URL). L'esistenza del file
    (originale o variante _fitted) viene verificata da _maybe_insert_image,
    che ha la logica di fallback su _fitted. Questa funzione è solo un filtro
    di forma: distingue "URL Pexels" da "path file system".
    """
    if not value:
        return False
    lowered = value.lower()
    if lowered.startswith(("http://", "https://", "file://", "ftp://")):
        return False
    # Path locale formalmente valido — l'esistenza la verifica il consumer
    return True


def _clean_caption(query: str | None) -> str:
    """Turn an image query into a human caption — no "[ ]" placeholder syntax.

    FIX #25 (2026-05-26): empty query → empty caption (blank shape, never a
    "[ immagine ]" stub). Otherwise capitalise the first letter and trim, so
    a real photo reads as "Cassetta primo soccorso", not "[ cassetta ... ]".
    """
    q = (query or "").strip()
    if not q:
        return ""
    return q[0].upper() + q[1:]


import re as _re

# FIX #27.2: marker testuali dei placeholder del template v2. Qualsiasi shape
# il cui testo inizia con uno di questi (o con '[') e che NON è stato riscritto
# va azzerato a fine render, così non resta MAI un placeholder visibile.
_PLACEHOLDER_MARKERS: tuple[str, ...] = (
    "Testo sezione",
    "Titolo Slide",
    "Primo punto",
    "Secondo punto",
    "Terzo punto",
    "Quarto punto",
    "Quinto punto",
    "Sottotitolo",
    "Didascalia",
    "Cosa hai imparato",
    "AUTO_SHAPE",
)
# Etichette DECORATIVE fisse del template (band, label sezione, checkmark): NON
# toccare — fanno parte del design, non sono placeholder di contenuto.
_DECORATIVE_SHAPE_NAMES: frozenset[str] = frozenset({
    "nx_recap_band_label", "nx_case_band_label",
    "nx_case_label_situazione", "nx_case_label_azione", "nx_case_label_risultato",
})


def _parse_case_study_sections(body: str) -> list[str]:
    """Split a CASE_STUDY body into [situazione, azione, risultato].

    Cascade: (1) '---' separator (the prompt-mandated format), (2) labelled
    sections "Situazione:/Azione:/Risultato:", (3) newline. Always returns the
    cleaned non-empty pieces in order. Label prefixes are stripped so the
    rendered text reads "un operaio salda…" not "Situazione: un operaio…"
    (the template already shows the SITUAZIONE/AZIONE/RISULTATO labels).
    """
    body = (body or "").strip()
    if not body:
        return []

    if "---" in body:
        parts = [p.strip() for p in body.split("---")]
    else:
        # Try labelled sections (case-insensitive) anywhere in the text.
        label_re = _re.compile(
            r"(?im)^\s*(situazione|azione|risultato|decisione|esito)\s*[:\-]\s*"
        )
        if label_re.search(body):
            # Split keeping the text after each label.
            chunks = _re.split(r"(?im)^\s*(?:situazione|azione|risultato|decisione|esito)\s*[:\-]\s*", body)
            parts = [c.strip() for c in chunks if c.strip()]
        else:
            parts = [ln.strip() for ln in body.split("\n")]

    # Strip a leading inline label if still present (e.g. "Situazione: ...").
    cleaned: list[str] = []
    for p in parts:
        p = _re.sub(
            r"^\s*(situazione|azione|risultato|decisione|esito)\s*[:\-]\s*",
            "", p, flags=_re.IGNORECASE,
        ).strip()
        if p:
            cleaned.append(p)
    return cleaned


def _blank_layout_placeholder_text(prs: Any) -> None:
    """FIX #27.7: svuota il testo placeholder degli shape NEI LAYOUT.

    Ogni slide eredita gli shape del proprio layout (showMasterSp). Se un layout
    ha `nx_title="Cosa hai imparato"`, quel testo appare SOTTO il titolo scritto
    nella slide → sovrapposizione. Svuotiamo i text-frame di tutti gli shape dei
    layout TRANNE le label decorative fisse (banda RIEPILOGO/CASO STUDIO, label
    sezione SITUAZIONE/AZIONE/RISULTATO, checkmark ✓) che devono restare.
    """
    keep = _DECORATIVE_SHAPE_NAMES | {
        "nx_recap_checkmark_710", "nx_recap_checkmark_720", "nx_recap_checkmark_730",
        "nx_recap_checkmark_740", "nx_recap_checkmark_750",
    }
    for layout in prs.slide_layouts:
        for sh in layout.shapes:
            if (sh.name or "") in keep:
                continue
            if sh.has_text_frame and sh.text_frame.text.strip():
                sh.text_frame.text = ""


def _blank_unwritten_placeholders(slide: Any) -> None:
    """FIX #27.2: azzera ogni text_frame ancora col placeholder del template.

    Garanzia anti-placeholder universale (decisione utente: "slide sempre piena,
    pulita, funzionale"). Scorre le shape della slide DOPO che i ruoli sono stati
    scritti: se un text_frame inizia con '[' o con un marker placeholder noto, e
    la shape NON è decorativa fissa, lo svuota. Non tocca shape già riscritte con
    contenuto reale (che non iniziano con i marker).
    """
    for shape in slide.shapes:
        name = shape.name or ""
        if name in _DECORATIVE_SHAPE_NAMES:
            continue
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        is_placeholder = text.startswith("[") or any(
            text.startswith(m) for m in _PLACEHOLDER_MARKERS
        )
        if is_placeholder:
            shape.text_frame.text = ""


def _set_shape_text(shape: Any, text: str) -> bool:
    """Replace text in a shape's text_frame.

    Returns True if the write succeeded, False if the shape has no text_frame.

    FIX #28.3 (2026-05-26): RIMOSSO il blocco autofit FIX #21 (run.font.size +
    MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE). I due meccanismi confliggevano (PowerPoint
    ricalcolava all'apertura, causando il "documento riparato"). Per i BODY
    placeholder veri del template v3 l'autofit lo fa <a:normAutofit/> nativamente.
    Per gli AUTO_SHAPE residui (RECAP text shapes, CASE_STUDY sezioni, decorativi)
    la dimensione è fissa nel layout — preferiamo testo che eventualmente esce dal
    box piuttosto che il workaround che rompe PowerPoint.
    """
    if not shape.has_text_frame:
        return False
    tf = shape.text_frame
    tf.text = text
    tf.word_wrap = True
    return True


def _write_shape_by_role(
    slide: Any,
    layout_type: SlideType,
    role: str,
    text: str,
    report: BuildReport,
) -> None:
    """Scrive `text` nello shape del `role` nel `layout_type` cercando per NAME.

    Template v2 (Claude Design) usa shape names canonici `nx_*` invece di
    indici fissi. Più robusto a future modifiche del template.
    """
    target_shape = _resolve_shape(slide, layout_type, role)
    if target_shape is None:
        return  # già loggato dal resolver
    if _set_shape_text(target_shape, text):
        report.shapes_written += 1


def _resolve_shape(slide: Any, layout_type: SlideType, role: str) -> Any | None:
    """FIX #28.3 (analista 2026-05-26): UNICO punto che conosce la differenza
    placeholder/AUTO_SHAPE. Il rendering logic non lo sa né gli importa.

    Target nella SHAPE_MAP:
      ("ph", idx) → cerca via slide.placeholders[idx] (stabile al round-trip)
      "nx_xxx"   → cerca via shape.name == "nx_xxx" (legacy per AUTO_SHAPE custom)
      None       → ruolo non mappato per quel layout (silente)
    """
    role_map = SHAPE_MAP.get(layout_type, {})
    target = role_map.get(role)
    if target is None:
        return None

    # Placeholder per idx (più robusto: PowerPoint può rinominare i placeholder al
    # round-trip, ma idx è stabile)
    if isinstance(target, tuple) and len(target) == 2 and target[0] == "ph":
        ph_idx = target[1]
        try:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == ph_idx:
                    return ph
        except Exception as exc:
            logger.warning(
                "placeholder_lookup_failed",
                layout=layout_type.value, role=role, idx=ph_idx, error=str(exc)[:120],
            )
            return None
        logger.warning(
            "placeholder_missing", layout=layout_type.value, role=role, idx=ph_idx
        )
        return None

    # AUTO_SHAPE per name (RECAP, CASE_STUDY, QUIZ, decorativi)
    if isinstance(target, str):
        for shape in slide.shapes:
            if shape.name == target:
                return shape
        logger.warning(
            "shape_missing", layout=layout_type.value, role=role, name=target
        )
        return None

    logger.warning(
        "shape_target_unknown", layout=layout_type.value, role=role, target=repr(target)
    )
    return None


def _write_shape_by_role_bullets(
    slide: Any,
    layout_type: SlideType,
    role: str,
    items: list[str],
    report: BuildReport,
) -> None:
    """FIX #28.3 (analista 2026-05-26): scrive una LISTA di bullet come paragrafi
    SEPARATI (NO `join("\\n")`), con `tf.clear()` + `add_paragraph()` per ogni
    elemento. Funziona uniforme sia su placeholder BODY (bullet ereditato dal
    master) sia su AUTO_SHAPE (testo nudo senza glifo — è il template a
    determinare il bullet, non il codice).
    """
    target_shape = _resolve_shape(slide, layout_type, role)
    if target_shape is None or not target_shape.has_text_frame:
        return
    tf = target_shape.text_frame
    tf.clear()
    if not items:
        return
    # Il primo paragrafo è già presente dopo tf.clear(); riempi quello, poi append.
    first = tf.paragraphs[0]
    first.text = items[0]
    for extra in items[1:]:
        p = tf.add_paragraph()
        p.text = extra
        # p.level = 0 implicito; serve solo se vuoi sub-bullet (qui no).
    tf.word_wrap = True
    report.shapes_written += 1


class SlideBuilderV2:
    """XML-level slide builder via python-pptx shape iteration (FASE 3).

    Stesso contratto pubblico di SlideBuilder v1 (build / output_path) per
    drop-in replacement in ProductionBuilder.
    """

    def __init__(
        self,
        brand_config: dict[str, Any] | None = None,
        template_path: Path = DEFAULT_TEMPLATE,
        output_dir: Path = DEFAULT_OUTPUT_DIR,
        layout_map: dict[SlideType, int] | None = None,
    ) -> None:
        self.brand_config = brand_config or {}
        self.template_path = template_path
        self.output_dir = output_dir
        self.layout_map = layout_map or DEFAULT_LAYOUT_MAP

    # ─── Public API ───

    def build(
        self,
        slides: list[SlideContent],
        course: dict[str, Any],
        image_map: dict[int, str],
    ) -> str:
        """Build a PPTX from slides + image_map and return output path."""
        report = BuildReport()
        prs = Presentation(str(self.template_path))
        # FIX #27.7 (2026-05-26): svuota il testo placeholder DEGLI SHAPE NEI
        # LAYOUT (es. nx_title="Cosa hai imparato", nx_body, nx_recap_text_X...).
        # Causa bug "titolo sovrapposto": il layout mantiene i suoi shape con
        # testo placeholder, che PowerPoint EREDITA e mostra SOTTO il testo
        # scritto nella slide → due titoli sovrapposti. Svuotando i layout una
        # volta, l'eredità non porta testo. Le label decorative fisse
        # (RIEPILOGO, CASO STUDIO, SITUAZIONE...) restano (whitelist).
        _blank_layout_placeholder_text(prs)
        # Rimuovi le slide di esempio del template (se presenti)
        self._strip_template_example_slides(prs)
        total_pages = len(slides)

        # FIX #26 (2026-05-26): enumerate gives the GLOBAL position used as the
        # image_map key — slide.index is module-local and collides. page_num
        # also uses pos (slide.index repeated per module gave duplicate page
        # numbers like "32 / 493" fourteen times).
        for pos, slide in enumerate(slides):
            self._render_slide(prs, pos, slide, image_map, total_pages, report)

        output_path = self._compute_output_path(course)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        logger.info(
            "pptx_built_v2",
            output=output_path,
            slides=report.slides_built,
            shapes_written=report.shapes_written,
            images_inserted=report.images_inserted,
            image_fallbacks=report.image_fallbacks,
            warnings=len(report.warnings),
        )
        return output_path

    # ─── Private helpers ───

    def _strip_template_example_slides(self, prs: Any) -> None:
        """Rimuove le eventuali slide di esempio dal template (non i layout)."""
        # nexus_master.pptx ha 0 slide e 8 layout — niente da rimuovere.
        # Ma se in futuro il template avrà slide demo, vanno via QUI.
        xml_slides = prs.slides._sldIdLst  # type: ignore[attr-defined]
        for sld_id in list(xml_slides):
            xml_slides.remove(sld_id)

    def _render_slide(
        self,
        prs: Any,
        pos: int,
        slide_content: SlideContent,
        image_map: dict[int, str],
        total_pages: int,
        report: BuildReport,
    ) -> None:
        layout_idx = self.layout_map.get(slide_content.slide_type)
        if layout_idx is None or layout_idx >= len(prs.slide_layouts):
            layout_idx = 1  # fallback CONTENT_TEXT
            report.warnings.append(
                f"slide {slide_content.index}: no layout for {slide_content.slide_type.value}, "
                f"falling back to CONTENT_TEXT"
            )
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        report.slides_built += 1

        # ─── CRITICAL (FASE 3 + FIX #20 2026-05-25): clona SOLO le AUTO_SHAPE
        # testuali dal layout, NON le PICTURE (logo brand). Il logo resta come
        # eredità del layer layout/master e viene renderizzato automaticamente.
        # showMasterSp resta default (=1) → eredità ON, testo scritto sulle
        # nostre shape clonate copre il testo placeholder originale del layout.
        self._clone_layout_text_shapes_only(slide, layout)

        # ─── Speaker notes (sempre, per audio TTS) ───
        if slide_content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = slide_content.speaker_notes

        page_num = f"{pos + 1} / {total_pages}"

        # ─── Branch per SlideType ───
        st = slide_content.slide_type
        if st == SlideType.TITLE:
            _write_shape_by_role(slide, st, "title", slide_content.title, report)
            # Subtitle: dal metadata corso (titolo + ore + normativa), NON dall'LLM
            # (TITLE non ha body per design FASE 1). Se normative_ref valorizzato,
            # lo usiamo come subtitle informativo.
            if slide_content.normative_ref and slide_content.normative_ref != "—":
                _write_shape_by_role(
                    slide, st, "subtitle", slide_content.normative_ref, report
                )

        elif st == SlideType.CONTENT_TEXT:
            # FIX #28.3 (2026-05-26): title + body via placeholder veri (v3 template).
            # body = list[str], scritto come paragrafi separati → bullet ereditato.
            _write_shape_by_role(slide, st, "title", slide_content.title, report)
            _write_shape_by_role_bullets(slide, st, "body", slide_content.bullets, report)
            if slide_content.normative_ref:
                _write_shape_by_role(
                    slide, st, "normative_ref", slide_content.normative_ref, report
                )
            _write_shape_by_role(slide, st, "page_num", page_num, report)

        elif st == SlideType.RECAP:
            # FIX #30.0-sexies (2026-05-26, conferma analista P1): RECAP è
            # un CONTENT_TEXT con bullet ✓ invece di •. Layout v4 ha singolo
            # BODY placeholder (idx=1, nx_recap_body) con lstStyle/lvl1pPr/
            # buChar="✓" ereditato. Stesso writer di CONTENT_TEXT: 1 paragrafo
            # per bullet → ereditarietà del glifo applicata per paragrafo.
            # MAI joinare con \n (sarebbe soft line break = 1 paragrafo solo =
            # 1 solo checkmark).
            _write_shape_by_role(slide, st, "title", slide_content.title, report)
            _write_shape_by_role_bullets(slide, st, "body", slide_content.bullets, report)
            _write_shape_by_role(slide, st, "page_num", page_num, report)

        elif st == SlideType.CONTENT_IMAGE:
            # FIX #28.3 (2026-05-26): title + body via placeholder veri (v3).
            # bullets:list[str] scritti come paragrafi → bullet ereditato.
            _write_shape_by_role(slide, st, "title", slide_content.title, report)
            _write_shape_by_role_bullets(slide, st, "body", slide_content.bullets, report)
            _write_shape_by_role(
                slide, st, "image_caption",
                _clean_caption(slide_content.image.query),
                report,
            )
            _write_shape_by_role(
                slide, st, "normative_ref", slide_content.normative_ref, report
            )
            _write_shape_by_role(slide, st, "page_num", page_num, report)
            # Inserimento immagine (se path locale disponibile)
            self._maybe_insert_image(slide, pos, image_map, report)

        elif st == SlideType.DIAGRAM:
            # FIX #28.3 (2026-05-26): DIAGRAM è body-less per design (analista).
            # La didascalia 1-2 bullet vive in nx_caption AUTO_SHAPE — NESSUN bullet
            # ereditato perché caption descrittiva ≠ body multi-bullet.
            _write_shape_by_role(slide, st, "title", slide_content.title, report)
            # Compatibilità: bullets:list[str] → caption come testo concatenato.
            # (1-2 bullet brevi joinati da '. ' per leggibilità didascalia.)
            caption_text = ". ".join(b.rstrip(".") for b in (slide_content.bullets or []) if b.strip())
            _write_shape_by_role(slide, st, "diagram_caption", caption_text, report)
            _write_shape_by_role(
                slide, st, "normative_ref", slide_content.normative_ref, report
            )
            _write_shape_by_role(slide, st, "page_num", page_num, report)
            # Inserimento PNG diagramma renderizzato (se disponibile)
            self._maybe_insert_image(slide, pos, image_map, report)

        elif st == SlideType.QUIZ:
            _write_shape_by_role(slide, st, "title", slide_content.title, report)
            options = slide_content.quiz_options or []
            roles = ["option_a", "option_b", "option_c", "option_d"]
            for i, role in enumerate(roles):
                if i < len(options):
                    letter = chr(ord("A") + i)
                    _write_shape_by_role(slide, st, role, f"{letter}. {options[i]}", report)
            # Marker risposta corretta — DISABILITATO (testo vuoto + spostato
            # fuori dalla slide). Bug pre-fix: testo "Risposta corretta: A"
            # (21 chars) dentro shape 23px causava wrap char-per-char in
            # verticale, producendo una striscia visibile attraverso le opzioni
            # quiz. Decisione: slide quiz pulita senza marker visivo
            # (l'info quiz_correct resta nel meta_json per consumer downstream).
            if slide_content.quiz_correct is not None and 0 <= slide_content.quiz_correct < 4:
                _write_shape_by_role(slide, st, "correct_marker", "", report)
                from pptx.util import Emu  # local import (pattern coerente con riga ~770)
                # Sposta il marker FUORI dalla slide (left/top negativi)
                for shape in slide.shapes:
                    if shape.name == "nx_correct_marker":
                        shape.left = Emu(-100000)
                        shape.top = Emu(-100000)
                        break

        elif st == SlideType.CASE_STUDY:
            # FIX #28.3 (2026-05-26): sezioni è VERA list[str] (3 elementi via
            # validator). Niente più parsing del body — distribuzione 1-a-1 sulle
            # 3 shape fisiche (nx_situazione/azione/risultato, AUTO_SHAPE come
            # da decisione analista: marker = label decorative, non bullet).
            _write_shape_by_role(slide, st, "title", slide_content.title, report)
            roles_3 = ["situazione", "azione", "risultato"]
            sezioni = slide_content.sezioni or []
            for i, role in enumerate(roles_3):
                text = sezioni[i] if i < len(sezioni) else ""
                _write_shape_by_role(slide, st, role, text, report)
            _write_shape_by_role(
                slide, st, "normative_ref", slide_content.normative_ref, report
            )
            _write_shape_by_role(slide, st, "page_num", page_num, report)

        elif st == SlideType.MODULE_OPEN:
            # FIX #30.2 (2026-05-26): bookend apertura modulo. Titolo riformulato
            # come "MODULO N" + sotto-titolo (titolo del modulo). I valori sono
            # derivati da module_index + module title del contesto, popolati nel
            # content_agent al momento dell'emissione bookend slot.
            module_num_text = slide_content.title or f"MODULO {slide_content.module_index + 1}"
            _write_shape_by_role(slide, st, "module_num", module_num_text, report)
            # Il "titolo modulo" reale lo mettiamo nel sotto-titolo via bullets[0]
            # se presente (convention); altrimenti vuoto.
            module_title_text = (slide_content.bullets[0] if slide_content.bullets else "")
            _write_shape_by_role(slide, st, "module_title", module_title_text, report)

        elif st == SlideType.MODULE_CLOSE:
            # FIX #30.2 (2026-05-26): bookend chiusura modulo. Riusa pattern
            # RECAP: 1 placeholder body con 5 ✓ ereditati dal template v4
            # (nx_recap_body lstStyle buChar="✓").
            close_num_text = f"RIEPILOGO MODULO {slide_content.module_index + 1}"
            _write_shape_by_role(slide, st, "module_num", close_num_text, report)
            _write_shape_by_role(slide, st, "module_title", slide_content.title, report)
            _write_shape_by_role_bullets(slide, st, "body", slide_content.bullets, report)

        elif st == SlideType.CLOSING:
            _write_shape_by_role(slide, st, "title", slide_content.title or "Grazie", report)

        # FIX #27.2: garanzia anti-placeholder universale. Dopo aver scritto i
        # ruoli, azzera ogni text_frame ancora col testo placeholder del template
        # (es. "Testo sezione azione." non sovrascritto, "[ AUTO_SHAPE... ]" del
        # box immagine non riempito). Slide sempre pulita.
        _blank_unwritten_placeholders(slide)

    def _clone_layout_text_shapes_only(self, slide: Any, layout: Any) -> None:
        """Deep-copy SOLO gli AUTO_SHAPE testuali del layout — NO PICTURE.

        FIX #20 (2026-05-25): rispetto al precedente _clone_layout_shapes:
        - PICTURE shapes (logo brand 'Image 0') NON vengono clonate.
          Il logo resta come eredità del layer layout/master, visibile
          automaticamente perché showMasterSp default = 1.
        - Solo <p:sp> (AutoShape testuali) vengono clonate. Quando scriviamo
          testo via SHAPE_MAP indici, sovrascriviamo il testo placeholder del
          layout senza toccare il logo.

        Risultato: testo per-slide editabile + logo brand sempre visibile +
        banda rosa/footer fissi del layout sempre visibili.
        """
        from pptx.oxml.ns import qn

        slide_spTree = slide.shapes._spTree  # type: ignore[attr-defined]
        layout_spTree = layout.shapes._spTree  # type: ignore[attr-defined]

        # Rimuovi solo gli <p:sp> già copiati da add_slide (placeholders).
        # NON tocchiamo <p:pic> esistenti (non dovrebbero essercene, ma safe).
        for existing in list(slide_spTree):
            if existing.tag == qn("p:sp"):
                slide_spTree.remove(existing)

        # Clona SOLO <p:sp> dal layout. Le <p:pic> (logo) restano ereditate.
        for shape_el in layout_spTree:
            if shape_el.tag == qn("p:sp"):
                slide_spTree.append(copy.deepcopy(shape_el))

    def _maybe_insert_image(
        self,
        slide: Any,
        pos: int,
        image_map: dict[int, str],
        report: BuildReport,
    ) -> None:
        """FASE 4: inserisce l'immagine in image_map[pos] nel box dedicato.

        FIX #30.0-sexies (2026-05-26, conferma analista P3): pattern preferenziale
        è `slide.placeholders[idx].insert_picture(path)` per PICTURE placeholder
        veri (type=18) — consuma il placeholder e restituisce PlaceholderPicture,
        center-crop al box. Fallback su nome canonico (nx_image_box/nx_diagram_box)
        per layout vecchi senza placeholder PICTURE, fallback finale euristica.
        """
        path = image_map.get(pos)
        if not _is_local_path(path):
            logger.warning("image_no_path", pos=pos, in_map=pos in image_map)
            report.image_fallbacks += 1
            return
        assert path is not None

        # FIX #30.0-septies (2026-05-26): path da image_map è il file ORIGINALE
        # (output/images/<uuid>.png) che il cleanup_tmp di production_builder
        # cancella dopo 1h. Il file persistente è la versione fit-to-box
        # (<uuid>_fitted.png). Provo direttamente quella se l'originale non c'è
        # (ipotesi 1: originale già rimosso, fit già fatto in rebuild precedente).
        # Se originale ESISTE: lo uso per fit_image_to_box che genera/aggiorna
        # il _fitted. Output usato per insert_picture: SEMPRE il _fitted.
        import os
        original_exists = os.path.exists(path)
        fitted_candidate = path.rsplit(".", 1)[0] + "_fitted.png"
        fitted_exists = os.path.exists(fitted_candidate)

        if not original_exists and not fitted_exists:
            logger.warning("image_file_missing", pos=pos, original=path,
                           fitted=fitted_candidate)
            report.image_fallbacks += 1
            return

        from pptx.enum.shapes import PP_PLACEHOLDER
        from pptx.util import Emu

        from app.services.image_service import fit_image_to_box

        # PATH 1 (preferito): PICTURE placeholder type=18 via slide.placeholders.
        # CONTENT_IMAGE (layout 3): idx=2  |  DIAGRAM (layout 4): idx=2
        picture_ph = None
        for ph in slide.placeholders:
            try:
                if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                    picture_ph = ph
                    break
            except Exception:
                continue

        if picture_ph is not None:
            # Usa _fitted se esiste, altrimenti fai fit ora dall'originale
            insert_path = fitted_candidate if fitted_exists else path
            try:
                picture_ph.insert_picture(insert_path)
                report.images_inserted += 1
                return
            except Exception as exc:
                logger.warning("placeholder_insert_picture_failed",
                               pos=pos, error=str(exc), path=insert_path,
                               fallback="named_box")
                # Continua sui fallback path

        # PATH 2 (fallback): AUTO_SHAPE per nome canonico (template legacy).
        _IMAGE_BOX_NAMES = ("nx_image_box", "nx_diagram_box", "nx_image", "nx_diagram")
        best_box = None
        for sh in slide.shapes:
            if (sh.name or "") in _IMAGE_BOX_NAMES:
                best_box = sh
                break

        # PATH 3 (fallback finale): euristica shape più grande senza testo.
        if best_box is None:
            best_area = 0
            for sh in slide.shapes:
                try:
                    w = int(sh.width or 0)
                    h = int(sh.height or 0)
                except Exception:
                    continue
                area = w * h
                if sh.shape_type is not None and "PICTURE" in str(sh.shape_type):
                    continue
                if area < int(2.0 * 914400) ** 2:
                    continue
                if sh.has_text_frame:
                    text = sh.text_frame.text.strip()
                    if text and not (text.startswith("[") or text.startswith("(")):
                        continue
                if area > best_area:
                    best_area = area
                    best_box = sh

        if best_box is None:
            logger.warning("image_no_box", pos=pos, shapes=[sh.name for sh in slide.shapes])
            report.image_fallbacks += 1
            return

        pic_shape = best_box
        left, top = pic_shape.left, pic_shape.top
        width, height = pic_shape.width, pic_shape.height

        # Fit-to-box: se _fitted già presente lo riuso, altrimenti chiamo
        # fit_image_to_box (richiede originale ancora presente).
        if fitted_exists:
            fitted_path = fitted_candidate
        elif original_exists:
            box_w_px = max(1, int(width / 914400 * 96))
            box_h_px = max(1, int(height / 914400 * 96))
            fitted_path = fit_image_to_box(path, box_w_px, box_h_px)
        else:
            fitted_path = None

        if fitted_path is None:
            logger.warning("image_insert_no_fitted", pos=pos)
            report.image_fallbacks += 1
            return

        # Rimuovo il box e inserisco l'immagine fitted nella stessa posizione.
        try:
            pic_el = pic_shape._element  # type: ignore[attr-defined]
            pic_el.getparent().remove(pic_el)
            slide.shapes.add_picture(
                fitted_path, Emu(left), Emu(top), width=Emu(width), height=Emu(height)
            )
            report.images_inserted += 1
        except Exception as exc:
            logger.warning("image_insert_failed", pos=pos, error=str(exc))
            report.image_fallbacks += 1

    def _compute_output_path(self, course: dict[str, Any]) -> str:
        course_id_raw = str(course.get("id") or course.get("course_type") or "course")
        # Strip path separators per sicurezza
        course_id = course_id_raw.replace("/", "_").replace("\\", "_")
        return str(self.output_dir / f"{course_id}.pptx")
