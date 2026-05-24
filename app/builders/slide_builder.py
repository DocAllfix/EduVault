"""SlideBuilder — synchronous PPTX assembler (BLUEPRINT §07.1 / §07.3).

The class is built around three contracts:

1. ``build(slides, course, image_map) -> str`` is **synchronous**: callers
   (``ProductionBuilder``, FASE 4.5) wrap it in ``asyncio.to_thread``.
   The Semaphore(1) that guards python-pptx + lxml (REI-3) lives in
   ``generation_service`` (FASE 5.1) — NOT here.

2. ``image_map: dict[int, str]`` carries ONLY local filesystem paths
   (BP §07 line 2148). Anything that looks like a URL is treated as missing
   and triggers the textual fallback.

3. Image insertion is wrapped in try/except per BP §07.1 line 2301-2312:
   a failing insert logs ``pptx_image_insert_failed`` and writes
   ``[Immagine non disponibile]`` into the picture placeholder. The whole
   build NEVER crashes because one image is broken/missing.

The template ``assets/templates/nexus_master.pptx`` is a HUMAN-AUTHORED
binary asset (#R8). Its slide_layouts must include the 8 BP §07.3 layouts
in this order — overridable per-instance via ``layout_map`` for tests that
use a synthetic template.
"""

from __future__ import annotations

import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import structlog
from pptx import Presentation

from app.models.core import SlideType
from app.models.pipeline import SlideContent

logger = structlog.get_logger()

DEFAULT_TEMPLATE = Path("assets/templates/nexus_master.pptx")
DEFAULT_OUTPUT_DIR = Path("output")
IMAGE_MISSING_FALLBACK = "[Immagine non disponibile]"

# BP §07.3 layout order. Overridable via ``SlideBuilder(layout_map=...)``
# so tests can target python-pptx's default 11-layout template.
DEFAULT_LAYOUT_MAP: dict[SlideType, int] = {
    SlideType.TITLE: 0,
    SlideType.CONTENT_TEXT: 1,
    SlideType.CONTENT_IMAGE: 2,
    SlideType.DIAGRAM: 3,  # declassified FIX-8 v1.0 but kept for resilience
    SlideType.QUIZ: 4,
    SlideType.CASE_STUDY: 5,
    SlideType.RECAP: 6,
    SlideType.CLOSING: 7,
}


@dataclass
class BuildReport:
    """Per-build counters for diagnostics. Not the canonical GenerationReport —
    that one is composed by ProductionBuilder (FASE 4.5) from this + others."""

    slides_built: int = 0
    images_inserted: int = 0
    image_fallbacks: int = 0
    layout_fallbacks: int = 0
    warnings: list[str] = field(default_factory=list)


def _is_local_path(value: str | None) -> bool:
    """True iff ``value`` is a non-empty local path on disk.

    URLs (http/https/file://) are explicitly rejected: BP §07 line 2148
    requires that ``image_map`` only carries local paths produced by the
    Image Service (FASE 4.3).
    """
    if not value:
        return False
    lowered = value.lower()
    if lowered.startswith(("http://", "https://", "file://", "ftp://")):
        return False
    return os.path.isfile(value)


def _find_placeholder_by_type(slide: Any, ptype_names: tuple[str, ...]) -> Any | None:
    """Return the first placeholder whose ``placeholder_format.type.name``
    matches one of ``ptype_names`` (case-sensitive). ``None`` if none."""
    for ph in slide.placeholders:
        ptype = ph.placeholder_format.type
        name = getattr(ptype, "name", None)
        if name in ptype_names:
            return ph
    return None


def _set_text(placeholder: Any, text: str) -> None:
    """Replace text in a placeholder without losing the formatting of the
    first run (python-pptx convention)."""
    tf = placeholder.text_frame
    tf.text = text


def _set_speaker_notes(slide: Any, notes: str) -> None:
    if not notes:
        return
    slide.notes_slide.notes_text_frame.text = notes


def _insert_image_with_fallback(
    slide: Any,
    image_path: str | None,
    report: BuildReport,
    slide_index: int,
) -> None:
    """Insert ``image_path`` into the first PICTURE placeholder.

    Behavior per BP §07.1 line 2301-2312:
    - missing/non-local path → write IMAGE_MISSING_FALLBACK into a body
      placeholder (or skip if no suitable placeholder exists)
    - insert fails (corrupt file, format mismatch) → same fallback, log warning
    """
    picture_ph = _find_placeholder_by_type(slide, ("PICTURE",))

    if not _is_local_path(image_path):
        report.image_fallbacks += 1
        logger.warning(
            "pptx_image_missing_or_remote", slide=slide_index, path=image_path
        )
        if picture_ph is not None:
            try:
                _set_text(picture_ph, IMAGE_MISSING_FALLBACK)
            except Exception:
                pass
        return

    if picture_ph is None:
        report.image_fallbacks += 1
        report.warnings.append(
            f"slide {slide_index}: layout has no PICTURE placeholder"
        )
        return

    try:
        picture_ph.insert_picture(image_path)
        report.images_inserted += 1
    except Exception as exc:
        report.image_fallbacks += 1
        logger.warning(
            "pptx_image_insert_failed",
            slide=slide_index,
            path=image_path,
            error=str(exc),
        )
        try:
            _set_text(picture_ph, IMAGE_MISSING_FALLBACK)
        except Exception:
            pass


def _populate_title_and_body(slide: Any, slide_content: SlideContent) -> None:
    title_ph = _find_placeholder_by_type(slide, ("TITLE", "CENTER_TITLE"))
    if title_ph is not None:
        _set_text(title_ph, slide_content.title)

    body_ph = _find_placeholder_by_type(slide, ("BODY", "OBJECT", "SUBTITLE"))
    if body_ph is not None and slide_content.body:
        _set_text(body_ph, slide_content.body)


def _populate_quiz(slide: Any, slide_content: SlideContent) -> None:
    """QUIZ layout: title + question (body) + up to 4 options.

    Options are appended to the body placeholder (one per line, prefixed with
    A./B./C./D.) with the correct answer marked. BP §07.3 LAYOUT 4 reserves
    4 distinct option placeholders, but we fall back to a single body
    placeholder when those are not present (which is the case in tests
    using python-pptx's default Title-and-Content layout).
    """
    _populate_title_and_body(slide, slide_content)

    if not slide_content.quiz_options:
        return

    correct_idx = slide_content.quiz_correct
    lines = []
    for i, opt in enumerate(slide_content.quiz_options):
        letter = chr(ord("A") + i)
        marker = " ✓" if i == correct_idx else ""
        lines.append(f"{letter}. {opt}{marker}")
    options_text = "\n".join(lines)

    body_ph = _find_placeholder_by_type(slide, ("BODY", "OBJECT"))
    if body_ph is not None:
        existing = body_ph.text_frame.text
        if existing:
            _set_text(body_ph, f"{existing}\n\n{options_text}")
        else:
            _set_text(body_ph, options_text)


class SlideBuilder:
    """Synchronous PPTX assembler. BP §07.1.

    The constructor opens the template once; ``build`` may be called many
    times (one per course) on the same instance because each call clones
    slide layouts into a fresh Presentation.
    """

    def __init__(
        self,
        brand_config: dict[str, Any] | None = None,
        template_path: Path = DEFAULT_TEMPLATE,
        output_dir: Path = DEFAULT_OUTPUT_DIR,
        layout_map: dict[SlideType, int] | None = None,
    ) -> None:
        if not template_path.is_file():
            raise FileNotFoundError(
                f"PPTX template not found at {template_path}. "
                "It is a human-authored binary asset — see BP §07.3 "
                "and Master Plan FASE 4 prerequisite (#R8)."
            )
        self.brand_config = brand_config or {}
        self.template_path = template_path
        self.output_dir = output_dir
        self.layout_map = layout_map or DEFAULT_LAYOUT_MAP

    def build(
        self,
        slides: list[SlideContent],
        course: dict[str, Any],
        image_map: dict[int, str],
    ) -> str:
        """Build the PPTX, write it to disk, return the absolute path.

        Synchronous on purpose (REI-3): ``ProductionBuilder`` wraps this in
        ``asyncio.to_thread`` under a single-permit semaphore.
        """
        prs = Presentation(str(self.template_path))
        report = BuildReport()
        max_layout_idx = len(prs.slide_layouts) - 1

        for slide_content in slides:
            layout_idx = self.layout_map.get(slide_content.slide_type, 1)
            if layout_idx > max_layout_idx:
                report.layout_fallbacks += 1
                report.warnings.append(
                    f"slide {slide_content.index}: layout {layout_idx} "
                    f"unavailable, falling back to 1"
                )
                layout_idx = min(1, max_layout_idx)
            layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(layout)

            stype = slide_content.slide_type
            if stype is SlideType.QUIZ:
                _populate_quiz(slide, slide_content)
            elif stype in (SlideType.CONTENT_IMAGE, SlideType.DIAGRAM):
                _populate_title_and_body(slide, slide_content)
                image_path = image_map.get(slide_content.index)
                _insert_image_with_fallback(
                    slide, image_path, report, slide_content.index
                )
            else:
                _populate_title_and_body(slide, slide_content)

            _set_speaker_notes(slide, slide_content.speaker_notes)
            report.slides_built += 1

        output_path = self._compute_output_path(course)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

        logger.info(
            "pptx_built",
            path=str(output_path),
            slides=report.slides_built,
            images_inserted=report.images_inserted,
            image_fallbacks=report.image_fallbacks,
            layout_fallbacks=report.layout_fallbacks,
        )
        return str(output_path)

    def _compute_output_path(self, course: dict[str, Any]) -> Path:
        course_id = course.get("id") or course.get("course_id") or "course"
        safe_id = str(course_id).replace(os.sep, "_").replace("/", "_")
        return self.output_dir / f"{safe_id}_corso.pptx"


__all__ = [
    "DEFAULT_LAYOUT_MAP",
    "DEFAULT_OUTPUT_DIR",
    "DEFAULT_TEMPLATE",
    "IMAGE_MISSING_FALLBACK",
    "BuildReport",
    "SlideBuilder",
]
