"""Integration tests — FASE 3 vast-hopping-sketch — SlideBuilderV2.

Verifica che SlideBuilderV2 (XML find/replace via deep-copy layout shapes)
scriva EFFETTIVAMENTE il testo LLM-generated nelle slide, risolvendo il GAP
shape_map documentato nel test E.03/E.04 (SlideBuilder v1 placeholder API
falliva sul template Claude Design AUTO_SHAPE).

Il test critico ``test_text_actually_present_all_types`` è la regression suite
del fix: ogni SlideType deve avere il proprio marker visibile nel PPTX riaperto.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from app.builders.slide_builder_v2 import SlideBuilderV2
from app.models.core import SlideType
from tests._helpers import make_slide

TEMPLATE = Path("assets/templates/nexus_master.pptx")

pytestmark = pytest.mark.skipif(
    not TEMPLATE.is_file(), reason="nexus_master.pptx template not present"
)


def _all_text(pptx_path: str) -> str:
    prs = Presentation(pptx_path)
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return " | ".join(parts)


def test_v2_writes_pptx_with_n_slides(tmp_path: Path) -> None:
    builder = SlideBuilderV2(template_path=TEMPLATE, output_dir=tmp_path)
    slides = [
        make_slide(SlideType.CONTENT_TEXT, index=0, title="Slide uno"),
        make_slide(SlideType.CONTENT_TEXT, index=1, title="Slide due"),
    ]
    out = builder.build(slides=slides, course={"id": "v2-count"}, image_map={})
    prs = Presentation(out)
    assert len(prs.slides) == 2


def test_v2_text_actually_present_all_types(tmp_path: Path) -> None:
    """REGRESSION del GAP E.04: il testo LLM DEVE essere visibile nel PPTX.

    Ogni SlideType ha un marker unico; tutti devono comparire nel pacchetto XML.
    """
    builder = SlideBuilderV2(template_path=TEMPLATE, output_dir=tmp_path)
    slides = [
        make_slide(SlideType.TITLE, index=0, title="MARKER_TITLE",
                   normative_ref="Corso 8 ore DM 388"),
        make_slide(SlideType.CONTENT_TEXT, index=1, title="MARKER_CT",
                   body="Bullet uno MARKER_CTB\nBullet due\nBullet tre"),
        make_slide(SlideType.QUIZ, index=2, title="MARKER_QUIZ domanda?",
                   quiz_options=["Opz A", "MARKER_OPT_B", "Opz C", "Opz D"],
                   quiz_correct=1),
        make_slide(SlideType.CASE_STUDY, index=3, title="MARKER_CASE",
                   body="Situazione MARKER_SIT breve --- Azione MARKER_AZ breve "
                        "--- Risultato MARKER_RIS breve"),
        make_slide(SlideType.RECAP, index=4, title="MARKER_RECAP",
                   body="Punto uno\nPunto due\nPunto tre"),
        make_slide(SlideType.CLOSING, index=5, title="Grazie MARKER_CLOSE"),
    ]
    out = builder.build(slides=slides, course={"id": "v2-markers"}, image_map={})
    combined = _all_text(out)

    for marker in (
        "MARKER_TITLE", "MARKER_CT", "MARKER_CTB", "MARKER_QUIZ",
        "MARKER_OPT_B", "MARKER_SIT", "MARKER_AZ", "MARKER_RIS",
        "MARKER_RECAP", "MARKER_CLOSE",
    ):
        assert marker in combined, f"{marker} NOT found in PPTX — GAP shape_map regression!"


def test_v2_quiz_marks_correct_answer(tmp_path: Path) -> None:
    builder = SlideBuilderV2(template_path=TEMPLATE, output_dir=tmp_path)
    slides = [
        make_slide(SlideType.QUIZ, index=0, title="Quale DPI per la testa?",
                   quiz_options=["Guanti", "Casco", "Cuffie", "Scarpe"],
                   quiz_correct=1),
    ]
    out = builder.build(slides=slides, course={"id": "v2-quiz"}, image_map={})
    combined = _all_text(out)
    assert "B. Casco" in combined
    assert "Risposta corretta: B" in combined


def test_v2_writes_speaker_notes(tmp_path: Path) -> None:
    builder = SlideBuilderV2(template_path=TEMPLATE, output_dir=tmp_path)
    notes = " ".join(["nota"] * 80)
    slide = make_slide(SlideType.CONTENT_TEXT, index=0, speaker_notes=notes)
    out = builder.build(slides=[slide], course={"id": "v2-notes"}, image_map={})
    prs = Presentation(out)
    assert "nota" in prs.slides[0].notes_slide.notes_text_frame.text
