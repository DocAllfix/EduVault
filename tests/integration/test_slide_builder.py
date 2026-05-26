"""SlideBuilder tests (FASE 4.2).

Uses a synthetic .pptx generated on the fly (python-pptx's default 11-layout
template) — the real ``assets/templates/nexus_master.pptx`` is human-authored
(#R8) and not part of the repo. Slide content is built directly with the
Pydantic models from ``app.models.pipeline`` so we exercise the same
contract the Content Agent (FASE 3.4) emits.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from app.builders.slide_builder import (
    DEFAULT_LAYOUT_MAP,
    IMAGE_MISSING_FALLBACK,
    SlideBuilder,
    _is_local_path,
)
from app.models.core import SlideType
from app.models.pipeline import ImageStrategy, SlideContent


# ─────────────── fixtures ───────────────


@pytest.fixture
def synthetic_template(tmp_path: Path) -> Path:
    """A 16:9 default Presentation saved to disk for SlideBuilder to consume."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    path = tmp_path / "synthetic_master.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def output_dir(tmp_path: Path) -> Path:
    out = tmp_path / "output"
    out.mkdir()
    return out


@pytest.fixture
def real_image(tmp_path: Path) -> Path:
    """Tiny valid PNG that python-pptx can insert without errors."""
    img = Image.new("RGB", (64, 64), color="red")
    path = tmp_path / "real.png"
    img.save(str(path), "PNG")
    return path


@pytest.fixture
def corrupt_image(tmp_path: Path) -> Path:
    """A file with a .png extension but garbage bytes — insert_picture fails."""
    path = tmp_path / "corrupt.png"
    path.write_bytes(b"not really a png")
    return path


# ─────────────── synthetic slide map ───────────────


# The default python-pptx template has these layouts:
#   0=Title Slide, 1=Title and Content, 2=Section Header,
#   3=Two Content, 4=Comparison, 5=Title Only, 6=Blank,
#   7=Content with Caption, 8=Picture with Caption, 9=..., 10=...
# We map our SlideType to layouts that have BOTH a title placeholder
# AND (for CONTENT_IMAGE/DIAGRAM) a PICTURE placeholder.
TEST_LAYOUT_MAP = {
    SlideType.TITLE: 0,  # Title Slide (CENTER_TITLE + SUBTITLE)
    SlideType.CONTENT_TEXT: 1,  # Title and Content (TITLE + OBJECT)
    SlideType.CONTENT_IMAGE: 8,  # Picture with Caption (TITLE + PICTURE + BODY)
    SlideType.DIAGRAM: 8,
    SlideType.QUIZ: 1,
    SlideType.CASE_STUDY: 1,
    SlideType.RECAP: 1,
    SlideType.CLOSING: 5,  # Title Only
}


def _slide(
    index: int,
    stype: SlideType,
    title: str = "Title",
    body: str | None = None,
    *,
    speaker_notes: str | None = None,
    quiz_options: list[str] | None = None,
    quiz_correct: int | None = None,
) -> SlideContent:
    """FASE 1 vast-hopping: delega a make_slide centralizzato per essere
    automaticamente conforme a LAYOUT_CONSTRAINTS (no soft-truncate)."""
    from tests._helpers import make_slide

    overrides: dict[str, object] = {
        "index": index,
        "module_index": 0,
        "normative_ref": "Art. 1, DM 388/2003",
        "source_chunk_ids": [],
    }
    # Title sotto i 70 char default
    if title and title != "Title":
        overrides["title"] = title
    if body is not None:
        overrides["body"] = body
    if speaker_notes is not None and speaker_notes != "":
        overrides["speaker_notes"] = speaker_notes
    if quiz_options is not None:
        overrides["quiz_options"] = quiz_options
    if quiz_correct is not None:
        overrides["quiz_correct"] = quiz_correct
    return make_slide(stype, **overrides)


def _builder(
    template: Path, output_dir: Path, layout_map: dict[SlideType, int] | None = None
) -> SlideBuilder:
    return SlideBuilder(
        brand_config={},
        template_path=template,
        output_dir=output_dir,
        layout_map=layout_map or TEST_LAYOUT_MAP,
    )


# ─────────────── 1. construction ───────────────


def test_default_layout_map_covers_all_slide_types() -> None:
    """Every SlideType the Content Agent can emit must have a layout idx."""
    for stype in SlideType:
        assert stype in DEFAULT_LAYOUT_MAP


def test_constructor_raises_if_template_missing(tmp_path: Path) -> None:
    missing = tmp_path / "nope.pptx"
    with pytest.raises(FileNotFoundError, match="#R8"):
        SlideBuilder(template_path=missing, output_dir=tmp_path)


# ─────────────── 2. _is_local_path helper ───────────────


def test_is_local_path_accepts_existing_file(real_image: Path) -> None:
    assert _is_local_path(str(real_image)) is True


def test_is_local_path_rejects_none_and_empty() -> None:
    assert _is_local_path(None) is False
    assert _is_local_path("") is False


def test_is_local_path_rejects_urls(real_image: Path) -> None:
    # Even if the file exists locally, a URL-shaped string is rejected
    for url in ("http://example.com/x.png", "https://x/y.png", "file:///tmp/x.png"):
        assert _is_local_path(url) is False


def test_is_local_path_rejects_nonexistent_path(tmp_path: Path) -> None:
    assert _is_local_path(str(tmp_path / "ghost.png")) is False


# ─────────────── 3. happy paths ───────────────


def test_build_writes_pptx_to_output_dir(
    synthetic_template: Path, output_dir: Path
) -> None:
    builder = _builder(synthetic_template, output_dir)
    slides = [_slide(0, SlideType.TITLE, title="Corso Primo Soccorso")]
    course = {"id": "course-abc-123"}

    out = builder.build(slides, course, image_map={})

    out_path = Path(out)
    assert out_path.exists()
    assert out_path.parent == output_dir
    assert out_path.name == "course-abc-123_corso.pptx"
    # The output file must be openable as a real PPTX
    prs = Presentation(str(out_path))
    assert len(prs.slides) == 1


def test_build_handles_all_slide_types(
    synthetic_template: Path, output_dir: Path
) -> None:
    builder = _builder(synthetic_template, output_dir)
    slides = [_slide(i, stype) for i, stype in enumerate(SlideType)]
    course = {"id": "multi-type"}

    out = builder.build(slides, course, image_map={})

    prs = Presentation(out)
    assert len(prs.slides) == len(SlideType)


def test_build_populates_title_and_body(
    synthetic_template: Path, output_dir: Path
) -> None:
    builder = _builder(synthetic_template, output_dir)
    slides = [
        _slide(
            0,
            SlideType.CONTENT_TEXT,
            title="Obblighi del datore di lavoro",
            body="Il datore di lavoro garantisce la sicurezza.",
        )
    ]
    out = builder.build(slides, {"id": "title-body"}, image_map={})

    prs = Presentation(out)
    s = prs.slides[0]
    texts = [ph.text_frame.text for ph in s.placeholders if ph.has_text_frame]
    assert any("Obblighi del datore di lavoro" in t for t in texts)
    assert any("Il datore di lavoro garantisce la sicurezza." in t for t in texts)


def test_build_writes_speaker_notes(
    synthetic_template: Path, output_dir: Path
) -> None:
    builder = _builder(synthetic_template, output_dir)
    # FASE 1: speaker_notes valido per CONTENT_TEXT (75-90 parole) + contiene "preposto"
    notes = (
        "Soffermarsi sul ruolo strategico del preposto in azienda secondo quanto "
        "previsto dall'articolo trentasette del decreto legislativo ottantuno del "
        "duemilaotto. Il preposto deve vigilare sull'attuazione concreta delle "
        "misure di sicurezza, segnalare prontamente i comportamenti non conformi "
        "tra i lavoratori, dare istruzioni operative chiare e specifiche per "
        "ciascuna mansione, comunicare immediatamente al dirigente le deficienze "
        "tecniche rilevate sui mezzi e sulle attrezzature. Tali obblighi sono "
        "sanzionati penalmente in caso di grave omissione, anche solo reiterata "
        "nel tempo. Quindi il ruolo del preposto è fondamentale."
    )
    slides = [_slide(0, SlideType.CONTENT_TEXT, speaker_notes=notes)]
    out = builder.build(slides, {"id": "notes"}, image_map={})

    prs = Presentation(out)
    notes_text = prs.slides[0].notes_slide.notes_text_frame.text
    assert "preposto" in notes_text


# ─────────────── 4. image insertion + fallback ───────────────


def test_build_inserts_local_image_successfully(
    synthetic_template: Path, output_dir: Path, real_image: Path
) -> None:
    builder = _builder(synthetic_template, output_dir)
    slides = [_slide(0, SlideType.CONTENT_IMAGE)]
    image_map = {0: str(real_image)}

    out = builder.build(slides, {"id": "img-ok"}, image_map=image_map)

    prs = Presentation(out)
    s = prs.slides[0]
    # A PICTURE placeholder that successfully received insert_picture
    # becomes a PlaceholderPicture with a non-empty image
    has_image = any(
        ph.placeholder_format.type.name == "PICTURE"
        and getattr(ph, "image", None) is not None
        for ph in s.placeholders
    )
    assert has_image, "expected at least one PICTURE placeholder with image data"


def test_build_falls_back_when_image_path_missing(
    synthetic_template: Path, output_dir: Path
) -> None:
    builder = _builder(synthetic_template, output_dir)
    slides = [_slide(0, SlideType.CONTENT_IMAGE)]
    # image_map is empty → no path for slide 0
    out = builder.build(slides, {"id": "img-missing"}, image_map={})

    prs = Presentation(out)
    s = prs.slides[0]
    all_text = " ".join(
        ph.text_frame.text for ph in s.placeholders if ph.has_text_frame
    )
    assert IMAGE_MISSING_FALLBACK in all_text


def test_build_falls_back_when_image_is_url(
    synthetic_template: Path, output_dir: Path
) -> None:
    builder = _builder(synthetic_template, output_dir)
    slides = [_slide(0, SlideType.CONTENT_IMAGE)]
    image_map = {0: "https://example.com/x.png"}  # URL — NOT a local path

    out = builder.build(slides, {"id": "img-url"}, image_map=image_map)

    prs = Presentation(out)
    s = prs.slides[0]
    all_text = " ".join(
        ph.text_frame.text for ph in s.placeholders if ph.has_text_frame
    )
    assert IMAGE_MISSING_FALLBACK in all_text


def test_build_falls_back_when_image_is_corrupt(
    synthetic_template: Path, output_dir: Path, corrupt_image: Path
) -> None:
    builder = _builder(synthetic_template, output_dir)
    slides = [_slide(0, SlideType.CONTENT_IMAGE)]
    image_map = {0: str(corrupt_image)}

    # insert_picture raises on bad image data — must be caught, build proceeds
    out = builder.build(slides, {"id": "img-corrupt"}, image_map=image_map)

    prs = Presentation(out)
    s = prs.slides[0]
    all_text = " ".join(
        ph.text_frame.text for ph in s.placeholders if ph.has_text_frame
    )
    assert IMAGE_MISSING_FALLBACK in all_text


def test_build_ignores_image_map_for_non_image_slide_types(
    synthetic_template: Path, output_dir: Path, real_image: Path
) -> None:
    """A CONTENT_TEXT slide must NOT pick up an image even if image_map
    contains its index — image insertion is gated on slide_type."""
    builder = _builder(synthetic_template, output_dir)
    slides = [_slide(0, SlideType.CONTENT_TEXT, body="Solo testo, nessuna immagine.")]
    image_map = {0: str(real_image)}

    out = builder.build(slides, {"id": "no-image"}, image_map=image_map)

    prs = Presentation(out)
    s = prs.slides[0]
    has_image = any(
        ph.placeholder_format.type.name == "PICTURE"
        and getattr(ph, "image", None) is not None
        for ph in s.placeholders
    )
    assert not has_image


# ─────────────── 5. quiz ───────────────


def test_build_quiz_renders_options_and_marks_correct(
    synthetic_template: Path, output_dir: Path
) -> None:
    builder = _builder(synthetic_template, output_dir)
    # FASE 1: QUIZ non ha body (è options-only). title=domanda, options=4.
    slides = [
        _slide(
            0,
            SlideType.QUIZ,
            title="Quante persone formano la squadra di primo soccorso?",
            quiz_options=["1", "2", "3 o piu", "Dipende dal RSPP"],
            quiz_correct=2,
        )
    ]
    out = builder.build(slides, {"id": "quiz-1"}, image_map={})

    prs = Presentation(out)
    s = prs.slides[0]
    body_text = " ".join(
        ph.text_frame.text for ph in s.placeholders if ph.has_text_frame
    )
    assert "A. 1" in body_text
    assert "B. 2" in body_text
    assert "C. 3 o piu" in body_text
    assert "D. Dipende dal RSPP" in body_text
    # The correct option (index 2 → C) is marked with a ✓
    assert "C. 3 o piu ✓" in body_text


def test_build_quiz_without_options_does_not_crash(
    synthetic_template: Path, output_dir: Path
) -> None:
    """FASE 1: il validator ora rigetta QUIZ senza options. Verifichiamo che
    una slide CONTENT_TEXT renderizzata sul layout QUIZ (mismatch) non crashi
    il SlideBuilder."""
    builder = _builder(synthetic_template, output_dir)
    # CONTENT_TEXT minimale (non QUIZ): il builder applica layout QUIZ comunque,
    # ma non crasha quando quiz_options=None.
    slides = [_slide(0, SlideType.CONTENT_TEXT, title="Domanda")]
    out = builder.build(slides, {"id": "quiz-empty"}, image_map={})
    prs = Presentation(out)
    assert len(prs.slides) == 1


# ─────────────── 6. layout out-of-range fallback ───────────────


def test_build_falls_back_to_layout_1_when_index_out_of_range(
    synthetic_template: Path, output_dir: Path
) -> None:
    """A layout_map pointing past the template's available layouts must
    fall back gracefully rather than crashing with IndexError."""
    over_map = dict(TEST_LAYOUT_MAP)
    over_map[SlideType.CONTENT_TEXT] = 99
    builder = _builder(synthetic_template, output_dir, layout_map=over_map)
    slides = [_slide(0, SlideType.CONTENT_TEXT)]

    out = builder.build(slides, {"id": "layout-oor"}, image_map={})
    prs = Presentation(out)
    assert len(prs.slides) == 1


# ─────────────── 7. output path safety ───────────────


def test_output_path_strips_separators_from_course_id(
    synthetic_template: Path, output_dir: Path
) -> None:
    builder = _builder(synthetic_template, output_dir)
    slides = [_slide(0, SlideType.TITLE)]
    # Path-separator characters in the id must not break out of output_dir
    course: dict[str, Any] = {"id": "ab/cd"}

    out = builder.build(slides, course, image_map={})
    out_path = Path(out)
    assert out_path.parent == output_dir
    assert "ab_cd_corso.pptx" == out_path.name


def test_output_path_falls_back_when_id_missing(
    synthetic_template: Path, output_dir: Path
) -> None:
    builder = _builder(synthetic_template, output_dir)
    slides = [_slide(0, SlideType.TITLE)]
    out = builder.build(slides, course={}, image_map={})
    assert Path(out).name == "course_corso.pptx"


# ─────────────── 8. resilience: image error does not abort the whole build ───────────────


def test_build_continues_after_image_failure(
    synthetic_template: Path,
    output_dir: Path,
    real_image: Path,
    corrupt_image: Path,
) -> None:
    """BP §07.1 invariant: one broken image does NOT kill the build."""
    builder = _builder(synthetic_template, output_dir)
    slides = [
        _slide(0, SlideType.CONTENT_IMAGE, title="OK"),
        _slide(1, SlideType.CONTENT_IMAGE, title="Broken"),
        _slide(2, SlideType.CONTENT_TEXT, title="After broken", body="Survives"),
    ]
    image_map = {0: str(real_image), 1: str(corrupt_image)}

    out = builder.build(slides, {"id": "resilient"}, image_map=image_map)

    prs = Presentation(out)
    assert len(prs.slides) == 3
    # Third slide must still have its title applied
    third_text = " ".join(
        ph.text_frame.text for ph in prs.slides[2].placeholders if ph.has_text_frame
    )
    assert "After broken" in third_text


# ─────────────── 9. structural smoke ───────────────


def test_image_map_only_accepts_local_paths_invariant() -> None:
    """Document BP §07 line 2148 constraint via a structural assertion:
    SlideBuilder MUST NOT accept anything other than local filesystem paths
    via image_map. URLs are explicitly rejected by _is_local_path.
    """
    for url in ("http://x", "https://x", "file:///x", "ftp://x"):
        assert _is_local_path(url) is False
    # And the public surface stays consistent
    assert IMAGE_MISSING_FALLBACK == "[Immagine non disponibile]"
