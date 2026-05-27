"""ProductionBuilder E2E tests (FASE 4.5).

PPTX is built REAL via SlideBuilder against a synthetic template (the human
template #R8 is not in the repo); PDF goes through Jinja2 REAL but WeasyPrint
is stubbed via ``sys.modules`` because GTK runtime is missing on Windows
(#R12, mirroring FASE 4.4 setup). 20 mock slides spanning 3 modules and
6 SlideType variants drive the happy path; pre-build guards and cleanup are
covered with focused unit-style tests.
"""

from __future__ import annotations

import os
import sys
import time
from pathlib import Path
from unittest.mock import AsyncMock, MagicMock, patch

import pytest
from pptx import Presentation
from pptx.util import Inches

# WeasyPrint stub before importing pdf_builder via production_builder.
if "weasyprint" not in sys.modules:  # pragma: no cover
    sys.modules["weasyprint"] = MagicMock()

from app.builders import production_builder as pb  # noqa: E402
from app.builders.pptx_validator import PptxValidator  # noqa: E402
from app.builders.production_builder import (  # noqa: E402
    CLEANUP_AGE_SECONDS,
    CLEANUP_PATTERNS,
    ProductionBuilder,
    check_disk_before_build,
    check_memory_before_build,
)
from app.builders.slide_builder import SlideBuilder  # noqa: E402
from app.models.core import SlideType  # noqa: E402
from app.models.pipeline import ImageStrategy, SlideContent  # noqa: E402


# ─────────────── shared fixtures ───────────────


@pytest.fixture
def synthetic_template(tmp_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    path = tmp_path / "synthetic_master.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def output_dir(tmp_path: Path) -> Path:
    out = tmp_path / "output"
    out.mkdir()
    return out


def _slide(
    index: int,
    *,
    module_index: int = 0,
    stype: SlideType = SlideType.CONTENT_TEXT,
    title: str | None = None,
    body: str | None = None,
) -> SlideContent:
    """FASE 1 vast-hopping: delega a make_slide centralizzato."""
    from tests._helpers import make_slide

    overrides: dict[str, object] = {
        "index": index,
        "module_index": module_index,
        "title": title or f"Slide {index}",
        "normative_ref": "Art. 1, DM 388/2003",
    }
    if body is not None:
        overrides["body"] = body
    return make_slide(stype, **overrides)


# Use a layout map valid for the python-pptx default template (mirrors
# test_slide_builder.TEST_LAYOUT_MAP).
TEST_LAYOUT_MAP = {
    SlideType.TITLE: 0,
    SlideType.CONTENT_TEXT: 1,
    SlideType.CONTENT_IMAGE: 8,
    SlideType.DIAGRAM: 8,
    SlideType.QUIZ: 1,
    SlideType.CASE_STUDY: 1,
    SlideType.RECAP: 1,
    SlideType.CLOSING: 5,
}


def _make_20_mock_slides() -> list[SlideContent]:
    """20 slides across 3 modules; types mixed but layout-compatible."""
    types = [
        SlideType.TITLE,
        SlideType.CONTENT_TEXT,
        SlideType.CONTENT_TEXT,
        SlideType.CONTENT_IMAGE,
        SlideType.QUIZ,
        SlideType.CASE_STUDY,
        SlideType.RECAP,
    ]
    slides: list[SlideContent] = []
    per_module = [7, 7, 6]
    idx = 0
    for module_index, count in enumerate(per_module):
        for j in range(count):
            stype = types[j % len(types)]
            slides.append(_slide(idx, module_index=module_index, stype=stype))
            idx += 1
    return slides


def _builder_with_synthetic_template(
    synthetic_template: Path, output_dir: Path
) -> ProductionBuilder:
    """A ProductionBuilder whose SlideBuilder points to the synthetic
    template (the real one #R8 is human-authored and absent)."""
    pb_inst = ProductionBuilder.__new__(ProductionBuilder)
    pb_inst.brand_config = {}
    pb_inst.slide_builder = SlideBuilder(
        template_path=synthetic_template,
        output_dir=output_dir,
        layout_map=TEST_LAYOUT_MAP,
    )
    from app.builders.pdf_builder import PdfBuilder

    pb_inst.pdf_builder = PdfBuilder(output_dir=output_dir)
    pb_inst.validator = PptxValidator()
    return pb_inst


# ─────────────── 1. pre-build guards ───────────────


def test_check_memory_passes_when_ample_ram() -> None:
    fake_mem = MagicMock(available=8 * 1024 * 1024 * 1024)  # 8GB
    with patch("app.builders.production_builder.psutil.virtual_memory", return_value=fake_mem):
        check_memory_before_build(100)


def test_check_memory_raises_when_estimated_exceeds_safety_ratio() -> None:
    # 100MB available → safety threshold = 60MB; 700 slides * 1.5MB = 1050MB
    fake_mem = MagicMock(available=100 * 1024 * 1024)
    with patch("app.builders.production_builder.psutil.virtual_memory", return_value=fake_mem):
        with pytest.raises(MemoryError, match="RAM insufficient"):
            check_memory_before_build(700)


def test_check_disk_passes_when_enough_space(tmp_path: Path) -> None:
    fake_usage = MagicMock(free=5 * 1024 * 1024 * 1024)  # 5GB
    with patch("app.builders.production_builder.shutil.disk_usage", return_value=fake_usage):
        check_disk_before_build(output_dir=tmp_path)


def test_check_disk_raises_when_below_threshold(tmp_path: Path) -> None:
    fake_usage = MagicMock(free=500 * 1024 * 1024)  # 500MB
    with patch("app.builders.production_builder.shutil.disk_usage", return_value=fake_usage):
        with pytest.raises(IOError, match="Insufficient disk space"):
            check_disk_before_build(output_dir=tmp_path)


# ─────────────── 2. PptxValidator ───────────────


def test_validator_returns_valid_when_count_matches(
    synthetic_template: Path, output_dir: Path
) -> None:
    builder = SlideBuilder(
        template_path=synthetic_template,
        output_dir=output_dir,
        layout_map=TEST_LAYOUT_MAP,
    )
    slides = [_slide(i) for i in range(3)]
    pptx_path = builder.build(slides, {"id": "v-ok"}, image_map={})

    result = PptxValidator().validate(pptx_path, slides)
    assert result["valid"] is True
    assert result["slide_count"] == 3
    assert result["warnings"] == []


def test_validator_flags_count_mismatch(
    synthetic_template: Path, output_dir: Path
) -> None:
    builder = SlideBuilder(
        template_path=synthetic_template,
        output_dir=output_dir,
        layout_map=TEST_LAYOUT_MAP,
    )
    slides = [_slide(i) for i in range(3)]
    pptx_path = builder.build(slides, {"id": "v-mismatch"}, image_map={})

    # Lie about the expected count: validator should report mismatch.
    extra = slides + [_slide(99)]
    result = PptxValidator().validate(pptx_path, extra)
    assert result["valid"] is False
    assert any("slide_count_mismatch" in w for w in result["warnings"])


def test_validator_handles_missing_pptx(tmp_path: Path) -> None:
    ghost = tmp_path / "ghost.pptx"
    result = PptxValidator().validate(str(ghost), [_slide(0)])
    assert result["valid"] is False
    assert any("pptx_missing" in w for w in result["warnings"])


# ─────────────── 3. ProductionBuilder.build — E2E 20 slides ───────────────


@pytest.mark.asyncio
async def test_production_build_end_to_end_20_slides(
    synthetic_template: Path, output_dir: Path
) -> None:
    builder = _builder_with_synthetic_template(synthetic_template, output_dir)
    slides = _make_20_mock_slides()
    assert len(slides) == 20  # sanity

    ws_callback = AsyncMock()
    fake_html = MagicMock()
    with patch("weasyprint.HTML", return_value=fake_html):
        pptx_path, pdf_path, report = await builder.build(
            slides=slides,
            course={"id": "e2e-20", "title": "Corso E2E", "duration_hours": 1},
            job_id="job-e2e-20",
            ws_callback=ws_callback,
            image_map={},
        )

    # 1. paths produced
    assert Path(pptx_path).is_file()
    assert pptx_path.endswith("e2e-20_corso.pptx")
    assert pdf_path.endswith("e2e-20_dispensa.pdf")
    # 2. real PPTX has 20 slides
    assert len(Presentation(pptx_path).slides) == 20
    # 3. PDF write_pdf invoked once with the computed path
    fake_html.write_pdf.assert_called_once_with(pdf_path)
    # 4. WebSocket progress steps (87, 92, 95) per BP §07.1 line 2250/2255/2260
    progress_steps = [call.args[1] for call in ws_callback.call_args_list]
    assert progress_steps == [87, 92, 95]
    # 5. Report mirrors the input shape
    assert report["total_slides"] == 20
    assert report["modules_completed"] == 3
    assert report["quiz_count"] >= 1
    assert report["modules_failed"] == 0
    assert report["normative_refs_count"] == 20  # all slides have a ref
    # No validation warnings when slide count matches.
    assert report["warnings"] == []


@pytest.mark.asyncio
async def test_production_build_propagates_validation_warnings(
    synthetic_template: Path, output_dir: Path
) -> None:
    """If the validator returns warnings, the report must surface them."""
    builder = _builder_with_synthetic_template(synthetic_template, output_dir)
    slides = [_slide(i) for i in range(3)]
    ws_callback = AsyncMock()

    fake_html = MagicMock()
    with patch("weasyprint.HTML", return_value=fake_html), patch.object(
        builder.validator,
        "validate",
        return_value={
            "valid": False,
            "slide_count": 3,
            "warnings": ["fake_warning:abc"],
        },
    ):
        _, _, report = await builder.build(
            slides=slides,
            course={"id": "warn", "title": "X"},
            job_id="job-warn",
            ws_callback=ws_callback,
            image_map={},
        )

    assert report["warnings"] == ["fake_warning:abc"]


@pytest.mark.asyncio
async def test_production_build_raises_on_memory_check_fail(
    synthetic_template: Path, output_dir: Path
) -> None:
    """A failing memory guard aborts BEFORE PPTX generation — no PPTX is
    written, no PDF is written, no WebSocket progress is sent."""
    builder = _builder_with_synthetic_template(synthetic_template, output_dir)
    ws_callback = AsyncMock()

    with patch.object(
        pb, "check_memory_before_build", side_effect=MemoryError("nope")
    ), patch("weasyprint.HTML") as html_cls:
        with pytest.raises(MemoryError, match="nope"):
            await builder.build(
                slides=[_slide(0)],
                course={"id": "mem-fail", "title": "X"},
                job_id="job-mem",
                ws_callback=ws_callback,
                image_map={},
            )

    ws_callback.assert_not_called()
    html_cls.assert_not_called()


# ─────────────── 4. _cleanup_tmp behavior ───────────────


def test_cleanup_removes_files_older_than_threshold(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """Only files older than CLEANUP_AGE_SECONDS are removed."""
    monkeypatch.chdir(tmp_path)
    (tmp_path / "output" / "diagrams").mkdir(parents=True)
    (tmp_path / "output" / "images").mkdir(parents=True)

    old = tmp_path / "output" / "diagrams" / "old.png"
    fresh = tmp_path / "output" / "images" / "fresh.png"
    old.write_bytes(b"old")
    fresh.write_bytes(b"fresh")

    # Backdate ``old`` by 2 hours, leave ``fresh`` brand new
    two_hours_ago = time.time() - 2 * 3600
    os.utime(old, (two_hours_ago, two_hours_ago))

    ProductionBuilder.__new__(ProductionBuilder)._cleanup_tmp()

    assert not old.exists(), "old tmp file should have been removed"
    assert fresh.exists(), "fresh file must be preserved"


def test_cleanup_swallows_oserror_silently(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """A file that disappears between glob and remove (race) must not crash."""
    monkeypatch.chdir(tmp_path)
    (tmp_path / "output" / "diagrams").mkdir(parents=True)
    f = tmp_path / "output" / "diagrams" / "racy.png"
    f.write_bytes(b"x")
    os.utime(f, (time.time() - 7200, time.time() - 7200))

    with patch("app.builders.production_builder.os.remove", side_effect=OSError("vanished")):
        # No exception bubbles up — silent swallow per BP §07.1 line 2283
        ProductionBuilder.__new__(ProductionBuilder)._cleanup_tmp()


# ─────────────── 5. structural meta-tests ───────────────


def test_cleanup_patterns_match_bp_07_1() -> None:
    """BP §07.1 line 2278 lists EXACTLY these three glob patterns."""
    assert CLEANUP_PATTERNS == [
        "output/tmp_*",
        "output/diagrams/*.png",
        "output/images/*.png",
    ]


def test_cleanup_age_is_one_hour() -> None:
    """BP §07.1 line 2277 — only files older than 1 hour are removed."""
    assert CLEANUP_AGE_SECONDS == 3600


def test_production_builder_does_not_hold_its_own_semaphore() -> None:
    """REI-3 invariant: the Semaphore(1) belongs to generation_service.

    A future refactor could be tempted to wrap build() in its own
    semaphore; this test makes the intent explicit and would fire if
    someone adds a class-level semaphore here.
    """
    builder = ProductionBuilder.__new__(ProductionBuilder)
    builder_attrs = {
        name
        for name in vars(ProductionBuilder).keys()
        if "semaphore" in name.lower() or "lock" in name.lower()
    }
    assert builder_attrs == set(), (
        f"REI-3 violation: ProductionBuilder owns concurrency primitives "
        f"({builder_attrs}). The Semaphore(1) lives in generation_service."
    )
    _ = builder  # placate ruff
